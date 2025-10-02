import json
import os
import re

import requests
from pydub import AudioSegment

from src import LANGUAGES_ABBR, get_openai_client, get_translator_client
from src.logger_module import get_logger

logger = get_logger("JVD")


def translate_text(text: str, target_language: str, source_language: str | None = "ja") -> str:
    logger.debug(f"ğŸŸ¢ Translating text: {text} from {source_language} to {target_language}")
    client = get_translator_client()

    if source_language is None:
        result = client.detect_language(text)
        source_language = result["language"]

    try:
        result = client.translate(
            text, target_language=target_language, source_language=source_language, format_="text"
        )
        logger.debug(f"ğŸŸ¢ Translation result: {result}")
        return result["translatedText"]
    except Exception as e:
        logger.error(f"ğŸŸ¢ Translation failed: {e}")
        return f"Error: {str(e)}"


def translate_to_all_languages(text: str, source_language: str | None = "ja") -> dict:
    logger.debug(f"ğŸŸ© Translating text to all languages: {text}")
    translations = {}
    for lang_code in LANGUAGES_ABBR.values():
        if lang_code.lower() == source_language.lower():
            continue
        translated_text = translate_text(text, lang_code, source_language)
        translations[lang_code] = translated_text
    order_list = ["EN", "ID", "ES", "VI", "FR", "NE", "BN", "ZH", "KO", "TL", "MY", "HI", "AR", "FA"]
    sorted_x = {key: translations[key] for key in order_list if key in translations}
    return sorted_x


def query_jisho(word: str) -> dict | None:
    logger.debug(f"ğŸ”´ Querying Jisho API for word: {word}")
    url = f"https://jisho.org/api/v1/search/words?keyword={word}"
    response = requests.get(url)
    if response.status_code == 200:
        res = response.json()
        data = res.get("data", [])
        if data:
            logger.debug(f"ğŸ”´ Jisho API response: {data[0]}")
            return data[0]  # Return the first matching entry
        else:
            logger.debug("ğŸ”´ No data found in Jisho API response")
            return None
    else:
        logger.error(f"ğŸ”´Jish o API request failed: {response.status_code}")
        return None


def query_kanji(kanji: str) -> dict | None:
    logger.debug(f"âšªï¸ Querying Kanji API for kanji: {kanji}")
    url = f"https://kanjiapi.dev/v1/kanji/{kanji}"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        if data:
            logger.debug(f"âšªï¸ Kanji API response: {data}")
            return data
        else:
            logger.debug("âšªï¸ No data found in Kanji API response")
            return None
    else:
        logger.error(f"âšªï¸ Kanji API request failed: {response.status_code}")
        return None


def extract_kanji(text) -> list[str]:
    return re.findall(r"[\u4E00-\u9FFF]", text)


class JPWord:
    def __init__(self, initial_word: str, initial_jlpt_level: int, ai_init: bool = False, raw_collocations: str = ""):
        logger.debug(f"ğŸ“– Initializing for word: {initial_word}")
        self.word = initial_word
        self.jlpt_level = initial_jlpt_level
        self.youtube_link = ""
        self.in_db = False
        self.jisho_data = {}
        self.kanji_list = []
        self.kanji_data = {}
        self.collocations = []
        self.ai_explanation = {}
        self.meaning_translations = {}
        if not ai_init:
            return
        self.jisho_data = query_jisho(self.word)
        if self.jisho_data:
            self.kanji_list = extract_kanji(self.word)
            for kanji in self.kanji_list:
                kanji_info = query_kanji(kanji)
                if kanji_info:
                    self.kanji_data[kanji] = kanji_info
        else:
            self.kanji_list = []
            self.kanji_data = {}
        self.collocations = self.get_collocations(raw_collocations)
        self.ai_explanation = self.get_ai_explanation()
        self.translate()
        self.save_to_json()

    def get_collocations(self, raw_collocations: str = "") -> dict:
        logger.info(f"ğŸ–¥ï¸ Getting collocations for word: {self.word}")
        if raw_collocations == "":
            response = get_openai_client().responses.create(
                model="gpt-4o-mini",
                input=[
                    {
                        "role": "system",
                        "content": [
                            {
                                "type": "input_text",
                                "text": "\n".join(
                                    [
                                        "Search on the web for finding the most common and useful collocations of the given Japanese word.",
                                        "A collocation is a natural combination (pattern) of words that are often used together (some usual patterns are: noun + verb, adjective + noun, verb + adverb, adverb + adjective, verb + noun, noun + noun, etc.).",
                                        "For each pattern provide 1 to 5 common collocations that are useful for language learners. Avoid rare or archaic collocations.",
                                        "Consider the collocation only (in kanji) and avoid extra explanations/translations.",
                                        "First write the pattern, then under it report the collocations with bullet point.",
                                    ]
                                ),
                            }
                        ],
                    },
                    {"role": "user", "content": [{"type": "input_text", "text": self.word}]},
                ],
                text={"format": {"type": "text"}},
                reasoning={},
                tools=[
                    {
                        "type": "web_search_preview",
                        "filters": None,
                        "search_context_size": "medium",
                        "user_location": {
                            "type": "approximate",
                            "city": None,
                            "country": None,
                            "region": None,
                            "timezone": None,
                        },
                    }  # type: ignore
                ],
                tool_choice={"type": "web_search_preview"},
                temperature=1,
                max_output_tokens=2024,
                top_p=1,
                store=False,
                include=["web_search_call.action.sources"],  # type: ignore
            )
            raw_collocations = response.output_text
            logger.debug(f"ğŸ–¥ï¸ raw collocations: {raw_collocations}")
        response = get_openai_client().responses.create(
            model="gpt-4o-mini",
            input=[
                {
                    "role": "system",
                    "content": [
                        {
                            "type": "input_text",
                            "text": 'you have to make the JSON format of the given text about collocations of a Japanese vocabulary.\nOutput JSON:\n{ "collocations": list of strings}',
                        }
                    ],
                },
                {"role": "user", "content": [{"type": "input_text", "text": raw_collocations}]},
            ],
            text={"format": {"type": "json_object"}},
            reasoning={},
            tools=[],
            tool_choice="none",
            temperature=1,
            max_output_tokens=2375,
            top_p=1,
            store=False,
            include=["web_search_call.action.sources"],  # type: ignore
        )
        return eval(response.output_text)["collocations"]

    def get_ai_explanation(self) -> dict:
        logger.info(f"ğŸ–¥ï¸ Getting AI explanation for word: {self.word}")
        response = get_openai_client().responses.create(
            model="gpt-5-mini",
            input=[
                {
                    "role": "developer",
                    "content": [
                        {
                            "type": "input_text",
                            "text": "\n".join(
                                [
                                    "You are a helpful assistant that explains Japanese words in detail with respect to the given data collected from Jisho.org and KanjiAPI.dev.",
                                    "Provide clear, concise, comprehensive and short explanations suitable for language learners.",
                                    "The explanation should be in spoken form, as if you are explaining to a friend.",
                                    "Make sure to strictly follow the output format specified in the user prompt. No extra text outside the JSON object.",
                                ]
                            ),
                        }
                    ],
                },
                {
                    "role": "user",
                    "content": "\n".join(
                        [
                            f"The target word is '{self.word}'. Use the following data:\n\nJisho Data: {self.jisho_data}\n\nKanji Data: {self.kanji_data if self.kanji_data else 'No kanji data available.'}\n\nCollocations: {self.collocations}\n\nProvide a comprehensive explanation suitable for a language learner.",
                            "## introduction_japanese: ",
                            "Without mentioning the meaning and reading of the word, name the situations, contexts, and occasions in which this word is used. The explanation should be in Japanese, short and in spoken form suitable for elementary level language learners. Start with the word itself.",
                            "Example: ã€Œã¯ã’ã—ã„ã€ã¯ã€æ„Ÿæƒ…ã€å¤©æ°—ã€å¼·èª¿ãŒå¿…è¦ãªçŠ¶æ³ã«ã¤ã„ã¦ã®ä¼šè©±ã§ã‚ˆãè¦‹ã‹ã‘ã‚‹å½¢å®¹è©ã§ã™ã€‚",
                            "## introduction_english: ",
                            "Translation of introduction_japanese in English. Write the word in its hiragana/katakana form. Start with 'The [adjective/noun/verb ...] [word] ...",
                            "Example: 'ã¯ã’ã—ã„ is an adjective commonly seen in conversations about emotions, weather, and situations that require emphasis.'",
                            "## meanings: ",
                            "List of all meanings grouped by nuance. Each nuance should have its own list of meanings. If there is no nuance, put all meanings in a single list. Each meaning should be short, concise, and in spoken form suitable for elementary level language learners, no extra explanations in parentheses. Avoid rare meanings.",
                            "## meaning_explanation_japanese: ",
                            "Short but complete explanation of literal meanings of the word in Japanese in spoken form suitable for elementary level language learners. Avoid using the word itself, instead use its synonyms or antonyms.",
                            "Example: ã€Œã¯ã’ã—ã„ã€ã¯ã€å¼·ã„æ„Ÿæƒ…ã‚„æ¿€ã—ã„å‹•ä½œã‚’è¡¨ã™è¨€è‘‰ã§ã™ã€‚ä¾‹ãˆã°ã€æ¿€ã—ã„é›¨ã‚„æ¿€ã—ã„è­°è«–ãªã©ã€ä½•ã‹ãŒéå¸¸ã«å¼·ã„ã“ã¨ã‚’ç¤ºã—ã¾ã™ã€‚",
                            "## meaning_explanation_english: ",
                            "Shortly provide the literal meanings of the word. Write the word in its hiragana/katakana form. The explanation should be spoken form and suitable for elementary level language learners.",
                            "Example: The word ã¯ã’ã—ã„ means 'intense', 'fierce', or 'violent'. This word can describe actions, emotions, or conditions that are of great intensity or force.",
                            "## kanji_details: ",
                            "For each kanji in the word, provide its meanings, readings (on'yomi and kun'yomi), and 1 or 2 (max) common words using this kanji except the main word. For each common word, provide its reading and meaning.",
                            "## kanji_explanation_english: ",
                            "Write in a natural, conversational transcript of a teacher explaining the kanji and its meanings. for each kanji (in the order it appears), compose one 3-4 short sentence paragraph that: 1. Describes the kanji's core meaning. 2. except the original word, Presents the 1-2 vocabularies that use that kanji, and how this kanji gives meaning in this vocabulary. Constrains: * Explanation field is the transcription of a speech. Don't use bullet points, parenthesis, new lines, titles, or anything similar. * Do not include the original word in the vocabs. * In the explanation field only insert the hiragana for of Japanese vocabs. No kanjis. * Explanation starts with English phrases such as:  The [first/second/...] kanji means ...",
                            "## synonyms: ",
                            "List the 1 (maximum 2) most commonly used synonym for the provided Japanese vocabulary word (no readings or any other extra text, perferebly in kanji). Excluding the original word.",
                            "The format for each synonym is: kanji : reading : meaning. For example: é€Ÿã„ : ã¯ã‚„ã„ : fast, quick",
                            "## synonyms_explanation: ",
                            'provide the English transcription of a very short explanation about the synonyms listed, including their nuances and meanings. If there were no synonyms in the list, say "No common synonyms found." Constraints 1. Only insert the hiragana for of Japanese vocabs. No kanjis. 2. Explanation starts with English phrases such as:  The most common synonym[s] of the [word] [are/is] ... 3. Very shortly explain the nuances of each synonym and antonym listed, and how they look like the original word.',
                            "## antonyms: ",
                            "List the 1 (maximum 2) most commonly used antonym for the provided Japanese vocabulary word (no readings or any other extra text, perferebly in kanji). Excluding the original word.",
                            "The format for each antonym is: kanji : reading : meaning. For example: é…ã„ : ãŠãã„ : slow",
                            "## antonyms_explanation: ",
                            'provide the English transcription of a very short explanation about the antonyms listed, including their nuances and meanings. If there were no antonyms in the list, say "No common antonyms found." Constraints 1. Only insert the hiragana for of Japanese vocabs. No kanjis. 2. Explanation starts with English phrases such as:  The most common antonym[s] of the [word] [are/is] ... 3. Very shortly explain the nuances of each synonym and antonym listed, and how they differ from the original word.',
                            "## Examples: ",
                            "With respect to the provided collocations, provide 5 - 7 example short and simple sentences using the target word in different contexts useful for language learners.",
                            "Each example should be in kanji and its furigana form. In the furigana form, provide the reading of kanjis in parenthesis right after the kanji. If there are no kanjis in the sentence, just write the sentence as is.",
                            "For example for the word 'ã¯ã’ã—ã„', an example sentence in kanji could be 'å½¼ã¯ã¯ã’ã—ã„æ„Ÿæƒ…ã‚’æŒã£ã¦ã„ã‚‹ã€‚' and its furigana would be 'å½¼(ã‹ã‚Œ)ã¯ã¯ã’ã—ã„æ„Ÿæƒ…(ã‹ã‚“ã˜ã‚‡ã†)ã‚’æŒ(ã‚‚)ã£ã¦ã„ã‚‹ã€‚'.",
                            "make sure the sentences are short, simple and easy to understand for elementary level language learners.",
                            "make sure you follow the rules for furigana.",
                            "## output_format: ",
                            "The output should be in JSON format:",
                            """
{
  "kanji": the word's kanji (if the kanji of the word is not commonly used, write the word itself - no kanji),
  "reading": the word's reading,
  "introduction_japanese": string,
  "introduction_english": string,
  "meanings": [[list of all meanings in nuance1], [list of all meanings in nuance2], ...],
  "meaning_explanation_japanese": string,
  "meaning_explanation_english": string,
  "kanji_explanation_english": string
  "kanji_details": [
    {
      "kanji": "the kanji character",
      "common_words":  # 1 or 2 common words using this kanji except the main word
        ["word1 in kanji (word1 in hiragana): meaning", "word2 in kanji (word2 in hiragana): meaning"]
    },
    ...
  ],
  "synonyms": ["synonym1 in kanji : synonym1 in hiragana : meaning", ...],
  "synonyms_explanation": string,
  "antonyms": ["antonym1 in kanji : antonym1 in hiragana : meaning", ...],
  "antonyms_explanation": string,
  "examples": [{"kanji": string, "furigana": string}, ...]
}
""",
                        ]
                    ),
                },
            ],
            text={"format": {"type": "json_object"}, "verbosity": "medium"},
            reasoning={
                "effort": "medium",
                "summary": None,
            },
            tools=[],
            tool_choice="none",
            store=False,
            include=["reasoning.encrypted_content", "web_search_call.action.sources"],  # type: ignore
        )
        logger.debug(f"ğŸŸ© AI explanation response: {response.output_text}")
        return eval(response.output_text)

    def translate(self):
        logger.info(f"ğŸŸ© Translating AI explanation for word: {self.word}")
        self.ai_explanation["meaning_translations"] = {}
        for nuance in self.ai_explanation.get("meanings", []):
            nuance_translations = translate_to_all_languages(nuance[0], source_language="en")
            self.ai_explanation["meaning_translations"][nuance[0]] = nuance_translations
            self.ai_explanation["meaning_translations"][nuance[0]]["EN"] = nuance[1:]
        for example in self.ai_explanation.get("examples", []):
            example["translations"] = translate_to_all_languages(example["kanji"])

    def to_dict(self) -> dict:
        return {
            "version": "0.2.0",
            "word": self.word,
            "jlpt_level": self.jlpt_level,
            "youtube_link": self.youtube_link,
            "in_db": self.in_db,
            "ai_explanation": self.ai_explanation,
            "jisho_data": self.jisho_data,
            "kanji_list": self.kanji_list,
            "kanji_data": self.kanji_data,
            "collocations": self.collocations,
        }

    def save_to_json(self) -> None:
        dir_path = os.path.join("output", self.word)
        if not os.path.exists(dir_path):
            os.makedirs(dir_path)
        file_path = os.path.join(dir_path, f"{self.word}.json")
        with open(file_path, "w", encoding="utf-8") as f:
            json.dump(self.to_dict(), f, ensure_ascii=False, indent=4)

    @staticmethod
    def load_from_json(word: str) -> "JPWord":
        # file_path = os.path.join("output", word, f"{word}.json")
        file_path = os.path.join("resources", "words", f"{word}.json")
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        obj = JPWord(data["word"], data["jlpt_level"], ai_init=False)
        obj.youtube_link = data["youtube_link"]
        obj.in_db = data["in_db"]
        obj.ai_explanation = data["ai_explanation"]
        obj.jisho_data = data["jisho_data"]
        obj.kanji_list = data["kanji_list"]
        obj.kanji_data = data["kanji_data"]
        obj.collocations = data["collocations"]
        return obj

    def __str__(self) -> str:
        from pprint import pformat

        return pformat(
            self.to_dict(),
            indent=1,
            width=120,
            compact=False,
            sort_dicts=False,
        )

    def pptx_generation(self, num_examples: int | None = 4) -> None:
        logger.info(f"ğŸŸ¦ Generating PPTX for word: {self.word}")
        file_name = f"./Output/{self.word}/{self.word} JLPT N{self.jlpt_level} Vocabulary.pptx"
        if os.path.exists(file_name):
            return

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import (
            MSO_AUTO_SIZE,
            MSO_VERTICAL_ANCHOR,
            PP_PARAGRAPH_ALIGNMENT,
        )
        from pptx.util import Inches, Pt

        prs = Presentation("resources/pptx_templates/template.pptx")

        # Title slide
        first_slide = prs.slides.add_slide(prs.slide_layouts[0])
        presentation_title = first_slide.shapes.title
        if presentation_title:
            presentation_title.text = self.ai_explanation.get("kanji", self.word)
            presentation_title.text_frame.paragraphs[0].font.size = Pt(160)
            presentation_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 95, 154)

        presentation_subtitle = first_slide.placeholders[1]
        presentation_subtitle.text = self.ai_explanation.get("reading", self.word)  # type: ignore
        presentation_subtitle.text_frame.paragraphs[0].font.size = Pt(100)  # type: ignore
        presentation_subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(192, 79, 21)  # type: ignore

        first_slide.shapes.add_movie(
            f"./output/{self.word}/audio/0_introduction.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Definitions slide
        definitions_slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = definitions_slide.shapes.add_textbox(Inches(5), Inches(0), Inches(40 / 3 - 5), Inches(7.5))
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        shape.text_frame.word_wrap = True
        for en, translations in self.ai_explanation["meaning_translations"].items():
            p = shape.text_frame.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            run.text = en + "\n"
            run.font.size = Pt(115)
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(33, 95, 154)
            if len(translations["EN"]):
                run = p.add_run()
                run.text = ", ".join(translations["EN"]) + "\n"
                run.font.size = Pt(32)
                run.font.color.rgb = RGBColor(192, 79, 21)
                run.font.name = "Berlin Sans FB"
                p.space_after = Pt(0)
                p.line_spacing = 0.9
            for code, translation in translations.items():
                if code.lower() == "en":
                    continue
                run_code = p.add_run()
                run_code.text = f"     {code}"
                run_code.font.size = Pt(16)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)

                # Add ordinary run for the translation
                run_translation = p.add_run()
                run_translation.text = f"{translation}"
                run_translation.font.size = Pt(20)
                run_translation.font.name = "Berlin Sans FB"
                run_translation.font.color.rgb = RGBColor(0, 0, 0)

        definitions_slide.shapes.add_movie(
            f"./output/{self.word}/audio/1_definition.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Add a slide for Kanji breakdown
        kanji_slide = prs.slides.add_slide(prs.slide_layouts[6])
        for i, kanji in enumerate(self.kanji_list):
            kanji_shape = kanji_slide.shapes.add_textbox(
                Inches(0), Inches(i * 7.5 / 2), Inches(40 / 3 * 0.3), Inches(7.5 / 2)
            )
            p = kanji_shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = kanji
            run.font.size = Pt(250)
            run.font.bold = True
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(33, 95, 154)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.space_after = Pt(0)
            p.line_spacing = 0.9

            kanji_explanation_shape = kanji_slide.shapes.add_textbox(
                Inches(40 / 3 * 0.3),
                Inches(i * 7.5 / 2),
                Inches(40 / 3 * 0.7),
                Inches(7.5 / 2),
            )
            kanji_explanation_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            kanji_explanation_shape.text_frame.word_wrap = True
            p = kanji_explanation_shape.text_frame.paragraphs[-1]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            run = p.add_run()
            run.text = "Meanings:"
            run.font.size = Pt(24)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(127, 127, 127)
            run = p.add_run()
            run.text = f" {', '.join(self.kanji_data[kanji]['meanings'][:4])}\n"
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = "Readings:"
            run.font.size = Pt(24)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(127, 127, 127)
            run = p.add_run()
            run.text = f" {', '.join(self.kanji_data[kanji]['kun_readings'] + self.kanji_data[kanji]['on_readings'])}\n"
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            for vocab in self.ai_explanation["kanji_details"][i]["common_words"]:
                run.text += f"\n{vocab}"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(192, 79, 21)
            p.space_after = Pt(0)
            p.line_spacing = 0.9

        kanji_slide.shapes.add_movie(
            f"./output/{self.word}/audio/2_kanji_explanation.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Add a slide for examples
        for i, example in enumerate(self.ai_explanation["examples"][:num_examples]):
            example_slide = prs.slides.add_slide(prs.slide_layouts[6])

            top_shape = example_slide.shapes.add_textbox(Inches(0), Inches(0), Inches(40 / 3), Inches(3))
            top_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            top_shape.text_frame.word_wrap = True
            top_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = top_shape.text_frame.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            run.text = example["kanji"]
            run.font.size = Pt(70)
            run.font.bold = True
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(33, 95, 154)
            run = p.add_run()
            run.text = f"\n{example['furigana']}"  # type: ignore
            run.font.size = Pt(33)
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(192, 79, 21)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.space_after = Pt(0)
            p.line_spacing = 0.9

            left_languages = [
                "EN",
                "ID",
                "ES",
                "VI",
                "FR",
                "NE",
                "AR",
            ]

            left_shape = example_slide.shapes.add_textbox(
                Inches(20 / 3 * 0.025), Inches(3), Inches(20 / 3 * 0.95), Inches(4.5)
            )
            left_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            left_shape.text_frame.word_wrap = True
            left_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for lang in left_languages:
                p = left_shape.text_frame.paragraphs[-1]
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                run_code = p.add_run()
                run_code.text = f"{lang}: "
                run_code.font.size = Pt(20)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)
                run_translation = p.add_run()
                run_translation.text = f"{example['translations'][lang]}\n"
                run_translation.font.size = Pt(25)
                run_translation.font.name = "Berlin Sans FB"
            right_shape = example_slide.shapes.add_textbox(
                Inches(20 / 3 * 1.025), Inches(3), Inches(20 / 3 * 0.95), Inches(4.5)
            )
            right_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            right_shape.text_frame.word_wrap = True
            right_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            right_languages = [
                "BN",
                "ZH",
                "KO",
                "TL",
                "MY",
                "HI",
                "FA",
            ]
            for lang in right_languages:
                p = right_shape.text_frame.paragraphs[-1]
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                run_code = p.add_run()
                run_code.text = f"{lang}: "
                run_code.font.size = Pt(20)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)
                run_translation = p.add_run()
                run_translation.text = f"{example['translations'][lang]}\n"  # type: ignore
                run_translation.font.size = Pt(25)
                run_translation.font.bold = False
                run_translation.font.name = "Berlin Sans FB"

            example_slide.shapes.add_movie(
                f"./output/{self.word}/audio/{3 + i}_example.wav",
                left=Pt(0),
                top=Pt(-50),
                width=Pt(50),
                height=Pt(50),
                mime_type="audio/x-wav",
            )

        # Add a slide for synonyms and antonyms
        synonyms_slide = prs.slides.add_slide(prs.slide_layouts[6])

        synonyms_shape = synonyms_slide.shapes.add_textbox(Inches(0), Inches(0), Inches(20 / 3), Inches(7.5))
        synonyms_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        synonyms_shape.text_frame.word_wrap = True
        p = synonyms_shape.text_frame.paragraphs[-1]
        run = p.add_run()
        run.text = "\nSYNONYMS\n\n"
        run.font.size = Pt(40)
        run.font.name = "Berlin Sans FB Demi"
        run.font.color.rgb = RGBColor(33, 95, 154)
        for synonym in self.ai_explanation.get("synonyms", []):
            kanji, reading, meaning = synonym.split(":")
            run = p.add_run()
            run.text = f"{kanji}\n"
            run.font.size = Pt(54)
            run.font.bold = True  # type: ignore
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(33, 95, 154)
            run = p.add_run()
            run.text = f"{reading}\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = f"{meaning}\n\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(0)
            p.line_spacing = 0.9
        synonyms_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        antonyms_shape = synonyms_slide.shapes.add_textbox(Inches(20 / 3), Inches(0), Inches(20 / 3), Inches(7.5))
        antonyms_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        antonyms_shape.text_frame.word_wrap = True
        p = antonyms_shape.text_frame.paragraphs[-1]
        run = p.add_run()
        run.text = "\nANTONYMS\n\n"
        run.font.size = Pt(40)
        run.font.name = "Berlin Sans FB Demi"
        run.font.color.rgb = RGBColor(192, 79, 21)
        for antonym in self.ai_explanation.get("antonyms", []):
            kanji, reading, meaning = antonym.split(":")
            run = p.add_run()
            run.text = f"{kanji}\n"
            run.font.size = Pt(54)
            run.font.bold = True  # type: ignore
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(192, 79, 21)
            run = p.add_run()
            run.text = f"{reading}\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = f"{meaning} \n\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(0)
            p.line_spacing = 0.9
        antonyms_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        synonyms_slide.shapes.add_movie(
            f"./output/{self.word}/audio/100_synonyms_antonyms.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        prs.save(file_name)

    def tts(self, num_examples: int | None = 4) -> None:
        os.makedirs(f"./Output/{self.word}/audio", exist_ok=True)

        # Generate audio for introduction
        introduction_en_audio_path = f"./output/{self.word}/audio/introduction_en.mp3"
        if not os.path.exists(introduction_en_audio_path):
            logger.info("ğŸ”Š Generating audio for english introduction")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['introduction_english']}",
                instructions=f"English mixed with Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(introduction_en_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        introduction_jp_audio_path = f"./output/{self.word}/audio/introduction_jp.mp3"
        if not os.path.exists(introduction_jp_audio_path):
            logger.info("ğŸ”Š Generating audio for japanese introduction")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['introduction_japanese']}",
                instructions=f"Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(introduction_jp_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        introduction_audio = AudioSegment.silent(duration=100)
        introduction_audio = AudioSegment.from_mp3(introduction_en_audio_path).apply_gain(12)
        introduction_audio += AudioSegment.silent(duration=500)
        introduction_audio += AudioSegment.from_mp3(introduction_jp_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/0_introduction.wav", "wb") as title_file:
            introduction_audio.export(title_file, format="wav")

        # Generate audio for meaning explanation
        definition_en_audio_path = f"./output/{self.word}/audio/definition_en.mp3"
        if not os.path.exists(definition_en_audio_path):
            logger.info("ğŸ”Š Generating audio for english definition")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['meaning_explanation_english']}",
                instructions=f"English mixed with Japanese. Calmly and gently. Correct pronunciation: {self.word}:{self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(definition_en_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        definition_jp_audio_path = f"./output/{self.word}/audio/definition_jp.mp3"
        if not os.path.exists(definition_jp_audio_path):
            logger.info("ğŸ”Š Generating audio for japanese definition")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['meaning_explanation_japanese']}",
                instructions=f"Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(definition_jp_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        definition_audio = AudioSegment.silent(duration=100)
        definition_audio = AudioSegment.from_mp3(definition_en_audio_path).apply_gain(12)
        definition_audio += AudioSegment.silent(duration=500)
        definition_audio += AudioSegment.from_mp3(definition_jp_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/1_definition.wav", "wb") as word_file:
            definition_audio.export(word_file, format="wav")

        # Generate audio for kanji explanations
        kanji_explanation_audio_path = f"./output/{self.word}/audio/kanji_explanation.mp3"
        if not os.path.exists(kanji_explanation_audio_path):
            logger.info("ğŸ”Š Generating audio for kanji explanation")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.ai_explanation["kanji_explanation_english"],
                instructions=f"English. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(kanji_explanation_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        kanji_explanation_audio = AudioSegment.from_mp3(kanji_explanation_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/2_kanji_explanation.wav", "wb") as kanji_file:
            kanji_explanation_audio.export(kanji_file, format="wav")

        # Generate audio for examples
        for i, example in enumerate(self.ai_explanation["examples"][:num_examples]):
            example_jp_audio_path = f"./output/{self.word}/audio/example_{i + 1}_jp.mp3"
            if not os.path.exists(example_jp_audio_path):
                logger.info(f"ğŸ”Š Generating audio for example {i + 1}")
                response = get_openai_client().audio.speech.create(
                    model="gpt-4o-mini-tts",
                    voice="coral",
                    input=example["kanji"],
                    instructions=f"Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                    response_format="mp3",
                    speed=0.95,
                )
                with open(example_jp_audio_path, "wb") as audio_file:
                    audio_file.write(response.content)
            example_en_audio_path = f"./output/{self.word}/audio/example_{i + 1}_en.mp3"
            if not os.path.exists(example_en_audio_path):
                logger.info(f"ğŸ”Š Generating audio for example {i + 1} translation")
                response = get_openai_client().audio.speech.create(
                    model="gpt-4o-mini-tts",
                    voice="coral",
                    input=example["translations"]["EN"],
                    instructions=f"English. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                    response_format="mp3",
                    speed=0.95,
                )
                with open(example_en_audio_path, "wb") as audio_file:
                    audio_file.write(response.content)
            example_jp_audio = AudioSegment.from_mp3(example_jp_audio_path).apply_gain(12)
            example_en_audio = AudioSegment.from_mp3(example_en_audio_path).apply_gain(12)

            example_audio = AudioSegment.silent(duration=100)
            example_audio += example_en_audio
            example_audio += AudioSegment.silent(duration=500)
            example_audio += example_jp_audio
            example_audio += AudioSegment.silent(duration=500)
            example_audio += example_jp_audio
            with open(f"./output/{self.word}/audio/{3 + i}_example.wav", "wb") as example_file:
                example_audio.export(example_file, format="wav")

        # Generate audio for synonyms and antonyms explanations
        synonyms_explanation_audio_path = f"./output/{self.word}/audio/synonyms_explanation.mp3"
        if not os.path.exists(synonyms_explanation_audio_path):
            logger.info("ğŸ”Š Generating audio for synonyms explanation")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.ai_explanation["synonyms_explanation"] + " " + self.ai_explanation["antonyms_explanation"],
                instructions=f"English mixed with Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(synonyms_explanation_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        synonyms_explanation_audio = AudioSegment.from_mp3(synonyms_explanation_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/100_synonyms_antonyms.wav", "wb") as syn_ant_file:
            synonyms_explanation_audio.export(syn_ant_file, format="wav")

    def show_in_streamlit(self, st, auth: dict | None = None) -> None:
        if self.word == self.ai_explanation.get("reading", ""):
            st.markdown(
                f"# {self.word} :orange-badge[N{self.jlpt_level}] :green-badge[{'Common' if self.jisho_data['is_common'] else 'Uncommon'}]"
            )
        else:
            st.markdown(
                f"# {self.word} ({self.ai_explanation['reading']}) :orange-badge[N{self.jlpt_level}] :green-badge[{'Common' if self.jisho_data['is_common'] else 'Uncommon'}]"
            )
        if self.youtube_link:
            st.video(self.youtube_link)

        st.markdown(self.ai_explanation["introduction_english"])
        st.markdown(self.ai_explanation["introduction_japanese"])
        st.markdown("### Meanings")
        st.markdown(self.ai_explanation["meaning_explanation_english"])
        st.markdown(self.ai_explanation["meaning_explanation_japanese"])
        for m, tr in self.ai_explanation["meaning_translations"].items():
            if isinstance(tr["EN"], list):
                tr["EN"].insert(0, m)
                tr["EN"] = ", ".join(tr["EN"])
            with st.container(border=1, horizontal=True):
                for k, v in tr.items():
                    if auth:
                        user_langs = [LANGUAGES_ABBR[lang] for lang in auth.get("preferred_languages", [])]
                        if k not in user_langs:
                            continue
                    st.markdown(f":gray-badge[{k}] {v}")

        col1, col2 = st.columns(2)
        with col1:
            st.markdown("### Synonyms")
            with st.container(border=1):
                for s in self.ai_explanation["synonyms"]:
                    st.markdown(s)
        with col2:
            st.markdown("### Antonyms")
            with st.container(border=1):
                for a in self.ai_explanation["antonyms"]:
                    st.markdown(a)

        st.markdown("### Kanji")
        for i, k in enumerate(self.kanji_list):
            cont = st.container(border=1)
            col1, col2 = cont.columns([2, 4])
            with col1:
                st.image(
                    f"https://raw.githubusercontent.com/KanjiVG/kanjivg/refs/heads/master/kanji/0{self.kanji_data[k]['unicode'].lower()}.svg",
                    width="stretch",
                )
            with col2:
                st.markdown(f":gray-badge[Meaning] {', '.join(self.kanji_data[k]['meanings'])}")
                st.markdown(f":gray-badge[On-yomi] {', '.join(self.kanji_data[k]['on_readings'])}")
                st.markdown(f":gray-badge[Kun-yomi] {', '.join(self.kanji_data[k]['kun_readings'])}")
                st.markdown("---")
                for common_word in self.ai_explanation["kanji_details"][i]["common_words"]:
                    st.markdown(common_word)

        st.markdown("### Examples")
        for ex in self.ai_explanation["examples"]:
            from src.utils import create_html_with_ruby

            with st.expander(ex["kanji"]):
                ruby = create_html_with_ruby(ex["furigana"])
                st.markdown(ruby, unsafe_allow_html=True)
                for key, value in ex["translations"].items():
                    if auth:
                        user_langs = [LANGUAGES_ABBR[lang] for lang in auth.get("preferred_languages", [])]
                        if key not in user_langs:
                            continue
                    st.markdown(f":gray-badge[{key}] {value}")


if __name__ == "__main__":
    # "ä¸Šé”": "åè© + å‹•è© *ï¼ˆã€Œã€œãŒï¼ã‚’ä¸Šé”ã™ã‚‹ï¼ã•ã›ã‚‹ï¼æœ›ã‚€ã€ãªã©ï¼‰ *æŠ€è¡“ãŒä¸Šé”ã™ã‚‹ *èƒ½åŠ›ãŒä¸Šé”ã™ã‚‹ *å®ŸåŠ›ãŒä¸Šé”ã™ã‚‹ *ç·´ç¿’ã—ã¦ä¸Šé”ã™ã‚‹ *ç¿’ã„äº‹ã‚’ä¸Šé”ã•ã›ã‚‹ *å‰¯è© + ä¸Šé”ã™ã‚‹ï¼ä¸Šé”ã•ã›ã‚‹ *ï¼ˆç¨‹åº¦ã‚„æ§˜å­ã‚’è¡¨ã™å‰¯è©ã¨çµ„ã‚€ï¼‰ *é£›èºçš„ã«ä¸Šé”ã™ã‚‹ *ç€å®Ÿã«ä¸Šé”ã™ã‚‹ *æ€¥é€Ÿã«ä¸Šé”ã™ã‚‹ *å°‘ã—ãšã¤ä¸Šé”ã™ã‚‹ *ç›®è¦šã¾ã—ãä¸Šé”ã™ã‚‹ *ä¸Šé” + ã® + åè© *ï¼ˆä¸Šé”ã—ãŸçµæœãƒ»æ€§è³ªã‚’è¡¨ã™åè©ï¼‰ *ä¸Šé”ã®éç¨‹ *ä¸Šé”ã®ç›®å®‰ *ä¸Šé”ã®é€Ÿåº¦ *ä¸Šé”ã®ç¨‹åº¦ *ä¸Šé”ã®è·¡ *ä¸Šé” + ãŒ + å½¢å®¹è© *ï¼ˆã€Œä¸Šé”ãŒæ—©ã„ï¼é…ã„ï¼è‘—ã—ã„ï¼é †èª¿ï¼ç›®ç«‹ã¤ã€ãªã©ï¼‰ *ä¸Šé”ãŒæ—©ã„ *ä¸Šé”ãŒé…ã„ *ä¸Šé”ãŒè‘—ã—ã„ *ä¸Šé”ãŒé †èª¿ã  *ä¸Šé”ãŒç›®ç«‹ã¤ *å‹•è© + ä¸Šé” *ï¼ˆä»–ã®å‹•è©ãŒã€Œä¸Šé”ã€ã«ã‹ã‹ã‚‹ãƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ *å®Ÿæ„Ÿã—ã¦ä¸Šé”ã‚’æ„Ÿã˜ã‚‹ *åŠªåŠ›ã—ã¦ä¸Šé”ã‚’å›³ã‚‹ *æ§˜å­ã‚’è¦‹ã¦ä¸Šé”ã‚’åˆ¤æ–­ã™ã‚‹ *ç¶™ç¶šã—ã¦ä¸Šé”ã‚’ç›®æŒ‡ã™ *æ¯”è¼ƒã—ã¦ä¸Šé”ã‚’è¦‹ã‚‹",
    # "ã†ã¾ã„": "å‰¯è© + ã†ã¾ã„ï¼ä¸Šæ‰‹ã„ï¼æ—¨ã„ *ï¼ˆæ§˜å­ãƒ»ç¨‹åº¦ã‚’å¼·ã‚ãŸã‚Šç·©ã‚ãŸã‚Šï¼‰ *ã¨ã¦ã‚‚ï¼ã™ã”ãã†ã¾ã„ *ã™ã”ãä¸Šæ‰‹ã„ *ã‹ãªã‚Šã†ã¾ã„ *ãšã„ã¶ã‚“ã†ã¾ã„ *æ¡ˆå¤–ã†ã¾ã„ *ã†ã¾ã„ï¼ä¸Šæ‰‹ã„ï¼æ—¨ã„ + åè© *ï¼ˆãã®â€œã†ã¾ã•â€ã‚’å¯¾è±¡åŒ–ã™ã‚‹ï¼‰ *ã†ã¾ã„æ–™ç† *ä¸Šæ‰‹ã„è¡¨ç¾ *ã†ã¾ã„è©± *ä¸Šæ‰‹ã„ä½¿ã„æ–¹ *ä¸Šæ‰‹ã„æ¼”æŠ€ *å‹•è© + ã†ã¾ãï¼ä¸Šæ‰‹ã«ï¼æ—¨ã + å‹•è© *ï¼ˆè¡Œç‚ºãƒ»å‹•ä½œã«å¯¾ã—ã¦ã€Œã†ã¾ãã€œã™ã‚‹ã€ãªã©ï¼‰ *ã†ã¾ãã‚„ã‚‹ *ä¸Šæ‰‹ã«è©±ã™ *ã†ã¾ãä¼ãˆã‚‹ *ä¸Šæ‰‹ã«ä½¿ã† *ã†ã¾ãã¾ã¨ã‚ã‚‹ *ã€œãŒ + ã†ã¾ã„ï¼ä¸Šæ‰‹ã„ï¼æ—¨ã„ *ï¼ˆä¸»èªã®èƒ½åŠ›ãƒ»æŠ€å·§ã‚’è¡¨ã™ï¼‰ *å½¼ã¯ä¸Šæ‰‹ã„ï¼ˆã ï¼‰ *å›ãŒã†ã¾ã„ *ã‚ã®äººãŒã†ã¾ã„ *èª°ãŒä¸Šæ‰‹ã„ã‹ *ãã‚ŒãŒã†ã¾ã„ï¼ˆä½¿ã‚ã‚Œã‚‹ï¼‰ *ã†ã¾ã„ + ï¼ˆå½¢å®¹è©ãƒ»å¥ï¼‰ *ï¼ˆã€Œã†ã¾ã„ã€ãŒä¿®é£¾èªå¥ã¨çµ„ã¾ã‚Œã‚‹å½¢ï¼‰ *ã†ã¾ã„å…·åˆã« *ã†ã¾ã„åŠ æ¸›ã« *ã†ã¾ã„èª¿å­ã§ *ã†ã¾ã„è©±ã«ä¹—ã‚‹ *ã†ã¾ã„ã‚¿ã‚¤ãƒŸãƒ³ã‚°ã§",
    # "ä¸å®‰": "å½¢å®¹è© + ä¸å®‰ * å¼·ã„ä¸å®‰ * å¤§ããªä¸å®‰ * é‡å¤§ãªä¸å®‰ * æ ¹å¼·ã„ä¸å®‰ * å¸¸ã«ä¸å®‰ * ä¸å®‰ + ã‚’ + å‹•è© * ä¸å®‰ã‚’æ„Ÿã˜ã‚‹ * ä¸å®‰ã‚’æŠ±ã * ä¸å®‰ã‚’è¦šãˆã‚‹ * ä¸å®‰ã‚’ç…½ã‚‹ * ä¸å®‰ã‚’å’Œã‚‰ã’ã‚‹ * ä¸å®‰ + ãŒ + å½¢å®¹è©ï¼å‹•è© * ä¸å®‰ãŒã‚ã‚‹ * ä¸å®‰ãŒæ¶ˆãˆã‚‹ * ä¸å®‰ãŒå‹Ÿã‚‹ * ä¸å®‰ãŒé«˜ã¾ã‚‹ * ä¸å®‰ãŒå’Œã‚‰ã * å‰¯è© + ä¸å®‰ * éå¸¸ã«ä¸å®‰ * ã‚„ã‚„ä¸å®‰ * ãŸã ä¸å®‰ * ä½•ã¨ãªãä¸å®‰ * å°‘ã—ä¸å®‰ * ä¸å®‰ + ã® + åè© * ä¸å®‰ã®ç¨® * ä¸å®‰ã®å£° * ä¸å®‰ã®åŸå›  * ä¸å®‰ã®å¿µ * ä¸å®‰ã®å…†ã—",
    # "ä¸¸ã„": "å½¢å®¹å‹•è© + åè© * ä¸¸ã„å½¢ * ä¸¸ã„é ­ * ä¸¸ã„ç›® * ä¸¸ã„é¡” * ä¸¸ã„æœº * ä¸¸ã„ + åè© + ã‚’ + å‹•è© * ä¸¸ã„çŸ³ã‚’æŠ•ã’ã‚‹ * ä¸¸ã„ã‚³ã‚¤ãƒ³ã‚’æ‰‹ã«å–ã‚‹ * ä¸¸ã„çš¿ã‚’æ´—ã† * ä¸¸ã„é¢¨èˆ¹ã‚’è†¨ã‚‰ã¾ã›ã‚‹ * ä¸¸ã„ãƒœãƒ¼ãƒ«ã‚’æ‰“ã¤ * å‰¯è©ãƒ»ä¿®é£¾èª + ä¸¸ã„ * éå¸¸ã«ä¸¸ã„ * æ¯”è¼ƒçš„ä¸¸ã„ * ä¸¸ãã¦å…‰ã‚‹ * ãã‚Œã„ã«ä¸¸ã„ * ã»ã¼ä¸¸ã„ * åè© + ãŒ + ä¸¸ã„ï¼ä¸¸ã * æœˆãŒä¸¸ã„ * é¡”ãŒä¸¸ã„ * ç›®ãŒä¸¸ã„ * ãŠè…¹ãŒä¸¸ã„ * é­šãŒä¸¸ã„ * æ…£ç”¨è¡¨ç¾ãƒ»æ¯”å–©çš„è¡¨ç¾ * ä¸¸ãåã‚ã‚‹ * ä¸¸ããªã‚‹ * ä¸¸ãæ²»ã‚ã‚‹ * ä¸¸ãè¨€ã† * ä¸¸ãåã¾ã‚‹",
    # "ä¹—å®¢": "åè© + å‹•è©ï¼å‹•è©ï¼‹åè© ä¹—å®¢ã‚’ä¹—ã›ã‚‹ ä¹—å®¢ã‚’é™ã‚ã™ ä¹—å®¢ãŒä¹—ã‚‹ ä¹—å®¢ãŒä¸‹è»Šã™ã‚‹ ä¹—å®¢ã‚’æ¡ˆå†…ã™ã‚‹ ã€œã« + ä¹—å®¢ ï¼ˆä¹—å®¢ã‚’ç›®çš„èªãƒ»å¯¾è±¡å ´é¢ã§å–ã‚‹è¡¨ç¾ï¼‰ åˆ—è»Šã«ä¹—å®¢ã‚’ä¹—ã›ã‚‹ é£›è¡Œæ©Ÿã«ä¹—å®¢ã‚’ä¹—ã›ã‚‹ ãƒã‚¹ã«ä¹—å®¢ã‚’è¼‰ã›ã‚‹ é›»è»Šã«ä¹—å®¢ãŒæ®ºåˆ°ã™ã‚‹ é§…ã«ä¹—å®¢ãŒã‚ãµã‚Œã‚‹ ä¹—å®¢ + ã® + åè© ä¹—å®¢æ•° ä¹—å®¢æ¡ˆå†… ä¹—å®¢ã‚µãƒ¼ãƒ“ã‚¹ ä¹—å®¢å±¤ ä¹—å®¢é‹é€ ä¹—å®¢ + ãŒ + å‹•è©ï¼å½¢å®¹è© ä¹—å®¢ãŒå¤šã„ ä¹—å®¢ãŒæ¸›ã‚‹ ä¹—å®¢ãŒå¢—ãˆã‚‹ ä¹—å®¢ãŒæ··é›‘ã™ã‚‹ ä¹—å®¢ãŒè‹¦æƒ…ã‚’è¨€ã† å‰¯è©ï¼ä¿®é£¾èª + ä¹—å®¢ å¤šãã®ä¹—å®¢ å°‘æ•°ã®ä¹—å®¢ ä¸»ãªä¹—å®¢ è¦³å…‰å®¢ä¹—å®¢ å·¨å¤§ãªä¹—å®¢æ•°",
    # "äºˆæœŸ": "åè© + å‹•è©ï¼ˆã€Œã€œã‚’äºˆæœŸã™ã‚‹ã€å‹ï¼‰ äº‹æ…‹ã‚’äºˆæœŸã™ã‚‹ çµæœã‚’äºˆæœŸã™ã‚‹ å¤‰åŒ–ã‚’äºˆæœŸã™ã‚‹ åå¿œã‚’äºˆæœŸã™ã‚‹ äº‹å‰ã«äºˆæœŸã™ã‚‹ å‰¯è© + å‹•è©ï¼ˆäºˆæœŸã¨ã¨ã‚‚ã«ä½¿ã‚ã‚Œã‚‹å‰¯è©ï¼‰ ååˆ†ã«äºˆæœŸã™ã‚‹ ãŠãŠã‚ˆãäºˆæœŸã™ã‚‹ å…¨ãäºˆæœŸã—ãªã„ ã¾ã£ãŸãäºˆæœŸã—ãªã‹ã£ãŸ äºˆæœŸã›ãš åè© + ã« + å‹•è©ï¼ˆã€Œã€œã«äºˆæœŸã•ã‚Œã‚‹ï¼ã€œã«äºˆæœŸã•ã‚Œãªã„ã€ãªã©ï¼‰ äºˆæœŸã«åã™ã‚‹ äºˆæœŸã«æ²¿ã† äºˆæœŸã«é”ã™ã‚‹ äºˆæœŸã«å±Šã äºˆæœŸã«é•ã‚ãš åè©ï¼‹å½¢å®¹è©ï¼ˆåè©ã€ŒäºˆæœŸã€ã¨çµ„ã‚€å½¢å®¹è©ä¿®é£¾ï¼‰ äºˆæœŸã›ã¬å‡ºæ¥äº‹ äºˆæœŸå¤–ã®åå¿œ äºˆæœŸã—ãŒãŸã„çµæœ äºˆæœŸã©ãŠã‚Šã®å±•é–‹ äºˆæœŸä»¥ä¸Šã®æˆæœ åè© + åè©ï¼ˆã€Œã€œã®äºˆæœŸã€ãªã©ã®æ§‹é€ ï¼‰ å°†æ¥ã®äºˆæœŸ æœŸå¾…ã¨ã®äºˆæœŸ äºˆæœŸã®ç¯„å›² äºˆæœŸã®é•ã„ äºˆæœŸã®å¤‰åŒ–",
    # "äºˆé˜²": "åè© + ã‚’ + å‹•è©ï¼ˆã€Œã€œã‚’äºˆé˜²ã™ã‚‹ï¼ã€œã‚’äºˆé˜²ã§ãã‚‹ã€å‹ï¼‰ æ„ŸæŸ“ã‚’äºˆé˜²ã™ã‚‹ ç—…æ°—ã‚’äºˆé˜²ã™ã‚‹ äº‹æ•…ã‚’äºˆé˜²ã™ã‚‹ ç½å®³ã‚’äºˆé˜²ã™ã‚‹ è™«æ­¯ã‚’äºˆé˜²ã™ã‚‹ åè© + ã® + åè©ï¼ˆã€Œã€œã®äºˆé˜²ã€ã€Œã€œäºˆé˜²ã€å‹ï¼‰ ç—…æ°—ã®äºˆé˜² æ„ŸæŸ“äºˆé˜² ç«ç½äºˆé˜² ç½å®³äºˆé˜² å¥åº·äºˆé˜²ï¼ˆâ€»ã‚ã¾ã‚Šæ—¥å¸¸èªã§ã¯ãªã„ãŒä½¿ã‚ã‚Œã‚‹ã“ã¨ã‚ã‚Šï¼‰ åè© + ã« + å‹•è©ï¼å½¢å®¹å‹•è©ï¼ˆã€Œã€œã«äºˆé˜²ç­–ã‚’è¬›ã˜ã‚‹ã€ã€Œã€œã«äºˆé˜²çš„ãªã€ãªã©ï¼‰ æ„ŸæŸ“ã«å¯¾ã—ã¦äºˆé˜²ç­–ã‚’è¬›ã˜ã‚‹ ç—…æ°—ã«å¯¾ã—ã¦äºˆé˜²æ‰‹æ®µã‚’ã¨ã‚‹ ç½å®³ã«å‚™ãˆã¦äºˆé˜²çš„æªç½®ã‚’ã¨ã‚‹ è¢«å®³ã«å¯¾ã™ã‚‹äºˆé˜²çš„å¯¾å¿œ ãƒªã‚¹ã‚¯ã«å¯¾ã™ã‚‹äºˆé˜²çš„å‡¦ç½® å½¢å®¹è© + åè©ï¼ˆã€Œã€œäºˆé˜²ã€ã€Œã€œçš„äºˆé˜²ã€ãªã©ä¿®é£¾èªï¼‰ äºˆé˜²æ¥ç¨® äºˆé˜²æªç½® äºˆé˜²çš„ äºˆé˜²åŠ¹æœ äºˆé˜²åŒ»ç™‚ å‰¯è©ï¼æ…£ç”¨å¥ + å‹•è©ï¼ˆäºˆé˜²ã‚’è¡¨ã™è¡¨ç¾è£œåŠ©èªï¼‰ ååˆ†ã«äºˆé˜²ã™ã‚‹ äº‹å‰ã«äºˆé˜²ã™ã‚‹ é©åˆ‡ã«äºˆé˜²ã™ã‚‹ å¾¹åº•ã—ã¦äºˆé˜²ã™ã‚‹ äºˆé˜²ã§ãã‚‹ï¼äºˆé˜²å¯èƒ½",
    # "äº‹ä»¶": "åè© + ã‚’ + å‹•è©ï¼ˆã€Œã€œã‚’äº‹ä»¶ã¨ã¿ãªã™ï¼ã€œã‚’æ‰±ã†ã€ãªã©ï¼‰ äº‹ä»¶ã‚’æ‰±ã† äº‹ä»¶ã‚’ç«‹ä»¶ã™ã‚‹ äº‹ä»¶ã‚’æœæŸ»ã™ã‚‹ äº‹ä»¶ã‚’å ±é“ã™ã‚‹ äº‹ä»¶ã‚’è§£æ±ºã™ã‚‹ åè© + ã® + åè©ï¼ˆã€Œã€œäº‹ä»¶ã€ã€Œã€œã®äº‹ä»¶ã€å‹ï¼‰ æ®ºäººäº‹ä»¶ äº¤é€šäº‹æ•…äº‹ä»¶ çªƒç›—äº‹ä»¶ é‡å¤§äº‹ä»¶ æœªè§£æ±ºäº‹ä»¶ åè© + ã« + å‹•è©ï¼å½¢å®¹è© äº‹ä»¶ã«ç™ºå±•ã™ã‚‹ äº‹ä»¶ã«å·»ãè¾¼ã¾ã‚Œã‚‹ äº‹ä»¶ã«é–¢ä¸ã™ã‚‹ äº‹ä»¶ã«è‡³ã‚‹ äº‹ä»¶ã«è¦‹èˆã‚ã‚Œã‚‹ å½¢å®¹è© + åè©ï¼ˆã€Œã€œäº‹ä»¶ã€ã€Œã€œãªäº‹ä»¶ã€å‹ï¼‰ è¡æ’ƒçš„ãªäº‹ä»¶ ç•°å¸¸ãªäº‹ä»¶ æ‚²æƒ¨ãªäº‹ä»¶ æƒ¨åŠ‡çš„ãªäº‹ä»¶ æ„å¤–ãªäº‹ä»¶ å‰¯è©ï¼æ¥ç¶šèª + åè©ï¼å‹•è©ï¼ˆã€Œã€œäº‹ä»¶ãŒèµ·ã“ã‚‹ã€ã€Œã€œã«ã‚ˆã‚Šäº‹ä»¶ãŒã€œã€ãªã©ï¼‰ çªå¦‚ã¨ã—ã¦äº‹ä»¶ãŒèµ·ã“ã‚‹ ç¤¾ä¼šå•é¡ŒåŒ–ã™ã‚‹äº‹ä»¶ æ¬¡ã€…ã¨äº‹ä»¶ãŒç™ºè¦šã™ã‚‹ äº‹ä»¶ãŒç›¸æ¬¡ã äº‹ä»¶ãŒè¡¨é¢åŒ–ã™ã‚‹",
    # "äº¤æ›": "åè© + ã‚’ + å‹•è©ï¼ˆã€œã‚’äº¤æ›ã™ã‚‹ï¼ã€œã‚’äº¤æ›ã§ãã‚‹ ç­‰ï¼‰ æƒ…å ±ã‚’äº¤æ›ã™ã‚‹ ååˆºã‚’äº¤æ›ã™ã‚‹ é€£çµ¡å…ˆã‚’äº¤æ›ã™ã‚‹ æ„è¦‹ã‚’äº¤æ›ã™ã‚‹ éƒ¨å“ã‚’äº¤æ›ã™ã‚‹ åè© + ã® + åè©ï¼ˆã€œã®äº¤æ›ï¼äº¤æ›ã€œ ç­‰ï¼‰ äº¤æ›æ¡ä»¶ äº¤æ›æ‰‹æ•°æ–™ äº¤æ›æ‰€ äº¤æ›æœŸ äº¤æ›åˆ¶åº¦ åè© + ã¨ + å‹•è©ï¼ˆã€œã¨äº¤æ›ã™ã‚‹ï¼äº¤æ›ã—åˆã† ç­‰ï¼‰ ç‰©ã¨äº¤æ›ã™ã‚‹ æ„è¦‹ã¨äº¤æ›ã—åˆã† æ›¸ç±ã¨äº¤æ›ã™ã‚‹ æƒ…å ±ã¨äº¤æ›ã™ã‚‹ è´ˆã‚Šç‰©ã¨äº¤æ›ã™ã‚‹ å½¢å®¹è© + åè©ï¼ˆäº¤æ›ã‚’ä¿®é£¾ã™ã‚‹å½¢å®¹èªã¨ã®çµ„ã¿åˆã‚ã›ï¼‰ ç›¸äº’äº¤æ› ç„¡å„Ÿäº¤æ› ç­‰ä¾¡äº¤æ› å³æ™‚äº¤æ› è‡ªç”±äº¤æ› å‹•è© + åè©ï¼ˆäº¤æ›ã™ã‚‹å¯¾è±¡ã‚’è¡¨ã™ï¼‰ äº¤æ›æ“ä½œ äº¤æ›æ¥­å‹™ äº¤æ›å–å¼• äº¤æ›å¯¾å¿œ äº¤æ›å‡¦ç† å‰¯è©ï¼è£œåŠ©èª + å‹•è©ï¼ˆäº¤æ›ã®ä»•æ–¹ãƒ»é »åº¦ã‚’è¡¨ã™ï¼‰ å®šæœŸçš„ã«äº¤æ›ã™ã‚‹ ç›¸äº’ã«äº¤æ›ã™ã‚‹ ç›´æ¥äº¤æ›ã™ã‚‹ éƒ¨åˆ†çš„ã«äº¤æ›ã™ã‚‹ é€Ÿã‚„ã‹ã«äº¤æ›ã™ã‚‹",
    # "ä»£é‡‘": "åè© + ã‚’ + å‹•è©ï¼ˆã€Œã€œä»£é‡‘ã‚’ã€œã™ã‚‹ã€å‹ï¼‰ ä»£é‡‘ã‚’æ”¯æ‰•ã† ä»£é‡‘ã‚’å—ã‘å–ã‚‹ ä»£é‡‘ã‚’è«‹æ±‚ã™ã‚‹ ä»£é‡‘ã‚’ç«‹ã¦æ›¿ãˆã‚‹ ä»£é‡‘ã‚’å›åã™ã‚‹ åè© + ã® + åè© æ”¯æ‰•ä»£é‡‘ è²©å£²ä»£é‡‘ å•†å“ä»£é‡‘ è«‹æ±‚ä»£é‡‘ å¾Œæ‰•ã„ä»£é‡‘ åè© + ã¨ + å‹•è©ï¼ˆäº¤æ›ãƒ»å¼•æ›ãªã©ï¼‰ å•†å“ã¨ä»£é‡‘ã‚’äº¤æ›ã™ã‚‹ ä»£é‡‘ã¨å¼•ãæ›ãˆã«æ¸¡ã™ ç‰©ã¨ä»£é‡‘ã‚’æ›ãˆã‚‹ ä»£é‡‘ã¨ç›¸æ®ºã™ã‚‹ ä»£é‡‘ã¨å¯¾è±¡ã‚’å¯¾å¿œã•ã›ã‚‹ å½¢å®¹è© + åè© å‰æ‰•ã„ä»£é‡‘ å¾Œæ‰•ã„ä»£é‡‘ ä¸€æ‹¬ä»£é‡‘ åˆ†å‰²ä»£é‡‘ é ã‚Šä»£é‡‘ å‰¯è©ï¼è£œåŠ©èª + åè©ï¼å‹•è© ã‚ã‚‰ã‹ã˜ã‚ä»£é‡‘ å…¨é¡ä»£é‡‘ ä¸€éƒ¨ä»£é‡‘ æ­£å¼ã«ä»£é‡‘ã‚’æ”¯æ‰•ã† ãŸã ã¡ã«ä»£é‡‘ã‚’æŒ¯è¾¼ã‚€",
    # "ä»²é–“": "åè© + ã‚’ + å‹•è©ï¼ˆã€Œã€œã‚’ä»²é–“ã«åŠ ãˆã‚‹ï¼ã€œã‚’ä»²é–“ã¨ã™ã‚‹ã€å‹ï¼‰ ä»²é–“ã‚’å¢—ã‚„ã™ ä»²é–“ã‚’å‹Ÿã‚‹ ä»²é–“ã‚’é›†ã‚ã‚‹ ä»²é–“ã‚’ã¤ãã‚‹ ä»²é–“ã‚’ä¿¡é ¼ã™ã‚‹ åè© + ã¨ + å‹•è©ï¼ˆã€Œã€œã¨ä»²é–“ã«ãªã‚‹ï¼ä»²é–“ã¨å…±ã«ã€œã™ã‚‹ã€å‹ï¼‰ å‹ã¨ä»²é–“ã«ãªã‚‹ ä»²é–“ã¨å”åŠ›ã™ã‚‹ ä»²é–“ã¨éã”ã™ ä»²é–“ã¨é€£æºã™ã‚‹ ä»²é–“ã¨åˆ†ã‹ã¡åˆã† åè© + ã® + åè©ï¼ˆã€Œã€œã®ä»²é–“ã€ã€Œä»²é–“ã®ã€œã€å‹ï¼‰ ä»²é–“ã®ä¸€äºº ä»²é–“æ„è­˜ ä»²é–“é–¢ä¿‚ ä»²é–“å¤–ã‚Œ ä»²é–“å…¥ã‚Š å½¢å®¹è© + åè©ï¼ˆä»²é–“ã‚’ä¿®é£¾ã™ã‚‹è¡¨ç¾ï¼‰ è¦ªã—ã„ä»²é–“ æœ¬å½“ã®ä»²é–“ å¿ƒå¼·ã„ä»²é–“ å¤šæ§˜ãªä»²é–“ æ–°ã—ã„ä»²é–“ å‰¯è©ï¼è£œåŠ©èª + å‹•è©ãƒ»åè©ï¼ˆä»²é–“ã¨å«ã‚ã‚‹è¡¨ç¾ï¼‰ æ°—è»½ã«ä»²é–“ã¨ ä¸€ç·’ã«ä»²é–“ã¨ ã‚‚ã¯ã‚„ä»²é–“ã§ã¯ãªã„ ä»²é–“ã¨ã—ã¦èªã‚ã‚‹ ä»²é–“ã§ã‚ã‚‹",
    # "ä»»ã›ã‚‹": "åè© + ã‚’ + åè© + ã« + ä»»ã›ã‚‹ï¼ˆã€œã‚’ã€œã«ä»»ã›ã‚‹å‹ï¼‰ ä»•äº‹ã‚’å…ˆè¼©ã«ä»»ã›ã‚‹ è²¬ä»»ã‚’éƒ¨ä¸‹ã«ä»»ã›ã‚‹ åˆ¤æ–­ã‚’ä»–äººã«ä»»ã›ã‚‹ æ±ºå®šã‚’å¹¹éƒ¨ã«ä»»ã›ã‚‹ ç®¡ç†ã‚’å°‚é–€æ¥­è€…ã«ä»»ã›ã‚‹ åè© + ã‚’ + å‹•è©å¥ + ã« + ä»»ã›ã‚‹ï¼ˆã€œã‚’ã€œã™ã‚‹ã“ã¨ã‚’ä»»ã›ã‚‹å‹ï¼‰ æƒé™¤ã‚’å…¨é¢çš„ã«ä»»ã›ã‚‹ äº¤æ¸‰ã‚’ä¸¸ã”ã¨ä»»ã›ã‚‹ ä¼ç”»ã‚’ã™ã¹ã¦ä»»ã›ã‚‹ ç‰‡ä»˜ã‘ã‚’å¾Œä»»è€…ã«ä»»ã›ã‚‹ æ‰‹é…ã‚’ç¾åœ°ã‚¹ã‚¿ãƒƒãƒ•ã«ä»»ã›ã‚‹ åè© + ã« + ä»»ã›ã‚‹ï¼ˆä¸»ä½“ã‚’ç¤ºã™è¡¨ç¾ï¼‰ éƒ¨ä¸‹ã«ä»»ã›ã‚‹ å°‚é–€å®¶ã«ä»»ã›ã‚‹ å®¶æ—ã«ä»»ã›ã‚‹ å‹äººã«ä»»ã›ã‚‹ ãƒ—ãƒ­ã«ä»»ã›ã‚‹ è£œåŠ©è¡¨ç¾ãƒ»å‰¯è© + ä»»ã›ã‚‹ï¼ˆä»»ã›ã‚‹ã®ä»•æ–¹ãƒ»åº¦åˆã„ã‚’è¡¨ã™ï¼‰ å®Œå…¨ã«ä»»ã›ã‚‹ æ€ã„åˆ‡ã£ã¦ä»»ã›ã‚‹ å®‰å¿ƒã—ã¦ä»»ã›ã‚‹ ã™ã¹ã¦ä»»ã›ã‚‹ ã»ã¨ã‚“ã©ä»»ã›ã‚‹ æ…£ç”¨å½¢ãƒ»å®šå‹è¡¨ç¾ ä»»ã›ã¦ãã ã•ã„ ä»»ã›ã£ãã‚Šã«ã™ã‚‹ï¼ä»»ã›ã£ãã‚Šã  ä»»ã›ã‚‹ã—ã‹ãªã„ ï½ã«ä»»ã›ã¦ãŠã ï½ã«ä»»ã•ã‚Œã‚‹",
    # "ä¼‘æš‡": "åè© + ã‚’ + å‹•è©ï¼ˆã€Œã€œã‚’å–ã‚‹ï¼ã€œã‚’ç”³è«‹ã™ã‚‹ã€å‹ï¼‰ ä¼‘æš‡ã‚’å–ã‚‹ ä¼‘æš‡ã‚’ç”³è«‹ã™ã‚‹ ä¼‘æš‡ã‚’å–å¾—ã™ã‚‹ ä¼‘æš‡ã‚’å»¶é•·ã™ã‚‹ ä¼‘æš‡ã‚’ä¸ãˆã‚‹ åè© + ä¸­ï¼é–“ï¼ˆçŠ¶æ…‹ãƒ»æœŸé–“ã‚’è¡¨ã™ï¼‰ ä¼‘æš‡ä¸­ ä¼‘æš‡æœŸé–“ä¸­ ä¼‘æš‡ã®é–“ é•·æœŸä¼‘æš‡ä¸­ åŠ´åƒä¼‘æš‡ä¸­ åè© + ã® + åè©ï¼ˆä¼‘æš‡ã‚’ä¿®é£¾ã™ã‚‹èªã¨ã®çµ„ã¿åˆã‚ã›ï¼‰ æœ‰çµ¦ä¼‘æš‡ å¹´æ¬¡ä¼‘æš‡ ç—…æ°—ä¼‘æš‡ ç”£ä¼‘ãƒ»è‚²ä¼‘ï¼ˆç”£å‰ãƒ»è‚²å…ä¼‘æš‡ï¼‰ ç‰¹åˆ¥ä¼‘æš‡ åè© + ã« + å‹•è©ï¼ˆç›®çš„ãƒ»ç†ç”±ã‚’å°ãæ§‹æ–‡ãªã©ï¼‰ ä¼‘æš‡ã«ã‚ã¦ã‚‹ ä¼‘æš‡ã«å½“ã¦ã‚‹ ä¼‘æš‡ã«å‡ºã‚‹ ä¼‘æš‡ã«å……ã¦ã‚‹ ä¼‘æš‡ã«å…¥ã‚‹ å½¢å®¹è©ï¼ä¿®é£¾èª + åè©ï¼ˆä¼‘æš‡ã‚’ä¿®é£¾ã™ã‚‹èªï¼‰ é•·æœŸä¼‘æš‡ çŸ­æœŸä¼‘æš‡ é€£ä¼‘ä¼‘æš‡ ç„¡çµ¦ä¼‘æš‡ å®Œå…¨ä¼‘æš‡",
    # "ä½œæ¥­": "åè© + ã‚’ + å‹•è©ï¼ˆã€œã‚’ä½œæ¥­ã™ã‚‹ï¼ã€œã‚’é€²ã‚ã‚‹ ç­‰ï¼‰ ä½œæ¥­ã‚’é€²ã‚ã‚‹ ä½œæ¥­ã‚’è¡Œã† ä½œæ¥­ã‚’é–‹å§‹ã™ã‚‹ ä½œæ¥­ã‚’å®Œäº†ã™ã‚‹ ä½œæ¥­ã‚’ä¸­æ–­ã™ã‚‹ åè© + ä¸­ï¼é–“ï¼ˆé€²è¡Œä¸­ãƒ»æœŸé–“ã‚’è¡¨ã™ï¼‰ ä½œæ¥­ä¸­ ä½œæ¥­ã®é€”ä¸­ ä½œæ¥­æœŸé–“ä¸­ ä½œæ¥­ã®åˆé–“ å¤œé–“ä½œæ¥­ä¸­ åè© + ã® + åè©ï¼ˆä½œæ¥­ã‚’ä¿®é£¾ã™ã‚‹èªã¨ã®çµ„ã¿åˆã‚ã›ï¼‰ ä½œæ¥­å·¥ç¨‹ ä½œæ¥­æ‰‹é † ä½œæ¥­åŠ¹ç‡ ä½œæ¥­ç’°å¢ƒ ä½œæ¥­è² è· å‰¯è©ï¼ä¿®é£¾èª + å‹•è©ï¼ˆä½œæ¥­ã™ã‚‹ã®ã—ã‹ãŸãƒ»ç¨‹åº¦ã‚’è¡¨ã™èªï¼‰ ä¸å¯§ã«ä½œæ¥­ã™ã‚‹ æ€¥ã„ã§ä½œæ¥­ã™ã‚‹ æ­£ç¢ºã«ä½œæ¥­ã™ã‚‹ æ‰‹æ—©ãä½œæ¥­ã™ã‚‹ å¾ã€…ã«ä½œæ¥­ã‚’é€²ã‚ã‚‹ åè© + ã« + å‹•è©ï¼ˆä½œæ¥­ã«é–¢ä¸ã™ã‚‹ãƒ»ä½œæ¥­å¯¾è±¡ã‚’è¡¨ã™è¡¨ç¾ï¼‰ ä½œæ¥­ã«æºã‚ã‚‹ ä½œæ¥­ã«å¾“äº‹ã™ã‚‹ ä½œæ¥­ã«å‚åŠ ã™ã‚‹ ä½œæ¥­ã«å–ã‚Šã‹ã‹ã‚‹ ä½œæ¥­ã«ç€æ‰‹ã™ã‚‹",
    # "ä½œç‰©": "åè© + ã‚’ + å‹•è©ï¼ˆï½ã‚’ä½œç‰©ã¨ã™ã‚‹ï¼ï½ã‚’è‚²ã¦ã‚‹ ç­‰ï¼‰  ä½œç‰©ã‚’è‚²ã¦ã‚‹  ä½œç‰©ã‚’åç©«ã™ã‚‹  ä½œç‰©ã‚’æ ½åŸ¹ã™ã‚‹  ä½œç‰©ã‚’æ¤ãˆã‚‹  ä½œç‰©ã‚’ä¿è­·ã™ã‚‹  åè© + ã® + åè©ï¼ˆï½ã®ä½œç‰©ï¼ä½œç‰©ã®ï½ï¼‰  è¾²ä½œç‰©ï¼ˆã®ä½œç‰©ï¼‰  ä½œç‰©ã®åé‡  ä½œç‰©ã®å“è³ª  ä½œç‰©ã®å‡ºè·  ä½œç‰©ã®ç”Ÿè‚²  å½¢å®¹è©ï¼ä¿®é£¾èª + åè©ï¼ˆä½œç‰©ã‚’ä¿®é£¾ã™ã‚‹èªï¼‰  é£Ÿç”¨ä½œç‰©  å·¥èŠ¸ä½œç‰©  è¾²æ¥­ä½œç‰©  ä¸»ä½œç‰©  å‰¯ä½œç‰©  åè© + ã« + å‹•è©ï¼ˆä½œç‰©ã«é–¢ã‚ã‚‹å‹•è©å¥ï¼‰  ä½œç‰©ã«è¢«å®³ã‚’ä¸ãˆã‚‹  ä½œç‰©ã«å½±éŸ¿ã‚’åŠã¼ã™  ä½œç‰©ã«é©ã™ã‚‹  ä½œç‰©ã«å¯¾å¿œã™ã‚‹  ä½œç‰©ã«æ „é¤Šã‚’ä¸ãˆã‚‹  å‰¯è©ï¼è£œåŠ©èª + å‹•è©ãƒ»åè©ï¼ˆé »åº¦ãƒ»ç¨‹åº¦ãƒ»çŠ¶æ…‹ã‚’è¡¨ã™ï¼‰  å¤šãã®ä½œç‰©  ã‚ˆã„ä½œç‰©  ååˆ†ãªä½œç‰©  è±Šä½œã®ä½œç‰©  ä½œç‰©å…¨ä½“",
    # "ä¿¡ç”¨": "åè© + ã‚’ + å‹•è©ï¼ˆã€Œã€œã‚’ä¿¡ç”¨ã™ã‚‹ï¼ã€œã‚’å¾—ã‚‹ã€å‹ï¼‰  äººã‚’ä¿¡ç”¨ã™ã‚‹  ç›¸æ‰‹ã‚’ä¿¡ç”¨ã™ã‚‹  è‡ªåˆ†ã‚’ä¿¡ç”¨ã™ã‚‹  å ±å‘Šã‚’ä¿¡ç”¨ã™ã‚‹  è¨¼è¨€ã‚’ä¿¡ç”¨ã™ã‚‹  åè© + ã® + åè©ï¼ˆä¿¡ç”¨ã‚’ä¿®é£¾ã™ã‚‹èªã€ä¿¡ç”¨ã®ç¨®é¡ãªã©ï¼‰  ä¿¡ç”¨åº¦  ä¿¡ç”¨èª¿æŸ»  ä¿¡ç”¨æ ¼ä»˜ã‘  ä¿¡ç”¨å–å¼•  ä¿¡ç”¨ä¿è¨¼  åè© + ã¨ + å‹•è©ï¼ˆä¿¡ç”¨ã‚’å‰æã¨ã™ã‚‹æ§‹æ–‡ï¼‰  ä¿¡ç”¨ã¨è£åˆ‡ã‚Š  ä¿¡ç”¨ã¨ä¿¡é ¼ã‚’çµã¶  ä¿¡ç”¨ã¨æœŸå¾…ã‚’å¯„ã›ã‚‹  ä¿¡ç”¨ã¨å®Ÿç¸¾ã‚’é‡ã­ã‚‹  ä¿¡ç”¨ã¨ç¾©å‹™ã‚’æœãŸã™  å½¢å®¹è©ï¼ä¿®é£¾èª + åè©ï¼ˆä¿¡ç”¨ã‚’å¼·èª¿ãƒ»æ¡ä»¶ã¥ã‘ã‚‹èªï¼‰  é«˜ã„ä¿¡ç”¨  å¼·ã„ä¿¡ç”¨  çµ¶å¯¾ã®ä¿¡ç”¨  ç¤¾ä¼šçš„ä¿¡ç”¨  é•·æœŸä¿¡ç”¨  å‰¯è©ï¼è£œåŠ©èª + å‹•è©ï¼ˆä¿¡ç”¨ã™ã‚‹ã®ã—ã‹ãŸãƒ»ç¨‹åº¦ã‚’è¡¨ã™èªï¼‰  å®‰æ˜“ã«ä¿¡ç”¨ã™ã‚‹  å®‰å¿ƒã—ã¦ä¿¡ç”¨ã™ã‚‹  å®Œå…¨ã«ä¿¡ç”¨ã™ã‚‹  åŠã°ä¿¡ç”¨ã™ã‚‹  ååˆ†ã«ä¿¡ç”¨ã™ã‚‹  å‹•è©å¥ + åè©ï¼ˆä¿¡ç”¨ã‚’å¯¾è±¡ã¨ã™ã‚‹å‹•è©ã¨ã®çµ„ã¿åˆã‚ã›ï¼‰  ä¿¡ç”¨ã‚’å¤±ã†  ä¿¡ç”¨ã‚’å›å¾©ã™ã‚‹  ä¿¡ç”¨ã‚’ç¯‰ã  ä¿¡ç”¨ã‚’æãªã†  ä¿¡ç”¨ã‚’é«˜ã‚ã‚‹",
    # "ä¿®æ­£": "åè© + ã‚’ + å‹•è©ï¼ˆã€œã‚’ä¿®æ­£ã™ã‚‹å‹ï¼‰  èª¤ã‚Šã‚’ä¿®æ­£ã™ã‚‹  æ–‡ç« ã‚’ä¿®æ­£ã™ã‚‹  è¨ˆç”»ã‚’ä¿®æ­£ã™ã‚‹  ãƒ‡ã‚¶ã‚¤ãƒ³ã‚’ä¿®æ­£ã™ã‚‹  ã‚¹ã‚±ã‚¸ãƒ¥ãƒ¼ãƒ«ã‚’ä¿®æ­£ã™ã‚‹  åè© + ã® + åè©ï¼ˆä¿®æ­£ã‚’ä¿®é£¾ï¼ä¿®æ­£å¯¾è±¡ã‚’ç¤ºã™èªï¼‰  ä¿®æ­£æ¡ˆ  ä¿®æ­£ç‚¹  ä¿®æ­£å¾Œ  ä¿®æ­£å‰  ä¿®æ­£ä½œæ¥­  è£œåŠ©èªï¼å‰¯è© + å‹•è©ï¼ˆä¿®æ­£ã™ã‚‹ã®ã—ã‹ãŸãƒ»åº¦åˆã„ã‚’è¡¨ã™èªï¼‰  å°‘ã—ä¿®æ­£ã™ã‚‹  è»½ãä¿®æ­£ã™ã‚‹  å¤§å¹…ã«ä¿®æ­£ã™ã‚‹  å†åº¦ä¿®æ­£ã™ã‚‹  æœ€çµ‚ä¿®æ­£ã™ã‚‹  åè© + ã« + å‹•è©ï¼ˆä¿®æ­£ã‚’åŠ ãˆã‚‹å¯¾è±¡ãƒ»ç›®çš„ã‚’ç¤ºã™æ§‹æ–‡ï¼‰  ï½ã«ä¿®æ­£ã‚’åŠ ãˆã‚‹  ä»•æ§˜ã«ä¿®æ­£ã‚’åŠ ãˆã‚‹  æ¡æ–‡ã«ä¿®æ­£ã‚’åŠ ãˆã‚‹  è¡¨ç¾ã«ä¿®æ­£ã‚’åŠ ãˆã‚‹  å†…å®¹ã«ä¿®æ­£ã‚’åŠ ãˆã‚‹  æ…£ç”¨è¡¨ç¾ãƒ»è¤‡åˆèª  è»Œé“ä¿®æ­£  ä¿®æ­£ãƒãƒ¼ã‚¸ãƒ§ãƒ³  ä¿®æ­£å¯¾å¿œ  ä¿®æ­£ãƒ•ã‚£ãƒ¼ãƒ‰ãƒãƒƒã‚¯  ä¿®æ­£å¯èƒ½",
    # "å€¤": "åè© + ã‚’ + å‹•è©ï¼ˆã€Œã€œã‚’ã€œã™ã‚‹ï¼ã€œã‚’ã€œã¨ã™ã‚‹ã€å‹ï¼‰  å€¤ã‚’ä»˜ã‘ã‚‹  å€¤ã‚’ä¸‹ã’ã‚‹  å€¤ã‚’ä¸Šã’ã‚‹  å€¤ã‚’è¦‹ç›´ã™  å€¤ã‚’è¨­å®šã™ã‚‹  åè© + ã® + åè©ï¼ˆå€¤ã‚’ä¿®é£¾ã™ã‚‹èªãƒ»å€¤ã®ç¨®é¡ãªã©ï¼‰  å€¤æ®µ  å€¤å¹…  å€¤ä¸Šã’  å€¤ä¸‹ã’  æœ€ä½å€¤ï¼æœ€é«˜å€¤  åè© + ãŒ + å‹•è©ï¼ˆå€¤ãŒã€œã™ã‚‹ãƒ»å€¤ãŒã€œã ã¨ã„ã†è¡¨ç¾ï¼‰  å€¤ãŒä¸ŠãŒã‚‹  å€¤ãŒä¸‹ãŒã‚‹  å€¤ãŒå¤‰å‹•ã™ã‚‹  å€¤ãŒå®‰å®šã™ã‚‹  å€¤ãŒé«˜ã„ï¼ä½ã„  å‹•è© + åè©ï¼ˆå‹•è©ã¨çµ„ã‚€åè©ã§å€¤ã‚’è¡¨ã™è¡¨ç¾ï¼‰  å€¤ã‚’è¡¨ç¤ºã™ã‚‹  å€¤ã‚’æ¯”è¼ƒã™ã‚‹  å€¤ã‚’è¨˜éŒ²ã™ã‚‹  å€¤ã‚’å ±å‘Šã™ã‚‹  å€¤ã‚’äºˆæ¸¬ã™ã‚‹  è£œåŠ©èªãƒ»å‰¯è© + å‹•è©ï¼ˆå€¤ã‚’æ‰±ã†å‹•è©è¡¨ç¾ã‚’ä¿®é£¾ã™ã‚‹èªï¼‰  é©æ­£ã«å€¤ã‚’ä»˜ã‘ã‚‹  å¤§å¹…ã«å€¤ã‚’ä¸‹ã’ã‚‹  æ€¥æ¿€ã«å€¤ãŒä¸ŠãŒã‚‹  ã‚ãšã‹ã«å€¤ã‚’å¤‰ãˆã‚‹  æ¯”è¼ƒçš„é«˜ã„å€¤",
    # "åƒ…ã‹": "å½¢å®¹å‹•è© + åè© â€œã‚ãšã‹ãª + åè©â€ ã®å½¢ã§ã€é‡ãƒ»ç¨‹åº¦ãƒ»æ™‚é–“ãªã©ã®å°‘ãªã•ã‚’è¡¨ã™ åƒ…ã‹ãªæ™‚é–“ åƒ…ã‹ãªå·® åƒ…ã‹ãªåˆ©ç›Š åƒ…ã‹ãªæœŸå¾… åƒ…ã‹ãªå¯èƒ½æ€§ å‰¯è©ï¼ˆåƒ…ã‹ã«ï¼‰ + å‹•è©ãƒ»å½¢å®¹è© â€œã‚ãšã‹ã« + ï½ã™ã‚‹ï¼ï½ã§ã‚ã‚‹â€ ã®å½¢ã§ã€ã»ã©ã‚ˆãï¼ã‚ãšã‹ã«å¤‰åŒ–ãƒ»çŠ¶æ…‹ã‚’ç¤ºã™ åƒ…ã‹ã«è¦‹ãˆã‚‹ åƒ…ã‹ã«ãšã‚Œã‚‹ åƒ…ã‹ã«é•ã† åƒ…ã‹ã«å‚¾ã åƒ…ã‹ã«èã“ãˆã‚‹ åè© + æ ¼åŠ©è© + åƒ…ã‹ åè©ã‚’ä¸»èªã‚„å¯¾è±¡ã¨ã—ã€ã€Œã€œãŒåƒ…ã‹ï¼ˆã§ã‚ã‚‹ï¼ã—ã‹ãªã„ï¼‰ã€ã®ã‚ˆã†ãªè¡¨ç¾ æ®‹ã‚ŠãŒåƒ…ã‹ å·®ãŒåƒ…ã‹ ä½™è£•ãŒåƒ…ã‹ åœ¨åº«ãŒåƒ…ã‹ å¤‰åŒ–ãŒåƒ…ã‹ æ•°è© + åƒ…ã‹ æ•°ã‚’é™å®šã—ã¦ã€ŒãŸã£ãŸâ€¦ã ã‘ã€ã®æ„å‘³åˆã„ã‚’å¼·ã‚ã‚‹å¥ åƒ…ã‹1äºº åƒ…ã‹æ•°åˆ† åƒ…ã‹æ•°æ—¥ åƒ…ã‹1å› åƒ…ã‹æ•°å€‹ åƒ…ã‹ + åŠ©è©ï¼ˆã€œã«ï¼ã€œã§ï¼ã€œã—ã‹ï¼‰ â€œåƒ…ã‹ã«â€ ã‚„ â€œåƒ…ã‹ã§ï¼ã—ã‹â€ ã®å½¢ã§ä¿®é£¾èªçš„ã«ç”¨ã„ã‚‹ åƒ…ã‹ã«å¤‰åŒ–ã™ã‚‹ åƒ…ã‹ã«é•ã† åƒ…ã‹ã§ã‚ã‚‹ åƒ…ã‹ã§ã—ã‹ãªã„",
    # "å…¸å‹": "åè© + ã® + å½¢å®¹è©ï¼åè© ï¼ˆã€Œå…¸å‹ã®ï½ã€ã¨ã„ã†è¡¨ç¾ï¼‰ å…¸å‹ã®ä¾‹ å…¸å‹ã®ãƒ‘ã‚¿ãƒ¼ãƒ³ å…¸å‹ã®äº‹ä¾‹ å…¸å‹ã®ç‰¹å¾´ å…¸å‹ã®å§¿ ï½ã¨ã—ã¦ + å…¸å‹ï¼ˆçš„ï¼‰ + åè© ï¼ˆä½•ã‹ã‚’å…¸å‹ä¾‹ã¨ã—ã¦ä½ç½®ã¥ã‘ã‚‹è¡¨ç¾ï¼‰ ï½ã¨ã—ã¦å…¸å‹ä¾‹ ï½ã¨ã—ã¦å…¸å‹åƒ ï½ã¨ã—ã¦å…¸å‹ãƒ¢ãƒ‡ãƒ« ï½ã¨ã—ã¦å…¸å‹çš„äº‹ä¾‹ ï½ã¨ã—ã¦å…¸å‹çš„ãƒ‘ã‚¿ãƒ¼ãƒ³ å½¢å®¹è© + å…¸å‹ï¼ˆä¾‹ãƒ»å§¿ãƒ»åƒãªã©ï¼‰ ï¼ˆå…¸å‹ã‚’ä¿®é£¾ã™ã‚‹è¨€è‘‰ä»˜ãè¡¨ç¾ï¼‰ æœ€ã‚‚å…¸å‹çš„ãªä¾‹ å¤å…¸çš„ã§å…¸å‹çš„ãªãƒ‘ã‚¿ãƒ¼ãƒ³ ã‚ã‹ã‚Šã‚„ã™ã„å…¸å‹åƒ æ¨™æº–çš„ãªå…¸å‹ä¾‹ ä¸€èˆ¬çš„ãªå…¸å‹åƒ ï½ã¯ + å…¸å‹ + ã  / ï½ãŒ + å…¸å‹ + ã  ï¼ˆè¿°èªã¨ã—ã¦ä½¿ã†è¡¨ç¾ï¼‰ ã“ã‚Œã¯å…¸å‹ã  å½¼ï¼å½¼å¥³ã¯å…¸å‹ã  ãã®æ„è¦‹ã¯å…¸å‹ã  ã“ã®ç¾è±¡ãŒå…¸å‹ã  ãã®ã‚±ãƒ¼ã‚¹ãŒå…¸å‹ã  å‹•è© + å…¸å‹ï¼ˆä¾‹ãƒ»å§¿ãªã©ï¼‰ ï¼ˆå‹•ä½œã¨ã¨ã‚‚ã«ã€Œå…¸å‹ã€ã‚’ä½¿ã†è¡¨ç¾ï¼‰ ç¤ºã™å…¸å‹ä¾‹ è¦‹ã›ã‚‹å…¸å‹åƒ æˆã‚‹ï¼ãªã‚‹å…¸å‹ä¾‹ æ¡ã‚‹å…¸å‹ä¾‹ å½¢æˆã™ã‚‹å…¸å‹åƒ",
    # "å†’é™º": "åè© + ã® + åè© ï¼ˆã€Œå†’é™ºã®ï½ã€ã¨ã„ã†è¡¨ç¾ï¼‰ å†’é™ºã®æ—… å†’é™ºã®å§‹ã¾ã‚Š å†’é™ºã®ç‰©èª å†’é™ºã®èˆå° å†’é™ºã®è¨˜éŒ² åè© + ã‚’ + å‹•è© ï¼ˆã€Œå†’é™ºã‚’ï½ã™ã‚‹ã€ãªã©ï¼‰ å†’é™ºã‚’ã™ã‚‹ å†’é™ºã‚’å§‹ã‚ã‚‹ å†’é™ºã‚’æ±‚ã‚ã‚‹ å†’é™ºã‚’é‡ã­ã‚‹ å†’é™ºã‚’æã‚Œã‚‹ ï½ã« + å†’é™ºï¼ˆã™ã‚‹ãƒ»å‡ºã‚‹ï¼‰ ï¼ˆç›®çš„åœ°ã‚„å ´é¢ã‚’ç¤ºã™ï¼‰ ï½ã«å†’é™ºã™ã‚‹ ï½ã«å†’é™ºã«å‡ºã‚‹ ï½ã«å†’é™ºå¿ƒã‚’æŠ±ã ï½ã«å†’é™ºã‚’æŒ‘ã‚€ ï½ã«å†’é™ºã‚’æŒã¡ã‹ã‘ã‚‹ å½¢å®¹è© + å†’é™ºï¼å†’é™ºçš„ ï¼ˆå†’é™ºã‚’ä¿®é£¾ã™ã‚‹è¡¨ç¾ï¼‰ åˆºæ¿€çš„ãªå†’é™º å±é™ºãªå†’é™º æœªçŸ¥ã®å†’é™º å£®å¤§ãªå†’é™º å†’é™ºçš„ãªè¨ˆç”» å†’é™º + åè© ï¼ˆã€Œå†’é™ºï½ã€ã®è¤‡åˆèªçš„ãªä½¿ã„æ–¹ï¼‰ å†’é™ºå°èª¬ å†’é™ºè­š å‡é™ºè€…ï¼ˆå†’é™ºè€…ï¼‰ å†’é™ºæ—…è¡Œ å†’é™ºå¿ƒ",
    # "å‡ºå¸­": "åè© + ã® + å‡ºå¸­ æ­£å¼ã®å‡ºå¸­ å…¨å“¡ã®å‡ºå¸­ æ¥è³“ã®å‡ºå¸­ å¿…è¦ãªå‡ºå¸­ è‡¨æ™‚ã®å‡ºå¸­ å‡ºå¸­ + ã‚’ + å‹•è© å‡ºå¸­ã‚’æ±‚ã‚ã‚‹ å‡ºå¸­ã‚’å–ã‚‹ å‡ºå¸­ã‚’ç¢ºèªã™ã‚‹ å‡ºå¸­ã‚’ä¿ƒã™ å‡ºå¸­ã‚’ä¿éšœã™ã‚‹ ï½ã« + å‡ºå¸­ï¼ˆã™ã‚‹ï¼ã§ãã‚‹ï¼è¦‹é€ã‚‹ï¼‰ ä¼šè­°ã«å‡ºå¸­ã™ã‚‹ æˆæ¥­ã«å‡ºå¸­ã™ã‚‹ å¼å…¸ã«å‡ºå¸­ã§ãã‚‹ ç·ä¼šã«å‡ºå¸­ã™ã‚‹ æˆæ¥­ã‚’å‡ºå¸­ã™ã‚‹ï¼ˆå°‘ã—é‡è¤‡ã ãŒæ…£ç”¨ï¼‰ å½¢å®¹è© + å‡ºå¸­ æ¬ å¸­ã¨å‡ºå¸­ å…¨å“¡å‡ºå¸­ å¤šæ•°å‡ºå¸­ æ­£å¼å‡ºå¸­ åèª‰å‡ºå¸­ å‡ºå¸­ + åè© å‡ºå¸­ç‡ å‡ºå¸­è€… å‡ºå¸­åç°¿ å‡ºå¸­è¨˜éŒ² å‡ºå¸­ç°¿",
    # "åˆ©ç›Š": "åè© + ã® + åè© åˆ©ç›Šã®è¿½æ±‚ åˆ©ç›Šã®ç¢ºä¿ åˆ©ç›Šã®é…åˆ† åˆ©ç›Šã®æœ€å¤§åŒ– åˆ©ç›Šã®é‚„å…ƒ åˆ©ç›Š + ã‚’ + å‹•è© åˆ©ç›Šã‚’å¾—ã‚‹ï¼ã‚ã’ã‚‹ åˆ©ç›Šã‚’è¿½æ±‚ã™ã‚‹ åˆ©ç›Šã‚’ç¢ºä¿ã™ã‚‹ åˆ©ç›Šã‚’åˆ†é…ã™ã‚‹ åˆ©ç›Šã‚’å†æŠ•è³‡ã™ã‚‹ ï½ã« + åˆ©ç›Šï¼ˆã™ã‚‹ï¼å¾—ã‚‹ï¼é‚„å…ƒã™ã‚‹ï¼‰ é¡§å®¢ã«åˆ©ç›Šã‚’é‚„å…ƒã™ã‚‹ å¸‚å ´ã«åˆ©ç›Šã‚’ã‚‚ãŸã‚‰ã™ ç¤¾ä¼šã«åˆ©ç›Šã‚’ã‚‚ãŸã‚‰ã™ æ ªä¸»ã«åˆ©ç›Šã‚’é‚„å…ƒã™ã‚‹ ä¼æ¥­ã«åˆ©ç›Šã‚’ã‚‚ãŸã‚‰ã™ å½¢å®¹è© + åˆ©ç›Šï¼åˆ©ç›Šçš„ ç´”ç²‹ãªåˆ©ç›Š å®Ÿè³ªçš„ãªåˆ©ç›Š è«å¤§ãªåˆ©ç›Š æ½œåœ¨çš„ãªåˆ©ç›Š çŸ­æœŸçš„ãªåˆ©ç›Š åˆ©ç›Š + åè© åˆ©ç›Šç‡ åˆ©ç›Šé¡ åˆ©ç›Šæ‹¡å¤§ åˆ©ç›Šæ§‹é€  åˆ©ç›Šåˆ†é…",
    # "åˆ¶åº¦": "åè© + ã® + åè© åˆ¶åº¦ã®å°å…¥ åˆ¶åº¦ã®æ”¹æ­£ åˆ¶åº¦ã®ç¢ºç«‹ åˆ¶åº¦ã®ç¶­æŒ åˆ¶åº¦ã®è¦‹ç›´ã— åˆ¶åº¦ + ã‚’ + å‹•è© åˆ¶åº¦ã‚’è¨­ã‘ã‚‹ åˆ¶åº¦ã‚’å°å…¥ã™ã‚‹ åˆ¶åº¦ã‚’æ”¹ã‚ã‚‹ï¼æ”¹æ­£ã™ã‚‹ åˆ¶åº¦ã‚’é©ç”¨ã™ã‚‹ åˆ¶åº¦ã‚’æ’¤å»ƒã™ã‚‹ ï½ã« + åˆ¶åº¦ï¼ˆã‚’ï¼ãŒï¼‰ å›½ã«åˆ¶åº¦ã‚’å°å…¥ã™ã‚‹ åŒæ§˜ã«åˆ¶åº¦ãŒå­˜åœ¨ã™ã‚‹ å¸‚å ´ã«åˆ¶åº¦ã‚’è¨­ã‘ã‚‹ ç¤¾ä¼šã«åˆ¶åº¦ã‚’æ ¹ä»˜ã‹ã›ã‚‹ çµ„ç¹”ã«åˆ¶åº¦ã‚’æ•´å‚™ã™ã‚‹ å½¢å®¹è© + åˆ¶åº¦ï¼åˆ¶åº¦çš„ å…¬çš„ãªåˆ¶åº¦ æ³•çš„ãªåˆ¶åº¦ æ–°ã—ã„åˆ¶åº¦ æ—¢å­˜ã®åˆ¶åº¦ æŸ”è»Ÿãªåˆ¶åº¦ åˆ¶åº¦ + åè© åˆ¶åº¦æ”¹é© åˆ¶åº¦è¨­è¨ˆ åˆ¶åº¦é‹ç”¨ åˆ¶åº¦æ”¹æ­£ åˆ¶åº¦åŸºæº–",
    # "åŠ æ¸›": "å‹•è© + åŠ æ¸›ï¼ˆã‚’ï¼ã™ã‚‹ï¼‰ åŠ æ¸›ã‚’ã™ã‚‹ åŠ æ¸›ã‚’èª¿æ•´ã™ã‚‹ åŠ æ¸›ã‚’è¦‹ã‚‹ åŠ æ¸›ã‚’è€ƒãˆã‚‹ åŠ æ¸›ã‚’ã¤ã‘ã‚‹ åè© + ã® + åŠ æ¸› å¡©ã®åŠ æ¸› ç«ã®åŠ æ¸› ç…§æ˜ã®åŠ æ¸› æ¹¯ã®åŠ æ¸› ç„¼ãåŠ æ¸› ï½å…·åˆ + åŠ æ¸›ï¼ï½ã‹ã’ã‚“ å…·åˆã®åŠ æ¸› ä½“ã®åŠ æ¸›ï¼ˆä¾‹ï¼šã€ŒãŠåŠ æ¸›ã¯ã„ã‹ãŒã§ã™ã‹ã€ï¼‰ ç—›ã¿ã®åŠ æ¸› æ°—åˆ†ã®åŠ æ¸› èª¿å­ã®åŠ æ¸› æ¥å°¾è¾çš„ç”¨æ³•ï¼ˆï½åŠ æ¸›ï¼‰ ï¼ˆå‹•è©ãƒ»çŠ¶æ…‹ã‚’è¡¨ã™èªã«ä»˜ã„ã¦ã€ã€Œï½ã®å…·åˆï¼ã‚„ã‚„ï½ãªæ§˜å­ã€ã‚’è¡¨ã™ï¼‰ ç„¼ãåŠ æ¸› ç„¼ã‘åŠ æ¸› æ›¸ãåŠ æ¸› æ­©ãåŠ æ¸› ã†ã¤ã‚€ãåŠ æ¸› å½¢å®¹è©ãƒ»å‰¯è© + åŠ æ¸› ã„ã„åŠ æ¸› é©åˆ‡ãªåŠ æ¸› ç¨‹ã‚ˆã„åŠ æ¸› éå¸¸ã«åŠ æ¸›ï¼ˆã‚ã¾ã‚Šä½¿ã‚ã‚Œãªã„ãŒæ–‡è„ˆã§ï¼‰ ã‚„ã‚„åŠ æ¸›",
    # "åŒ–ç²§": "åè© + ã® + åè© åŒ–ç²§ã®ä»•æ–¹ åŒ–ç²§ã®ç¨®é¡ åŒ–ç²§ã®åŠ¹æœ åŒ–ç²§ã®é“å…· åŒ–ç²§ã®æŠ€è¡“ åŒ–ç²§ + ã‚’ + å‹•è© åŒ–ç²§ã‚’ã™ã‚‹ åŒ–ç²§ã‚’è½ã¨ã™ åŒ–ç²§ã‚’ç›´ã™ åŒ–ç²§ã‚’æ¿ƒãã™ã‚‹ åŒ–ç²§ã‚’è–„ãã™ã‚‹ ï½ã« + åŒ–ç²§ï¼ˆã™ã‚‹ï¼ã•ã‚Œã‚‹ï¼‰ é¡”ã«åŒ–ç²§ã™ã‚‹ é ¬ã«åŒ–ç²§ã•ã‚Œã‚‹ ç›®å…ƒã«åŒ–ç²§ã‚’æ–½ã™ äººã«åŒ–ç²§ã‚’ã™ã‚‹ æ‹¡å¼µï¼šå£ã«åŒ–ç²§ã‚’æ–½ã™ å½¢å®¹è© + åŒ–ç²§ / åŒ–ç²§çš„ æ¿ƒã„åŒ–ç²§ è–„ã„åŒ–ç²§ ãã‚Œã„ãªåŒ–ç²§ è‡ªç„¶ãªåŒ–ç²§ æ´¾æ‰‹ãªåŒ–ç²§ åŒ–ç²§ + åè© åŒ–ç²§å“ åŒ–ç²§é“å…· åŒ–ç²§æ°´ åŒ–ç²§å®¤ åŒ–ç²§ç­†",
    # "åŒ»ç™‚": "åè© + ã® + åè© åŒ»ç™‚ã®ç¾å ´ åŒ»ç™‚ã®åˆ¶åº¦ åŒ»ç™‚ã®è²»ç”¨ åŒ»ç™‚ã®è³ª åŒ»ç™‚ã®æä¾› åŒ»ç™‚ + ã‚’ + å‹•è© åŒ»ç™‚ã‚’å—ã‘ã‚‹ åŒ»ç™‚ã‚’æä¾›ã™ã‚‹ åŒ»ç™‚ã‚’æ”¹å–„ã™ã‚‹ åŒ»ç™‚ã‚’æ”¯ãˆã‚‹ åŒ»ç™‚ã‚’æ”¹é©ã™ã‚‹ ï½ã« + åŒ»ç™‚ï¼ˆã‚’ï¼ãŒï¼‰ åœ°åŸŸã«åŒ»ç™‚ã‚’æä¾›ã™ã‚‹ å›½ã«åŒ»ç™‚åˆ¶åº¦ãŒã‚ã‚‹ ç¤¾ä¼šã«åŒ»ç™‚ã®æ ¼å·®ãŒã‚ã‚‹ é«˜é½¢è€…ã«åŒ»ç™‚ãŒå¿…è¦ã  ç½å®³æ™‚ã«åŒ»ç™‚ãŒä¸è¶³ã™ã‚‹ å½¢å®¹è© + åŒ»ç™‚ï¼åŒ»ç™‚çš„ é«˜åº¦ãªåŒ»ç™‚ å…ˆé€²çš„ãªåŒ»ç™‚ å…¬çš„ãªåŒ»ç™‚ ç·Šæ€¥åŒ»ç™‚ å¿…è¦ä¸å¯æ¬ ãªåŒ»ç™‚ åŒ»ç™‚ + åè© åŒ»ç™‚æ©Ÿé–¢ åŒ»ç™‚å¾“äº‹è€… åŒ»ç™‚è²» åŒ»ç™‚æŠ€è¡“ åŒ»ç™‚ã‚µãƒ¼ãƒ“ã‚¹",
    # "å„ä»‹": "å½¢å®¹è© + åè© å„ä»‹ãªå•é¡Œ å„ä»‹ãªç›¸æ‰‹ å„ä»‹ãªè©± å„ä»‹ãªçŠ¶æ³ å„ä»‹ãªäº‹æ…‹ åè© + ã‚’ + å‹•è© å„ä»‹ã‚’å¼•ãå—ã‘ã‚‹ å„ä»‹ã‚’ã‹ã‘ã‚‹ å„ä»‹ã‚’ã¿ã‚‹ å„ä»‹ã‚’ã•ã›ã‚‹ å„ä»‹ã‚’é¿ã‘ã‚‹ ï½ã« + å„ä»‹ï¼ˆãªã‚‹ï¼å·»ãè¾¼ã¾ã‚Œã‚‹ï¼ãªã‚‹ï¼‰ è¦ªã®å„ä»‹ã«ãªã‚‹ äººé–“é–¢ä¿‚ã«å„ä»‹ã«ãªã‚‹ å•é¡Œã«å„ä»‹ã«ãªã‚‹ ãƒˆãƒ©ãƒ–ãƒ«ã«å·»ãè¾¼ã¾ã‚Œã¦å„ä»‹ã«ãªã‚‹ ã‚ã‚‹äº‹æŸ„ã«å„ä»‹ã«ãªã‚‹ å‰¯è©ãƒ»æ¥ç¶šèª + å„ä»‹ï¼ˆã«ï¼ã§ï¼‰ éå¸¸ã«å„ä»‹ã«ãªã‚‹ ã¡ã‚‡ã£ã¨å„ä»‹ã§ã‚ã‚‹ ã¨ã¦ã‚‚å„ä»‹ã«æ€ã† ã‚„ã£ã‹ã„ã«ã—ã¦ã—ã¾ã† å„ä»‹ã§ãŸã¾ã‚‰ãªã„",
    # "åã‚ã‚‹": "å‹•è© + åã‚ã‚‹ï¼ç´ã‚ã‚‹ æˆåŠŸã‚’åã‚ã‚‹ åˆ©ç›Šã‚’åã‚ã‚‹ å­¦è²»ã‚’ç´ã‚ã‚‹ ç¨é‡‘ã‚’ç´ã‚ã‚‹ ç´›äº‰ã‚’åã‚ã‚‹ åè© + ã‚’ + åã‚ã‚‹ï¼ç´ã‚ã‚‹ é‹¼æã‚’å€‰åº«ã«åã‚ã‚‹ æ›¸ç‰©ã‚’ç›®éŒ²ã«åã‚ã‚‹ æ³¨æ–‡ã®å“ã‚’ç´ã‚ã‚‹ å£²ä¸Šã‚’å¸‚ã«ç´ã‚ã‚‹ æœ¬ã‚’æœ¬æ£šã«åã‚ã‚‹ åã‚ã‚‹ + åè© åã‚ã‚‹æˆæœ åã‚ã‚‹å‹åˆ© ç´ã‚ã‚‹é‡‘é¡ ç´ã‚ã‚‹ç¾©å‹™ åã‚ã‚‹åˆ©ç›Š ï½ã«ï¼ï½ã¸ + åã‚ã‚‹ï¼ç´ã‚ã‚‹ é¡§å®¢ã«åˆ©ç›Šã‚’åã‚ã‚‹ å›½ã¸ç¨é‡‘ã‚’ç´ã‚ã‚‹ æ›¸åº«ã«è³‡æ–™ã‚’åã‚ã‚‹ å·¥å ´ã«è£½å“ã‚’åã‚ã‚‹ æ¸‹æ»ã‚’åã‚ã‚‹ï¼ˆæš´å‹•ãªã©ã‚’ï¼‰ å½¢å®¹è©ãƒ»å‰¯è© + åã‚ã‚‹ï¼ç´ã‚ã‚‹ ç„¡äº‹ã«åã‚ã‚‹ ãã¡ã‚“ã¨ç´ã‚ã‚‹ æ—©ãåã‚ã‚‹ é©åˆ‡ã«ç´ã‚ã‚‹ ç¢ºå®Ÿã«åã‚ã‚‹",
    # "å±ã‚‹": "åè© + ã‚’ + å±ã‚‹ å­ã©ã‚‚ã‚’å±ã‚‹ éƒ¨ä¸‹ã‚’å±ã‚‹ ç”Ÿå¾’ã‚’å±ã‚‹ çŠ¬ã‚’å±ã‚‹ è¦ªã‚’å±ã‚‹ ï¼ˆãŸã ã—ä½¿ã‚ã‚Œã‚‹å ´é¢ã¯å°‘ãªã‚ï¼‰ ï½ã« + å±ã‚‹ï¼å±ã‚‰ã‚Œã‚‹ å­ã©ã‚‚ã«å±ã‚‹ èª¤ã‚Šã‚’ã—ãŸäººã«å±ã‚‹ æ•™å¸«ã«å±ã‚‰ã‚Œã‚‹ è¦ªã«å±ã‚‰ã‚Œã‚‹ å…ˆç”Ÿã«å±ã‚‰ã‚Œã‚‹ å±ã‚‹ + åè©ãƒ»è£œèª å±ã‚‹å£° å±ã‚‹ç†ç”± å±ã‚‹æ…‹åº¦ å±ã‚‹å£èª¿ å±ã‚‹å¿…è¦æ€§ å‰¯è©ãƒ»ä¿®é£¾èª + å±ã‚‹ å³ã—ãå±ã‚‹ ãã¤ãå±ã‚‹ å¼·ãå±ã‚‹ å„ªã—ãå±ã‚‹ æ™‚æŠ˜å±ã‚‹ å±ã‚‹ + å‹•è©ãƒ»å¥ å±ã£ã¦ã—ã¾ã† å±ã£ã¦ã‚‚ã‚‰ã† å±ã£ã¦ã‚„ã‚‹ å±ã‚‰ã‚Œã¦åçœã™ã‚‹ å±ã‚‰ã‚Œã¦æ°—ã¥ã",
    # "åˆå›³": "åè© + ã® + åˆå›³ å‡ºç™ºã®åˆå›³ é–‹å§‹ã®åˆå›³ çµ‚äº†ã®åˆå›³ åˆå›³ã®ã‚¿ã‚¤ãƒŸãƒ³ã‚° åˆå›³ã®æ„å‘³ åˆå›³ + ã‚’ + å‹•è© åˆå›³ã‚’å‡ºã™ åˆå›³ã‚’é€ã‚‹ åˆå›³ã‚’å¿˜ã‚Œã‚‹ åˆå›³ã‚’ã™ã‚‹ åˆå›³ã‚’è¦‹è½ã¨ã™ ï½ã® + åˆå›³ + ã«ï¼ˆãªã‚‹ï¼ãªã‚‹ã‚‚ã®ï¼‰ ãã‚Œã‚’åˆå›³ã«ã—ã¦å§‹ã‚ã‚‹ ãã®è¨€è‘‰ãŒåˆå›³ã«ãªã£ãŸ ç¬›ãŒåˆå›³ã«ãªã£ã¦ä¸€æ–‰ã«å‹•ã ç›®é…ã›ãŒåˆå›³ã«ãªã£ãŸ æ‰‹ã‚’æŒ™ã’ã‚‹ã®ãŒåˆå›³ã«ãªã‚‹ å½¢å®¹è©ãƒ»ä¿®é£¾èª + åˆå›³ å°ã•ãªåˆå›³ å¾®å¦™ãªåˆå›³ æ˜ç¢ºãªåˆå›³ ç„¡è¨€ã®åˆå›³ è¦‹ãˆã«ãã„åˆå›³",
    # "å‘¨å›²": "åè© + ã® + åè© å‘¨å›²ã®äººã€… å‘¨å›²ã®ç’°å¢ƒ å‘¨å›²ã®çŠ¶æ³ å‘¨å›²ã®ç›® å‘¨å›²ã®äººãŸã¡ ï½ã® + å‘¨å›² + ã«ï¼ã§ï¼ã¸ï¼ã‹ã‚‰ å‘¨å›²ã«æ³¨æ„ã™ã‚‹ å‘¨å›²ã§è¦³å¯Ÿã™ã‚‹ å‘¨å›²ã‹ã‚‰è¦‹ãˆã‚‹ å‘¨å›²ã¸å£°ã‚’ã‹ã‘ã‚‹ å‘¨å›²ã«æº¶ã‘è¾¼ã‚€ å‘¨å›² + ã‚’ + å‹•è© å‘¨å›²ã‚’è¦‹æ¸¡ã™ å‘¨å›²ã‚’æ°—ã«ã™ã‚‹ å‘¨å›²ã‚’è¦³å¯Ÿã™ã‚‹ å‘¨å›²ã‚’é¨’ãŒã›ã‚‹ å‘¨å›²ã‚’å–ã‚Šå·»ã å½¢å®¹è© + å‘¨å›²ï¼å‘¨å›²ã® åºƒã„å‘¨å›² ç‹­ã„å‘¨å›² è¿‘ã„å‘¨å›² å¿«é©ãªå‘¨å›² é™ã‹ãªå‘¨å›²",
    # "ä¸‹ã‚ã™": "åè© + ä¸‹ã‚ã™ï¼é™ã‚ã™ï¼å¸ã™ *ï¼ˆä½•ã‚’ã€Œä¸‹ã‚ã™ï¼é™ã‚ã™ï¼å¸ã™ã€ã‹ï¼‰ *è·ç‰©ã‚’ä¸‹ã‚ã™ *ãŠé‡‘ã‚’ä¸‹ã‚ã™ *ãƒ–ãƒ©ã‚¤ãƒ³ãƒ‰ã‚’ä¸‹ã‚ã™ *çœ‹æ¿ã‚’ä¸‹ã‚ã™ï¼é™ã‚ã™ *å•†å“ã‚’å¸ã™ *ä¸‹ã‚ã™ï¼é™ã‚ã™ + åè© *ï¼ˆã€Œä¸‹ã‚ã™ã€ã—ãŸã‚‚ã®ã‚’æŒ‡ã™ï¼‰ *ä¸‹ã‚ã—ãŸè·ç‰© *ä¸‹ã‚ã—ãŸãƒ–ãƒ©ã‚¤ãƒ³ãƒ‰ *ä¸‹ã‚ã—ãŸçœ‹æ¿ *å¸ã—ãŸå•†å“ *ä¸‹ã‚ã—ãŸãŠé‡‘ *å‰¯è© + ä¸‹ã‚ã™ï¼é™ã‚ã™ï¼å¸ã™ *ï¼ˆå‹•ä½œã®æ§˜å­ãƒ»ç¨‹åº¦ã‚’è¡¨ã™ï¼‰ *ã‚†ã£ãã‚Šä¸‹ã‚ã™ *çœŸã£ã™ãä¸‹ã‚ã™ *ä¸€æ°—ã«ä¸‹ã‚ã™ *ã“ã£ãã‚Šä¸‹ã‚ã™ *å¤§é‡ã«å¸ã™ *ä¸‹ã‚ã™ï¼é™ã‚ã™ + åŠ©è©å¥ *ï¼ˆç›®çš„ãƒ»å ´æ‰€ãªã©ã‚’è¡¨ã™èªã¨ã®çµ„åˆã›ï¼‰ *ï½ã‹ã‚‰ä¸‹ã‚ã™ï¼ˆä¾‹ï¼šæ£šã‹ã‚‰ä¸‹ã‚ã™ï¼‰ *ï½ã«ä¸‹ã‚ã™ï¼ˆä¾‹ï¼šåºŠã«ä¸‹ã‚ã™ï¼‰ *ï½ã‹ã‚‰é™ã‚ã™ï¼ˆä¾‹ï¼šãƒã‚¹ã‹ã‚‰é™ã‚ã™ï¼‰ *ï½ã¸å¸ã™ï¼ˆä¾‹ï¼šå¸å…ˆã«å¸ã™ï¼‰ *ï½ã¾ã§ä¸‹ã‚ã™ *æ…£ç”¨è¡¨ç¾ãƒ»æ…£ç”¨å¥çš„ç”¨æ³• *ï¼ˆæ¯”å–©çš„ãƒ»æ…£ç”¨çš„ãªä½¿ã„æ–¹ï¼‰ *è‚©ã®è·ã‚’ä¸‹ã‚ã™ *èƒ¸ã‚’ãªã§ä¸‹ã‚ã™ *äººã‚’ã“ãä¸‹ã‚ã™ *æ ¹ã‚’ä¸‹ã‚ã™ *é«ªã‚’ä¸‹ã‚ã™",
    # "ä»˜è¿‘": "åè© + ã® + åè©ï¼ˆï½ã®ä»˜è¿‘ï¼‰ é§…ã®ä»˜è¿‘ å…¬åœ’ã®ä»˜è¿‘ å­¦æ ¡ã®ä»˜è¿‘ å•†åº—è¡—ã®ä»˜è¿‘ é§è»Šå ´ã®ä»˜è¿‘ åè© + ã« + å‹•è©ï¼ˆä»˜è¿‘ã«ï½ãŒã‚ã‚‹ï¼ä»˜è¿‘ã«ï½ãŒè¦‹ãˆã‚‹ ç­‰ï¼‰ ä»˜è¿‘ã«ã‚ã‚‹ ä»˜è¿‘ã«è¦‹ãˆã‚‹ ä»˜è¿‘ã«ä½ç½®ã™ã‚‹ ä»˜è¿‘ã«ç‚¹åœ¨ã™ã‚‹ ä»˜è¿‘ã«å­˜åœ¨ã™ã‚‹ åè© + ã§ + å‹•è©ï¼ˆï½ä»˜è¿‘ã§ï½ã™ã‚‹ï¼‰ é§…ä»˜è¿‘ã§å¾…ã¡åˆã‚ã›ã™ã‚‹ å…¬åœ’ä»˜è¿‘ã§æ•£æ­©ã™ã‚‹ å•†åº—è¡—ä»˜è¿‘ã§è²·ã„ç‰©ã™ã‚‹ äº¤å·®ç‚¹ä»˜è¿‘ã§æ¸‹æ»ã™ã‚‹ é“è·¯ä»˜è¿‘ã§ç«‹ã¡æ­¢ã¾ã‚‹ å‰¯è© + åè©ï¼å‹•è©ï¼ˆå¤šå°‘ä¿®é£¾èªã¨ã¨ã‚‚ã«ä½¿ã‚ã‚Œã‚‹ï¼‰ ã™ãä»˜è¿‘ è¿‘ã„ä»˜è¿‘ ãã®ä»˜è¿‘ å‘¨è¾ºä»˜è¿‘ ä»˜è¿‘ä¸€å¸¯ å½¢å®¹è© + åè©ï¼ˆä¿®é£¾èªï¼‹ä»˜è¿‘ï¼‰ å‘¨è¾ºä»˜è¿‘ è¿‘è¾ºä»˜è¿‘ ç›´è¿‘ä»˜è¿‘ æœ€å¯„ã‚Šä»˜è¿‘ ä½å®…ä»˜è¿‘",
    # "å‘½ä»¤": "åè© + å‹•è© ã€Œå‘½ä»¤ + å‹•è©ã€ã®å½¢ã§ã€ã€Œå‘½ä»¤ã‚’ï¼ˆï½ã™ã‚‹ï¼å‡ºã™ï¼å—ã‘ã‚‹ï¼ä¸‹ã™â€¦ï¼‰ã€ãªã© å‘½ä»¤ã‚’å‡ºã™ å‘½ä»¤ã‚’å—ã‘ã‚‹ å‘½ä»¤ã‚’ä¸‹ã™ å‘½ä»¤ã‚’ä¼ãˆã‚‹ å‘½ä»¤ã‚’å®ˆã‚‹ å‹•è© + å‘½ä»¤ å‹•è©ã¨ã€Œå‘½ä»¤ã€ãŒçµã³ã¤ãå½¢ã€‚ã€Œï½å‘½ä»¤ã€ã€Œå‘½ä»¤ï½ã€ã€Œï½ã‚’å‘½ä»¤ã™ã‚‹ã€ãªã© å¼·ã„å‘½ä»¤ã‚’ï¼ˆå‡ºã™ï¼ä¸‹ã™ï¼‰ ä¸Šã‹ã‚‰ã®å‘½ä»¤ ä¸å½“ãªå‘½ä»¤ é•æ³•ãªå‘½ä»¤ å‘½ä»¤æ¨©ã‚’æŒã¤ åè© + ã® + å‘½ä»¤ æ‰€æœ‰ãƒ»ä¿®é£¾ã®å½¢ã§ã€Œï½ã®å‘½ä»¤ã€ã€Œå‘½ä»¤ã®ï½ã€ ä¸Šå¸ã®å‘½ä»¤ è¦ªã®å‘½ä»¤ è¦å‰‡ã®å‘½ä»¤ è»ã®å‘½ä»¤ çµ„ç¹”ã®å‘½ä»¤ å‘½ä»¤ + ã«ã‚ˆã‚‹ï¼ã«ã‚ˆã£ã¦ ã€Œæ‰‹æ®µãƒ»åŸå› ãƒ»æ ¹æ‹ ã€ã‚’ç¤ºã™å½¢ã§ä½¿ã† å‘½ä»¤ã«ã‚ˆã‚‹å‡¦åˆ† å‘½ä»¤ã«ã‚ˆã£ã¦å®Ÿè¡Œã•ã‚Œã‚‹ å‘½ä»¤ã«ã‚ˆã‚‹æ··ä¹± å‘½ä»¤ã«ã‚ˆã‚‹å¤‰æ›´ å‘½ä»¤ã«ã‚ˆã‚‹å½±éŸ¿ å‘½ä»¤ + ã‚’ + åè© å‘½ä»¤ãŒå¯¾è±¡ãƒ»å†…å®¹ã‚’ç¤ºã™å½¢ å‘½ä»¤ã‚’éµå®ˆï¼ˆã˜ã‚…ã‚“ã—ã‚…ï¼‰ã™ã‚‹ å‘½ä»¤ã‚’æ”¹è¨‚ã™ã‚‹ å‘½ä»¤ã‚’å–æ¶ˆã™ å‘½ä»¤ã‚’å¼·åŒ–ã™ã‚‹ å‘½ä»¤ã‚’é‚è¡Œã™ã‚‹",
    # "å›å¾©": "åè© + åŠ©è©ï¼ˆå›å¾© + ã‚’ï¼ãŒï¼ã«ï¼‰ï¼‹å‹•è© ï¼ˆå¯¾è±¡ï¼‹å›å¾©ã™ã‚‹ï¼å›å¾©ã‚’ï½ã™ã‚‹ï¼å›å¾©ãŒï½ã™ã‚‹ ãªã©ï¼‰ å¥åº·ãŒå›å¾©ã™ã‚‹ ä½“åŠ›ãŒå›å¾©ã™ã‚‹ å‚™è“„ã‚’å›å¾©ã™ã‚‹ ä¿¡ç”¨ã‚’å›å¾©ã™ã‚‹ åèª‰ã‚’å›å¾©ã™ã‚‹ ï½ãŒ + å›å¾© + ã™ã‚‹ ï¼ˆä¸»èªãŒ â€œã€œãŒå›å¾©ã™ã‚‹â€ ã®å½¢ï¼‰ å¤©å€™ãŒå›å¾©ã™ã‚‹ æ™¯æ°—ãŒå›å¾©ã™ã‚‹ ãƒ€ã‚¤ãƒ¤ï¼ˆé‹è¡Œï¼‰ãŒå›å¾©ã™ã‚‹ æ°—åŠ›ãŒå›å¾©ã™ã‚‹ æ„è­˜ãŒå›å¾©ã™ã‚‹ å›å¾© + ã® + åè© ï¼ˆå›å¾©ã‚’ä¿®é£¾ã™ã‚‹èªå¥ï¼‰ å›å¾©åŠ› å›å¾©æœŸ å›å¾©éç¨‹ å›å¾©å‚¾å‘ å›å¾©é€Ÿåº¦ å›å¾© + ã‚’ + å‹•è© ï¼ˆå›å¾©ã‚’è¡Œç‚ºã¨ã—ã¦è¡¨ã™èªå¥ï¼‰ å›å¾©ã‚’å¾…ã¤ å›å¾©ã‚’æ—©ã‚ã‚‹ å›å¾©ã‚’å›³ã‚‹ å›å¾©ã‚’ä¿ƒã™ å›å¾©ã‚’ä¿è¨¼ã™ã‚‹ ä»–å‹•è© + å›å¾© + ã™ã‚‹ ï¼ˆï½ãŒå›å¾©ã‚’ â€œã™ã‚‹â€ / â€œã•ã›ã‚‹â€ ã®æ§‹æ–‡ï¼‰ å›å¾©ã•ã›ã‚‹ å›å¾©ã™ã‚‹ å›å¾©ã•ã›ã¦ãã‚Œã‚‹ å›å¾©ã§ãã‚‹ å›å¾©ã•ã›ã‚ˆã†ã¨ã™ã‚‹",
    # "å›°é›£": "åè© + åŠ©è©ï¼å‹•è©ï¼ˆå›°é›£ + ã‚’ï¼ãŒï¼ã« + å‹•è©ï¼‰ å›°é›£ã‚’ä¹—ã‚Šè¶Šãˆã‚‹ å›°é›£ã‚’å…‹æœã™ã‚‹ å›°é›£ã‚’æŠ±ãˆã‚‹ å›°é›£ã«ç›´é¢ã™ã‚‹ å›°é›£ãŒç”Ÿã˜ã‚‹ å½¢å®¹è©ï¼é€£ä½“ä¿®é£¾ + å›°é›£ éå¸¸ã«å›°é›£ æ¥µã‚ã¦å›°é›£ æ¥µåº¦ã®å›°é›£ å¤§ããªå›°é›£ æ·±åˆ»ãªå›°é›£ å›°é›£ + ã® + åè© å›°é›£ã• å›°é›£åº¦ å›°é›£æ€§ å›°é›£ãªçŠ¶æ³ å›°é›£ãªèª²é¡Œ å›°é›£ + ã‚’ + å‹•è© å›°é›£ã‚’æ„Ÿã˜ã‚‹ å›°é›£ã‚’äºˆæƒ³ã™ã‚‹ å›°é›£ã‚’å…‹æœã™ã‚‹ å›°é›£ã‚’å¼•ãå—ã‘ã‚‹ å›°é›£ã‚’ç†è§£ã™ã‚‹",
    # "åŸºæœ¬": "å½¢å®¹è© + åŸºæœ¬ ï¼ˆã€ŒåŸºæœ¬ã€ã‚’ä¿®é£¾ã™ã‚‹èªï¼‰ æœ€ã‚‚åŸºæœ¬ ã”ãåŸºæœ¬ åŸºæœ¬çš„ å®Ÿè³ªçš„ãªåŸºæœ¬ æ ¹æœ¬çš„ãªåŸºæœ¬ åŸºæœ¬ + ã® + åè© ï¼ˆåŸºæœ¬ã‚’ä¿®é£¾èªã¨ã™ã‚‹å½¢ï¼‰ åŸºæœ¬æ¦‚å¿µ åŸºæœ¬åŸå‰‡ åŸºæœ¬æ§‹é€  åŸºæœ¬æ–¹é‡ åŸºæœ¬å§¿å‹¢ åè© + ã‚’ + åŸºæœ¬ + ã¨ã™ã‚‹ ï¼ˆï½ã‚’åŸºæœ¬ã¨ã™ã‚‹ï¼‰ å®‰å…¨ã‚’åŸºæœ¬ã¨ã™ã‚‹ ä¿¡é ¼ã‚’åŸºæœ¬ã¨ã™ã‚‹ åˆ©ç›Šã‚’åŸºæœ¬ã¨ã™ã‚‹ ç¶™ç¶šã‚’åŸºæœ¬ã¨ã™ã‚‹ åŠ¹ç‡ã‚’åŸºæœ¬ã¨ã™ã‚‹ åŸºæœ¬ + ã« + å‹•è©ï¼å½¢å®¹è© ï¼ˆåŸºæœ¬ã«â—‹â—‹ã™ã‚‹ï¼åŸºæœ¬ã«â—‹â—‹ã§ã‚ã‚‹ï¼‰ åŸºæœ¬ã«å¿ å®Ÿã§ã‚ã‚‹ åŸºæœ¬ã«å¾“ã† åŸºæœ¬ã«ç«‹ã¡æˆ»ã‚‹ åŸºæœ¬ã«ç…§ã‚‰ã™ åŸºæœ¬ã«åŸºã¥ã å‹•è© + åŸºæœ¬ + ã‚’ + åè©ï¼å‹•è© ï¼ˆï½ã‚’åŸºæœ¬ã¨ã™ã‚‹ã€ï½ã‚’åŸºæœ¬ã«ï½ã™ã‚‹ï¼‰ ï½ã‚’åŸºæœ¬ã«æ®ãˆã‚‹ ï½ã‚’åŸºæœ¬ã«è€ƒãˆã‚‹ ï½ã‚’åŸºæœ¬ã¨ã—ã¦æ¡ç”¨ã™ã‚‹ ï½ã‚’åŸºæœ¬ã«è¨­è¨ˆã™ã‚‹ ï½ã‚’åŸºæœ¬ã«å®Ÿè¡Œã™ã‚‹",
    # "å¥³å„ª": "å½¢å®¹è©ï¼é€£ä½“ä¿®é£¾ + å¥³å„ª å¥³å„ªã‚’ä¿®é£¾ã™ã‚‹è¨€è‘‰ è‹¥æ‰‹å¥³å„ª ãƒ™ãƒ†ãƒ©ãƒ³å¥³å„ª åå¥³å„ª äººæ°—å¥³å„ª å®ŸåŠ›æ´¾å¥³å„ª å¥³å„ª + ã® + åè© å¥³å„ªãŒæ‰€æœ‰ãƒ»æ‰€å±ãƒ»æ€§è³ªã‚’ç¤ºã™èª å¥³å„ªã®åµ å¥³å„ªã®åµï¼ˆæ„å‘³ï¼šã“ã‚Œã‹ã‚‰å¥³å„ªã«ãªã‚‹äººï¼‰ å¥³å„ªã®ç´ è³ª å¥³å„ªã®é“ å¥³å„ªã®çµŒæ­´ åè© + ã‚’ + å¥³å„ª å¯¾è±¡ã‚’å¥³å„ªã«ã™ã‚‹ãƒ»æ‰±ã†èª ã€œã‚’å¥³å„ªã«æŠœæ“¢ã™ã‚‹ ã€œã‚’å¥³å„ªã¨ã—ã¦èµ·ç”¨ã™ã‚‹ ã€œã‚’å¥³å„ªã¨ã—ã¦å£²ã‚Šå‡ºã™ ã€œã‚’å¥³å„ªã¨ã—ã¦èªã‚ã‚‰ã‚Œã‚‹ å¥³å„ª + ã‚’ + å‹•è© å¥³å„ªã«é–¢ã™ã‚‹å‹•ä½œèª å¥³å„ªã‚’ç›®æŒ‡ã™ å¥³å„ªã‚’æ¼”ã˜ã‚‹ï¼ˆæ¯”å–©çš„ã«ï¼‰ å¥³å„ªã‚’è‚²ã¦ã‚‹ å¥³å„ªã‚’æ”¯ãˆã‚‹ ï½ã¨ / ï½ã¨ã—ã¦ + å¥³å„ª ç«‹å ´ãƒ»æ¯”è¼ƒãƒ»å½¹å‰²ã‚’ç¤ºã™èª æ˜ ç”»å¥³å„ª èˆå°å¥³å„ª ãƒˆãƒƒãƒ—å¥³å„ª å›½éš›çš„å¥³å„ª å¥³å„ªã¨ã—ã¦æˆé•·ã™ã‚‹",
    # "å®Œå…¨": "å½¢å®¹è©ï¼é€£ä½“ä¿®é£¾ + å®Œå…¨ å®Œå…¨ã‚’ä¿®é£¾ã™ã‚‹è¨€è‘‰ ã¾ã£ãŸãå®Œå…¨ çœŸã®å®Œå…¨ ã»ã¼å®Œå…¨ å®Œå…¨ç„¡æ¬ ãª å®Œå…¨ãª å®Œå…¨ + ã® + åè© å®Œå…¨ã‚’ä¿®é£¾èªã¨ã™ã‚‹åè© å®Œå…¨æ€§ å®Œå…¨ç‰ˆ å®Œå…¨çŠ¯ç½ª å®Œå…¨ä¸€è‡´ å®Œå…¨ç‡ƒç„¼ åè© + ã‚’ + å®Œå…¨ + ã« + å‹•è© ï¼ˆï½ã‚’å®Œå…¨ã« â—‹â—‹ã™ã‚‹ï¼‰ éšœå®³ã‚’å®Œå…¨ã«å…‹æœã™ã‚‹ ä»»å‹™ã‚’å®Œå…¨ã«é‚è¡Œã™ã‚‹ åˆæ„ã‚’å®Œå…¨ã«å½¢æˆã™ã‚‹ æ•´å‚™ã‚’å®Œå…¨ã«è¡Œã† åˆ¶åº¦ã‚’å®Œå…¨ã«æ•´ãˆã‚‹ å®Œå…¨ + ã« + å‹•è©ï¼å½¢å®¹è© ï¼ˆâ€œå®Œå…¨ã«ï½ã™ã‚‹ï¼ï½ã§ã‚ã‚‹â€ã®æ§‹æ–‡ï¼‰ å®Œå…¨ã«æº€è¶³ã™ã‚‹ å®Œå…¨ã«ç†è§£ã™ã‚‹ å®Œå…¨ã«æ¶ˆå¤±ã™ã‚‹ å®Œå…¨ã«è§£æ±ºã™ã‚‹ å®Œå…¨ã«æ”¯é…ã™ã‚‹ ï½ãŒ + å®Œå…¨ + ã« + ï½ã™ã‚‹ ï¼ˆä¸»èª + å®Œå…¨ã« + å‹•è©ï¼‰ ã‚·ã‚¹ãƒ†ãƒ ãŒå®Œå…¨ã«ç¨¼åƒã™ã‚‹ æ¨©åŠ›ãŒå®Œå…¨ã«æŒæ¡ã•ã‚Œã‚‹ æ”¯é…ãŒå®Œå…¨ã«è¡Œãå±Šã ä¾å­˜ãŒå®Œå…¨ã«è§£æ¶ˆã™ã‚‹ æ¤œæŸ»ãŒå®Œå…¨ã«çµ‚äº†ã™ã‚‹",
    # "å®—æ•™": "å½¢å®¹è©ï¼é€£ä½“ä¿®é£¾ + å®—æ•™ å®—æ•™ã‚’ä¿®é£¾ã™ã‚‹èª ä¸–ç•Œå®—æ•™ å®—æ•™çš„ ä¼çµ±å®—æ•™ å®—æ•™é–“ å›½å®¶å®—æ•™ å®—æ•™ + ã® + åè© å®—æ•™ã‚’ä¿®é£¾èªã¨ã™ã‚‹èªå¥ å®—æ•™æ³•äºº å®—æ•™è¦³ å®—æ•™ä¿¡ä»° å®—æ•™å„€å¼ å®—æ•™å›£ä½“ åè© + ã‚’ + å®—æ•™ ï¼ˆï½ã‚’å®—æ•™ã«ã™ã‚‹ï¼æ‰±ã†æ§‹æ–‡ãªã©ï¼‰ ã‚­ãƒªã‚¹ãƒˆæ•™ã‚’å®—æ•™ã¨ã™ã‚‹ å®—æ•™ã‚’å•é¡Œè¦–ã™ã‚‹ å®—æ•™ã‚’ä¿¡ã˜ã‚‹ å®—æ•™ + ã‚’ + å‹•è© å®—æ•™ã«é–¢ã™ã‚‹å‹•ä½œèª å®—æ•™ã‚’ä¿¡ä»°ã™ã‚‹ å®—æ•™ã‚’åºƒã‚ã‚‹ å®—æ•™ã‚’å¦å®šã™ã‚‹ å®—æ•™ã‚’å°Šé‡ã™ã‚‹ å®—æ•™ã‚’æ”¹é©ã™ã‚‹ ï½ã¨ / ï½ã¨ã—ã¦ + å®—æ•™ ç«‹å ´ãƒ»æ¯”è¼ƒãƒ»å½¹å‰²ã‚’ç¤ºã™èª å®—æ•™ã¨ã—ã¦ å®—æ•™æ³•äººã¨ã—ã¦ å®—æ•™è‰² å®—æ•™æ–½è¨­ å®—æ•™æ´»å‹•",
    # "å¹…": "åè© + å‹•è© å¹…ãŒåºƒãŒã‚‹ å¹…ãŒç‹­ã¾ã‚‹ å¹…ã‚’æŒã¤ å¹…ã‚’åºƒã’ã‚‹ å¹…ã‚’åˆ¶é™ã™ã‚‹ å½¢å®¹è© + åè© åºƒã„å¹… ç‹­ã„å¹… ååˆ†ãªå¹… æœ€å¤§ã®å¹… é©åˆ‡ãªå¹… åè© + åè© å¹…æ–¹å‘ å¹…å“¡ï¼ˆã¯ã°ã„ã‚“ï¼‰ å¹…åºƒã• å¹…åˆ¶é™ å¹…å¯„ã› å‹•è© + åè© å¢—ã™å¹… æ¸›ã‚‹å¹… å¤‰åŒ–ã®å¹… è¨±å®¹ã®å¹… å½±éŸ¿ã®å¹…",
    # "å»¶æœŸ": "åè© + å‹•è© å»¶æœŸã™ã‚‹ å»¶æœŸã«ãªã‚‹ å»¶æœŸã‚’æ±ºå®šã™ã‚‹ å»¶æœŸã‚’ç™ºè¡¨ã™ã‚‹ å»¶æœŸã‚’è¦è«‹ã™ã‚‹ å½¢å®¹è© + åè© æ€¥ãªå»¶æœŸ ä¸å¯é¿ãªå»¶æœŸ å¤§å¹…ãªå»¶æœŸ è‡¨æ™‚ã®å»¶æœŸ é•·æœŸã®å»¶æœŸ åè© + åè© å»¶æœŸæ±ºå®š å»¶æœŸé€šçŸ¥ å»¶æœŸæªç½® å»¶æœŸæœŸé–“ å»¶æœŸç†ç”± å‹•è© + åè© è¡Œäº‹ã®å»¶æœŸ ä¼šè­°ã®å»¶æœŸ å‡ºç™ºã®å»¶æœŸ ç™ºè¡¨ã®å»¶æœŸ ç· åˆ‡ã®å»¶æœŸ",
    # "å»ºè¨­": "åè© ï¼‹ å‹•è©ï¼ˆï¼ã€Œï½ã‚’å»ºè¨­ã™ã‚‹ã€ãªã©ï¼‰ å»ºè¨­ã™ã‚‹ å»ºè¨­ã‚’è¡Œã† å»ºè¨­ã‚’é€²ã‚ã‚‹ å»ºè¨­ã«ç€æ‰‹ã™ã‚‹ å»ºè¨­ãŒé€²ã‚€ å½¢å®¹è©ï¼ä¿®é£¾èª ï¼‹ åè©ï¼ˆã€Œï½ãªå»ºè¨­ã€ï¼‰ å¤§è¦æ¨¡å»ºè¨­ æ€¥é€Ÿãªå»ºè¨­ å®‰å…¨ãªå»ºè¨­ å…±åŒå»ºè¨­ æŒç¶šå¯èƒ½ãªå»ºè¨­ åè© + åè©ï¼ˆã€Œï½ã®å»ºè¨­ã€ã€Œå»ºè¨­ï½ã€ãªã©ï¼‰ å»ºè¨­å·¥äº‹ å»ºè¨­è²»ç”¨ å»ºè¨­è¨ˆç”» å»ºè¨­æ¥­è€… å»ºè¨­ç¾å ´ å‹•è©å¥ï¼å¥å‹•è© ï¼‹ åè©ï¼ˆå‹•è©ï¼‹ç›®çš„èªçš„ãªçµ„ã¿åˆã‚ã›ï¼‰ åœŸåœ°ã‚’å»ºè¨­ã™ã‚‹ï¼ˆâ€»ãŸã ã—ã€ŒåœŸåœ°ã‚’å»ºè¨­ã™ã‚‹ã€ã¯ã‚„ã‚„ä¸è‡ªç„¶ã€‚ã€Œå®…åœ°ã‚’é€ æˆã—ã¦å»ºè¨­ã™ã‚‹ã€ãªã©æ–‡è„ˆæ¬¡ç¬¬ï¼‰ æ©‹ã‚’å»ºè¨­ã™ã‚‹ é“è·¯ã‚’å»ºè¨­ã™ã‚‹ ãƒ€ãƒ ã‚’å»ºè¨­ã™ã‚‹ æ–½è¨­ã‚’å»ºè¨­ã™ã‚‹",
    # "å¼•é€€": "åè© + å‹•è©ï¼ˆï¼ã€Œã€œã‚’å¼•é€€ã™ã‚‹ï¼ã€œãŒå¼•é€€ã™ã‚‹ã€ãªã©ï¼‰ å¼•é€€ã™ã‚‹ å¼•é€€ã‚’è¡¨æ˜ã™ã‚‹ å¼•é€€ã‚’æ±ºæ–­ã™ã‚‹ å¼•é€€ãŒå›ã‹ã‚Œã‚‹ å¼•é€€ã‚’è¿ãˆã‚‹ å½¢å®¹è©ï¼ä¿®é£¾èª + åè©ï¼ˆï¼ã€Œã€œãªå¼•é€€ã€ï¼‰ çªç„¶ã®å¼•é€€ ä¸æœ¬æ„ãªå¼•é€€ å¾…æœ›ã®å¼•é€€ å††æº€ãªå¼•é€€ æ—©æœŸã®å¼•é€€ åè© + åè©ï¼ˆï¼ã€Œã€œã®å¼•é€€ã€ã€Œå¼•é€€ã€œã€ãªã©ï¼‰ å¼•é€€è¡¨æ˜ å¼•é€€è©¦åˆ å¼•é€€ä¼šè¦‹ å¼•é€€ã‚»ãƒ¬ãƒ¢ãƒ‹ãƒ¼ å¼•é€€å¾Œ å‹•è©å¥ï¼å¥å‹•è© + åè©ï¼ˆå‹•è©ï¼‹ç›®çš„èªãƒ»è£œèªã¨ã—ã¦ã®çµ„ã¿åˆã‚ã›ï¼‰ ç¾å½¹ã‚’å¼•é€€ã™ã‚‹ ã‚¹ãƒãƒ¼ãƒ„ç•Œã‚’å¼•é€€ã™ã‚‹ å½¹è·ã‚’å¼•é€€ã™ã‚‹ å¼•é€€å¾Œã®ç”Ÿæ´» å¼•é€€ã‚’æ’¤å›ã™ã‚‹",
    # "æ‚©ã‚€": "åè© + å‹•è©ï¼ˆï¼ã€Œï½ã«æ‚©ã‚€ï¼ï½ã§æ‚©ã‚€ã€ãªã©ï¼‰ å°†æ¥ã«æ‚©ã‚€ å•é¡Œã«æ‚©ã‚€ ç—…æ°—ã§æ‚©ã‚€ äººé–“é–¢ä¿‚ã§æ‚©ã‚€ è‡ªåˆ†ã®æ€§æ ¼ã«æ‚©ã‚€ å½¢å®¹è©ï¼ä¿®é£¾èª + åè©ï¼ˆï¼ã€Œï½ãªæ‚©ã¿ã€ï¼‰ æ·±åˆ»ãªæ‚©ã¿ æ—¥å¸¸çš„ãªæ‚©ã¿ å°ã•ãªæ‚©ã¿ é•·å¹´ã®æ‚©ã¿ å¿ƒã®æ‚©ã¿ åè© + åè©ï¼ˆï¼ã€Œï½ã®æ‚©ã¿ï¼æ‚©ã¿ï½ã€ãªã©ï¼‰ æ‚©ã¿ã®ç¨® æ‚©ã¿ç›¸è«‡ æ‚©ã¿è§£æ¶ˆ æ‚©ã¿äº‹ æ‚©ã¿å¤šã å‹•è©å¥ï¼å¥å‹•è© + åè©ï¼ˆå‹•è©ï¼‹ç›®çš„èªçš„ãªçµ„ã¿åˆã‚ã›ï¼‰ æ‚©ã¿ã‚’æŠ±ãˆã‚‹ æ‚©ã¿ã‚’èã æ‚©ã¿ã‚’è©±ã™ æ‚©ã¿ã‚’å…±æœ‰ã™ã‚‹ æ‚©ã¿ã‚’æ‰“ã¡æ˜ã‘ã‚‹",
    # "æ‚ªé­”": "åè© + å‹•è©ï¼å‹•è©å¥ï¼ˆï½ãŒï¼ï½ã‚’ï½ã™ã‚‹ï¼‰ æ‚ªé­”ãŒã•ã•ã‚„ã æ‚ªé­”ã‚’å¬å–šã™ã‚‹ æ‚ªé­”ã«å–ã‚Šæ†‘ã‹ã‚Œã‚‹ æ‚ªé­”ã‚’è¿½ã„æ‰•ã† æ‚ªé­”ã¨å¥‘ç´„ã™ã‚‹ ä¿®é£¾èªï¼å½¢å®¹è© + åè©ï¼ˆï½ãªæ‚ªé­”ï¼‰ æš—é»’ã®æ‚ªé­” å¼·å¤§ãªæ‚ªé­” é‚ªæ‚ªãªæ‚ªé­” å¤ã®æ‚ªé­” é†œã„æ‚ªé­” åè© + åè©ï¼ˆï½ã®æ‚ªé­”ï¼æ‚ªé­”ï½ãªã©ï¼‰ æ‚ªé­”ã®ä½¿ã„ æ‚ªé­”ã®å¥‘ç´„ æ‚ªé­”ã®ã•ã•ã‚„ã æ‚ªé­”ã®åŠ› æ‚ªé­”ã®å½± å‹•è©å¥ + åè©ï¼ˆå‹•è©ï¼‹ç›®çš„èªã¨ã—ã¦ä½¿ã†ãƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ æ‚ªé­”ã‚’å‘¼ã¶ æ‚ªé­”ã‚’å°ã˜ã‚‹ æ‚ªé­”ã‚’å€’ã™ æ‚ªé­”ã‚’å°å°ã™ã‚‹ æ‚ªé­”ã‚’é€€ã‘ã‚‹",
    # "æ„æ€": "æ„æ€ã‚’ä¼ãˆã‚‹ æ„æ€ã‚’è¡¨ç¤ºã™ã‚‹ æ„æ€ã‚’ç¤ºã™ æ„æ€ã‚’ç¢ºèªã™ã‚‹ æ„æ€ã‚’å°Šé‡ã™ã‚‹ å½¢å®¹è©ï¼ä¿®é£¾èª + åè©ï¼ˆï¼ã€Œï½ãªæ„æ€ã€ï¼‰ æ˜ç¢ºãªæ„æ€ å¼·ã„æ„æ€ å›ºã„æ„æ€ è‡ªåˆ†ã®æ„æ€ è‡ªè¦šçš„ãªæ„æ€ åè© + åè©ï¼ˆï¼ã€Œï½ã®æ„æ€ï¼æ„æ€ï½ã€ãªã©ï¼‰ æ„æ€è¡¨ç¤º æ„æ€æ±ºå®š æ„æ€ç–é€š æ„æ€è¡¨ç¤ºæ›¸ï¼ˆã¾ãŸã¯ æ„æ€è¡¨æ˜æ›¸ï¼‰ æ„æ€ç¢ºèª å‹•è©å¥ï¼å¥å‹•è© + åè©ï¼ˆå‹•è©ï¼‹ç›®çš„èªçš„ãªçµ„ã¿åˆã‚ã›ï¼‰ æ„æ€ã‚’å›ºã‚ã‚‹ æ„æ€ã‚’æ›²ã’ã‚‹ æ„æ€ã‚’æŒã¤ æ„æ€è¡¨ç¤ºã‚’è¡Œã† æ„æ€æ±ºå®šã‚’ä¸‹ã™",
    # "æ†§ã‚Œã‚‹": "å‹•è© + æ†§ã‚Œã‚‹ ï¼ˆã‚ã‚‹å¯¾è±¡ã«æ†§ã‚Œã‚’æŠ±ããƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ â—‹â—‹ã«æ†§ã‚Œã‚‹ è‹¥ã„æ™‚ä»£ã«æ†§ã‚Œã‚‹ é ã„å›½ã«æ†§ã‚Œã‚‹ éƒ½ä¼šç”Ÿæ´»ã«æ†§ã‚Œã‚‹ ã‚ã®äººã«æ†§ã‚Œã‚‹ åè© + ã«æ†§ã‚Œã‚‹ ï¼ˆåè©ã‚’å¯¾è±¡ã«æ†§ã‚Œã‚‹è¡¨ç¾ï¼‰ èŠ¸èƒ½äººã«æ†§ã‚Œã‚‹ æœ‰åäººã«æ†§ã‚Œã‚‹ ã‚ã®äººã«æ†§ã‚Œã‚‹ ãƒ¢ãƒ‡ãƒ«ã«æ†§ã‚Œã‚‹ å…ˆè¼©ã«æ†§ã‚Œã‚‹ æ†§ã‚Œã‚‹ + å¯¾è±¡ã‚’è¡¨ã™èªå¥ ï¼ˆæ†§ã‚Œã‚‹å¯¾è±¡ã‚’è£œè¶³ã™ã‚‹èªå¥ã¨å…±èµ·ï¼‰ æ†§ã‚Œã‚‹å­˜åœ¨ æ†§ã‚Œã‚‹ç”Ÿæ´» æ†§ã‚Œã‚‹å¤¢ æ†§ã‚Œã‚‹ãƒ¢ãƒ‡ãƒ« æ†§ã‚Œã‚‹äººç‰© ã€œã«æ†§ã‚Œã‚’æŠ±ãï¼æŒã¤ ï¼ˆæ†§ã‚Œã‚’åè©å½¢ã§è¿°ã¹ã‚‹è¨€ã„æ–¹ï¼‰ æ†§ã‚Œã‚’æŠ±ã æ†§ã‚Œã‚’æŒã¤ æ†§ã‚Œã®å¿µ æ†§ã‚Œã®çš„ æ†§ã‚Œã®äºº æ†§ã‚Œã‚‹ + çŠ¶æ…‹ï¼å‹•è©è£œåŠ©è¡¨ç¾ ï¼ˆæ†§ã‚Œã®æ°—æŒã¡ã®çŠ¶æ…‹ã‚’è¡¨ã™èªã¨å…±èµ·ï¼‰ æ†§ã‚Œã¦ã‚„ã¾ãªã„ æ†§ã‚Œã¦ã„ã‚‹ æ†§ã‚Œã¦æ­¢ã¾ãªã„ æ†§ã‚Œã¦ã„ãŸ æ†§ã‚Œã¦ä»•æ–¹ãŒãªã„",
    # "æ‰‹è¡“": "åè© + ã‚’ + å‹•è©ï¼ˆç›®çš„èªï¼‹å‹•è©å‹ï¼‰ æ‰‹è¡“ã‚’å—ã‘ã‚‹ æ‰‹è¡“ã‚’ã™ã‚‹ æ‰‹è¡“ã‚’è¡Œã† æ‰‹è¡“ã‚’æ–½ã™ æ‰‹è¡“ã‚’ä¸­æ­¢ã™ã‚‹ åè© + ã® + åè©ï¼ˆæ‰‹è¡“ã«é–¢ã‚ã‚‹å±æ€§ãƒ»é–¢ä¿‚èªï¼‰ æ‰‹è¡“æ—¥ æ‰‹è¡“ä¸­ æ‰‹è¡“å¾Œ æ‰‹è¡“å‰ æ‰‹è¡“å®¤ å½¢å®¹è© + æ‰‹è¡“ ï¼ˆæ‰‹è¡“ã‚’ä¿®é£¾ã™ã‚‹èªï¼‰ ç·Šæ€¥æ‰‹è¡“ å¤§æ‰‹è¡“ å°æ‰‹è¡“ å¤–ç§‘æ‰‹è¡“ å†…è¦–é¡æ‰‹è¡“ å‹•è© + ã®æ‰‹è¡“ / ï½ã«ã‚ˆã‚‹æ‰‹è¡“ ï½ã«ã‚ˆã‚‹æ‰‹è¡“ ï½ã‚’ç›®çš„ã¨ã—ãŸæ‰‹è¡“ ï½ã‚’ä¼´ã†æ‰‹è¡“ ï½å¾Œã®æ‰‹è¡“ ï½å‰ã®æ‰‹è¡“ æ‰‹è¡“ + ã«é–¢ã™ã‚‹èªå¥ æ‰‹è¡“æˆåŠŸ æ‰‹è¡“å¤±æ•— æ‰‹è¡“åˆä½µç—‡ æ‰‹è¡“çµŒé æ‰‹è¡“è¨ˆç”»",
    # "æ•™ç§‘æ›¸": "åè© + ã® + æ•™ç§‘æ›¸ å›½èªã®æ•™ç§‘æ›¸ æ•°å­¦ã®æ•™ç§‘æ›¸ è‹±èªã®æ•™ç§‘æ›¸ ç†ç§‘ã®æ•™ç§‘æ›¸ æ­´å²ã®æ•™ç§‘æ›¸ æ•™ç§‘æ›¸ + ã® + åè© æ•™ç§‘æ›¸ä»£ æ•™ç§‘æ›¸ä»£é‡‘ æ•™ç§‘æ›¸ç·¨é›† æ•™ç§‘æ›¸æœ¬æ–‡ æ•™ç§‘æ›¸æ¡æŠ å‹•è© + æ•™ç§‘æ›¸ æ•™ç§‘æ›¸ã‚’ä½¿ã† æ•™ç§‘æ›¸ã‚’èª­ã‚€ æ•™ç§‘æ›¸ã‚’è²·ã† æ•™ç§‘æ›¸ã‚’é…å¸ƒã™ã‚‹ æ•™ç§‘æ›¸ã‚’æ”¹è¨‚ã™ã‚‹ æ•™ç§‘æ›¸ + ã‚’ + å‹•è© æ•™ç§‘æ›¸ã‚’è£œè¶³ã™ã‚‹ æ•™ç§‘æ›¸ã‚’æŒå‚ã™ã‚‹ æ•™ç§‘æ›¸ã‚’æå‡ºã™ã‚‹ æ•™ç§‘æ›¸ã‚’è¿”å´ã™ã‚‹ æ•™ç§‘æ›¸ã‚’æŒã¡é‹ã¶ å½¢å®¹è© + æ•™ç§‘æ›¸ æ–°ã—ã„æ•™ç§‘æ›¸ åˆ†ã‹ã‚Šã‚„ã™ã„æ•™ç§‘æ›¸ åšã„æ•™ç§‘æ›¸ å°å‹ã®æ•™ç§‘æ›¸ é«˜ä¾¡ãªæ•™ç§‘æ›¸",
    # "æ–°é®®": "å½¢å®¹è© + åè© ï¼ˆã€Œæ–°é®®ãªã€œã€ã¨ã„ã†å½¢ã§åè©ã‚’ä¿®é£¾ã™ã‚‹ï¼‰ æ–°é®®ãªç©ºæ°— æ–°é®®ãªé‡èœ æ–°é®®ãªé­š æ–°é®®ãªæ„Ÿè¦š æ–°é®®ãªå°è±¡ åè© + ãŒ + å½¢å®¹è© ï¼ˆä¸»èªãŒã€Œæ–°é®®ã€ã§ã‚ã‚‹ï¼‰ ç©ºæ°—ãŒæ–°é®®ã  é­šãŒæ–°é®®ã  æ„Ÿè¦šãŒæ–°é®®ã  é¦™ã‚ŠãŒæ–°é®®ã  å‡ºæ¥äº‹ãŒæ–°é®®ã  å‹•è© + æ–°é®®ï¼ˆã«ï¼ã ï¼‰ ï¼ˆå‹•ä½œã¨çµã³ã¤ãè¡¨ç¾ï¼‰ æ„Ÿã˜ãŒæ–°é®®ã  è¦‹ã‚‹ã¨æ–°é®®ã  æ€ã„ã¤ããŒæ–°é®®ã  æ°—åˆ†ãŒæ–°é®®ã«ãªã‚‹ ç›®ã«æ–°é®®ã  æ–°é®® + åè©ï¼ˆã€Œã€œã®æ–°é®®ã•ã€ãªã©ï¼‰ ï¼ˆæ–°é®®ã•ã€æ€§è³ªã‚’è¡¨ã™è¡¨ç¾ï¼‰ æ–°é®®ã• æ–°é®®å‘³ æ–°é®®æ„Ÿ æ–°é®®åº¦ æ–°é®®æ€§ å‹•è© + ã® + æ–°é®®ï¼ˆã‚’ï¼ãŒï¼‰ ï¼ˆã€Œã€œã®æ–°é®®ã‚’æ„Ÿã˜ã‚‹ã€ãªã©ï¼‰ å‘³ã®æ–°é®®ã‚’æ„Ÿã˜ã‚‹ ç©ºæ°—ã®æ–°é®®ã‚’å‘³ã‚ã† æ–°é®®ã®æ„Ÿè¦šã‚’å¾—ã‚‹ æ–°é®®ã®æ¯å¹ã‚’æ„Ÿã˜ã‚‹ æ–°é®®ã®æ„Ÿå‹•ã‚’è¦šãˆã‚‹",
    # "æ—¥å¸¸": "åè© + ã® + æ—¥å¸¸ å¹³å‡¡ã®æ—¥å¸¸ å¿™ã—ã„æ—¥å¸¸ å®‰å®šã—ãŸæ—¥å¸¸ æ™®é€šã®æ—¥å¸¸ æ…£ã‚ŒãŸæ—¥å¸¸ æ—¥å¸¸ + åè© æ—¥å¸¸ç”Ÿæ´» æ—¥å¸¸ä¼šè©± æ—¥å¸¸æ¥­å‹™ æ—¥å¸¸é¢¨æ™¯ æ—¥å¸¸æ„Ÿè¦š å‹•è© + æ—¥å¸¸ æ—¥å¸¸ã‚’é€ã‚‹ æ—¥å¸¸ã‚’å–¶ã‚€ æ—¥å¸¸ã‚’éã”ã™ æ—¥å¸¸ã«æˆ»ã‚‹ æ—¥å¸¸ã«æº¶ã‘è¾¼ã‚€ æ—¥å¸¸ + ã« + å‹•è©ï¼å½¢å®¹è© æ—¥å¸¸ã«æˆ»ã‚‹ æ—¥å¸¸ã«éã”ã™ æ—¥å¸¸ã«æ¬ ã‹ã›ãªã„ æ—¥å¸¸ã«åˆºæ¿€ã‚’æ±‚ã‚ã‚‹ æ—¥å¸¸ã«å¤‰åŒ–ã‚’ã¤ã‘ã‚‹ å½¢å®¹è© + æ—¥å¸¸ ä½•æ°—ãªã„æ—¥å¸¸ å¿™ã—ã„æ—¥å¸¸ å˜èª¿ãªæ—¥å¸¸ å¹³å‡¡ãªæ—¥å¸¸ å¹¸ã›ãªæ—¥å¸¸",
    # "æ˜‡ã‚‹": "åè© + ãŒ / ã‚’ + æ˜‡ã‚‹ ï¼ˆä¸»èªãŒæ˜‡ã‚‹ã€å¯¾è±¡ã‚’é€šã˜ã¦æ˜‡ã‚‹ï¼‰ å¤ªé™½ãŒæ˜‡ã‚‹ æœˆãŒæ˜‡ã‚‹ ç…™ãŒæ˜‡ã‚‹ æ°—çƒãŒæ˜‡ã‚‹ éœ§ãŒæ˜‡ã‚‹ æ˜‡ã‚‹ + ç›®çš„èª / åˆ°é”ç‚¹ ï¼ˆæ˜‡ã£ã¦ã„ãå¯¾è±¡ãƒ»åˆ°é”å…ˆï¼‰ åœ°ä½ã«æ˜‡ã‚‹ ä½ãŒæ˜‡ã‚‹ æ˜‡çµ¦ã™ã‚‹ / æ˜‡çµ¦ã‚’æ˜‡ã‚‹ï¼ˆæ˜‡çµ¦ã‚’ä½¿ã†ã“ã¨ãŒå¤šã„ï¼‰ æ˜‡é€²ã™ã‚‹ / è·ã«æ˜‡ã‚‹ éšæ®µã‚’æ˜‡ã‚‹ï¼ˆæ©Ÿæ¢°ãªã©ã‚’ç”¨ã„ã‚‹æ–‡è„ˆã§ï¼‰ ï½ã«æ˜‡ã‚‹ ï¼ˆã€œã¨ã„ã†æ®µéšãƒ»å ´æ‰€ã«æ˜‡ã‚‹æ„å‘³ï¼‰ ã‚¨ãƒ¬ãƒ™ãƒ¼ã‚¿ãƒ¼ã«æ˜‡ã‚‹ å±•æœ›å°ã«æ˜‡ã‚‹ å¤©ã«æ˜‡ã‚‹ ç¥æ®¿ã«æ˜‡ã‚‹ èˆå°ã«æ˜‡ã‚‹ åè© + ã® + æ˜‡ã‚‹ / ï½ã®æ˜‡ã‚‹ ï¼ˆåè©ï¼‹ã€Œæ˜‡ã‚‹ã€ã‚’çµ„ã¿è¾¼ã‚€è¡¨ç¾ï¼‰ æ˜‡ç«œ ä¸Šæ˜‡ï¼ˆã€Œæ˜‡ã€ã®èªã‚’ä½¿ã£ãŸæ´¾ç”Ÿèªï¼‰ æ˜‡å¤© æ˜‡é™ï¼ˆæ˜‡ã‚‹ï¼‹é™ã‚Šã‚‹ã®è¤‡åˆèªï¼‰ æ˜‡æ ¼ è¡¨ç¾ãƒ»æ…£ç”¨å¥ ï¼ˆã€Œæ˜‡ã‚‹ã€ã‚’å«ã‚€å®šå‹è¡¨ç¾ï¼‰ å¤©ã«ã‚‚æ˜‡ã‚‹æ°—æŒã¡ æœæ—¥ãŒæ˜‡ã‚‹ å¤•æ—¥ãŒæ˜‡ã‚‹ï¼ˆã‚ã¾ã‚Šä¸€èˆ¬çš„ã§ã¯ãªã„ãŒè©©çš„è¡¨ç¾ï¼‰ æ˜‡çµ¦ãƒ»æ˜‡é€² æ˜‡æ®µï¼ˆæ­¦é“ãªã©ã§éšç´šãŒä¸ŠãŒã‚‹ï¼‰",
    # "æš®ã‚‰ã—": "åè© + ã® + æš®ã‚‰ã— å¹³ç©ãªæš®ã‚‰ã— è±Šã‹ãªæš®ã‚‰ã— è³ªç´ ãªæš®ã‚‰ã— å®‰å®šã—ãŸæš®ã‚‰ã— å¿«é©ãªæš®ã‚‰ã— æš®ã‚‰ã— + ã® + åè© æš®ã‚‰ã—å‘ã æš®ã‚‰ã—ã¶ã‚Š æš®ã‚‰ã—ã¶ã‚Šã‚’è¦‹ã›ã‚‹ æš®ã‚‰ã—ã¶ã‚Šã‚’ä¼ãˆã‚‹ æš®ã‚‰ã—å‘ãã‚’æ”¹å–„ã™ã‚‹ å‹•è© + æš®ã‚‰ã— æš®ã‚‰ã—ã‚’å–¶ã‚€ æš®ã‚‰ã—ã‚’æ”¯ãˆã‚‹ æš®ã‚‰ã—ã‚’å®ˆã‚‹ æš®ã‚‰ã—ã‚’æ¥½ã—ã‚€ æš®ã‚‰ã—ã‚’ç«‹ã¦ã‚‹ æš®ã‚‰ã— + ã‚’ + å‹•è© æš®ã‚‰ã—ã‚’æ”¹å–„ã™ã‚‹ æš®ã‚‰ã—ã‚’è±Šã‹ã«ã™ã‚‹ æš®ã‚‰ã—ã‚’å¤‰ãˆã‚‹ æš®ã‚‰ã—ã‚’æ”¯ãˆã‚‹ æš®ã‚‰ã—ã‚’åˆ‡ã‚Šè©°ã‚ã‚‹ å½¢å®¹è© + æš®ã‚‰ã— é•·ã„æš®ã‚‰ã— æš®ã‚‰ã—ã‚„ã™ã„ æš®ã‚‰ã—ã«ãã„ æš®ã‚‰ã—å‘ãè‰¯ã„ æš®ã‚‰ã—å‘ãæ‚ªã„",
    # "æ§‹ã†": "å‹•è© + ã€Œæ§‹ã†ã€ æ°—ã« æ§‹ã† ã‚ã¾ã‚Š æ§‹ã† å…¨ã æ§‹ã‚ãªã„ï¼æ§‹ã†ã¾ã„ å¹²æ¸‰ã—ã¦ æ§‹ã† ä¸–è©±ã‚’ æ§‹ã† ã€Œæ§‹ã†ã€ + åè©ãƒ»å¥ æ§‹ã† äºº æ§‹ã† æ…‹åº¦ æ§‹ã† å¿…è¦ æ§‹ã† ä¾¡å€¤ æ§‹ã† æ™‚é–“ å‰¯è© + ã€Œæ§‹ã†ã€ å…¨ç„¶ æ§‹ã‚ãªã„ å°‘ã—ã‚‚ æ§‹ã‚ãªã„ ã¾ã£ãŸã æ§‹ã‚ãªã„ æ±ºã—ã¦ æ§‹ã†ãª åˆ¥ã« æ§‹ã‚ãªã„ ã€Œã€œã«æ§‹ã†ã€ ï¼ˆå¯¾è±¡ã‚’ç¤ºã™æ ¼åŠ©è©ã€Œã«ã€ã‚’ä¼´ã†ã‚±ãƒ¼ã‚¹ï¼‰ äºº ã«æ§‹ã† å­ã©ã‚‚ ã«æ§‹ã† ä»–äºº ã«æ§‹ã† çŠ¬ ã«æ§‹ã† å­˜åœ¨ ã«æ§‹ã†",
    # "æ­¦å™¨": "åè© + ã€Œæ­¦å™¨ã€ ä¸»è¦ æ­¦å™¨ é‡ç«å™¨ æ­¦å™¨ é è·é›¢ æ­¦å™¨ è¿‘æ¥ æ­¦å™¨ æ ¸ æ­¦å™¨ ã€Œæ­¦å™¨ã€ + åè©ãƒ»å¥ æ­¦å™¨ ã¨ã—ã¦ æ­¦å™¨ ä½¿ç”¨ æ­¦å™¨ ä¿æœ‰ æ­¦å™¨ æ‹¡æ•£ æ­¦å™¨ è£½é€  å‹•è© + ã€Œæ­¦å™¨ã€ æ­¦å™¨ã‚’ æŒã¤ æ­¦å™¨ã‚’ æºãˆã‚‹ æ­¦å™¨ã‚’ ä½¿ç”¨ã™ã‚‹ æ­¦å™¨ã‚’ æ”¾æ£„ã™ã‚‹ æ­¦å™¨ã‚’ æ”¹è‰¯ã™ã‚‹ å‰¯è©ãƒ»ä¿®é£¾èª + ã€Œæ­¦å™¨ã€ å¼·åŠ›ãª æ­¦å™¨ è‡´å‘½çš„ãª æ­¦å™¨ åŠ¹æœçš„ãª æ­¦å™¨ å±é™ºãª æ­¦å™¨ é«˜æ€§èƒ½ãª æ­¦å™¨",
    # "æ­©é“": "å½¢å®¹è© + æ­©é“ ç‹­ã„ æ­©é“ åºƒã„ æ­©é“ å®‰å…¨ãª æ­©é“ æ­©ãã‚„ã™ã„ æ­©é“ è»Šé“ã¨åˆ†é›¢ã•ã‚ŒãŸ æ­©é“ åè© + æ­©é“ æ­©é“ å¹… æ­©é“ è¨­ç½® æ­©é“ æ•´å‚™ æ­©é“ èˆ—è£… æ­©é“ åŒºé–“ å‹•è© + æ­©é“ æ­©é“ã‚’ æ­©ã æ­©é“ã‚’ æ•´å‚™ã™ã‚‹ æ­©é“ã‚’ æ‹¡å¼µã™ã‚‹ æ­©é“ã‚’ ç¢ºä¿ã™ã‚‹ æ­©é“ã‚’ é€šè¡Œã™ã‚‹ ã€Œã€œã®æ­©é“ã€ã€Œæ­©é“ã®ã€œã€ éƒ½å¸‚ ã®æ­©é“ å­¦æ ¡ ã®æ­©é“ å•†åº—è¡— ã®æ­©é“ æ­©é“ ã®æ®µå·® æ­©é“ ã®ç¸çŸ³",
    # "æ¶ˆé˜²": "åè© + ã€Œæ¶ˆé˜²ã€ æ¶ˆé˜²ç½² æ¶ˆé˜²å£« æ¶ˆé˜²è»Š æ¶ˆé˜²å›£ æ¶ˆé˜²æ³• ã€Œæ¶ˆé˜²ã€ + åè©ãƒ»å¥ æ¶ˆé˜² æ´»å‹• æ¶ˆé˜² è¨­å‚™ æ¶ˆé˜² è¨“ç·´ æ¶ˆé˜² è¨ˆç”» æ¶ˆé˜² çµ„ç¹” å‹•è© + ã€Œæ¶ˆé˜²ã€ æ¶ˆé˜²ã‚’ å¼·åŒ–ã™ã‚‹ æ¶ˆé˜²ã‚’ æ•´å‚™ã™ã‚‹ æ¶ˆé˜²ã‚’ æ”¯æ´ã™ã‚‹ æ¶ˆé˜²ã‚’ æ‹¡å……ã™ã‚‹ æ¶ˆé˜²ã‚’ ç®¡ç†ã™ã‚‹ ã€Œã€œã®æ¶ˆé˜²ã€ï¼ã€Œæ¶ˆé˜²ã®ã€œã€ åœ°åŸŸ ã®æ¶ˆé˜² å¸‚ç”ºæ‘ ã®æ¶ˆé˜² éƒ½å¸‚ ã®æ¶ˆé˜² æ¶ˆé˜² ã®è²¬ä»» æ¶ˆé˜² ã®å½¹å‰²",
    # "æ¸…æ½”": "å½¢å®¹è© + ã€Œæ¸…æ½”ã€ æ¸…æ½” ãª è¡£æœ æ¸…æ½” ãª éƒ¨å±‹ æ¸…æ½” ãª ç’°å¢ƒ æ¸…æ½” ãª å°æ‰€ æ¸…æ½” ãª ç©ºæ°— åè© + ã€Œæ¸…æ½”ã€ æ¸…æ½” æ„Ÿ æ¸…æ½” åº¦ æ¸…æ½” ã• æ¸…æ½” åŸºæº– æ¸…æ½” æ„è­˜ å‹•è© + ã€Œæ¸…æ½”ã€ æ¸…æ½”ã‚’ ä¿ã¤ æ¸…æ½”ã« ã™ã‚‹ æ¸…æ½”ã‚’ ä¿ã¦ã‚‹ æ¸…æ½”ã‚’ å¿ƒãŒã‘ã‚‹ æ¸…æ½”ã‚’ ç¶­æŒã™ã‚‹ ã€Œã€œã®æ¸…æ½”ã€ï¼ã€Œæ¸…æ½”ã®ã€œã€ ä½“ ã®æ¸…æ½” æ‰‹ ã®æ¸…æ½” è¡£æœ ã®æ¸…æ½” é£Ÿå™¨ ã®æ¸…æ½” å®¤å†… ã®æ¸…æ½”",
    # "æ¸›ã‚‹": "åè© + æ¸›ã‚‹ äººå£ãŒ æ¸›ã‚‹ å£²ä¸ŠãŒ æ¸›ã‚‹ åå…¥ãŒ æ¸›ã‚‹ æ”¯å‡ºãŒ æ¸›ã‚‹ æ•°ãŒ æ¸›ã‚‹ ã€Œã€œãŒæ¸›ã‚‹ã€ ï¼ˆä¸»èªã‚’è¡¨ã™ã€ŒãŒã€ã‚’ä¼´ã†ãƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ æ£®æ—ãŒ æ¸›ã‚‹ æ£®ãŒ æ¸›ã‚‹ æ£®æ—é¢ç©ãŒ æ¸›ã‚‹ é£²é…’é‡ãŒ æ¸›ã‚‹ ä½¿ç”¨é‡ãŒ æ¸›ã‚‹ å‰¯è© + æ¸›ã‚‹ / æ¸›ã£ã¦ã„ã‚‹ ã ã‚“ã ã‚“ æ¸›ã‚‹ å¾ã€…ã« æ¸›ã‚‹ å¤§å¹…ã« æ¸›ã‚‹ æ€¥æ¿€ã« æ¸›ã‚‹ è‘—ã—ã æ¸›ã£ã¦ã„ã‚‹ ã€Œã€œã‚’æ¸›ã‚‰ã™ã€å¯¾å¿œï¼ˆé–¢é€£ã‚³ãƒ­ã‚±ãƒ¼ã‚·ãƒ§ãƒ³ï¼‰ ï¼ˆã€Œæ¸›ã‚‹ã€ã¯è‡ªå‹•è©ã ãŒã€ä»–å‹•è©ã€Œæ¸›ã‚‰ã™ã€ã¨å¯¾ã«ãªã‚‹è¡¨ç¾ã‚‚åˆã‚ã›ã¦çŸ¥ã£ã¦ãŠãã¨ä¾¿åˆ©ï¼‰ ã‚³ã‚¹ãƒˆã‚’ æ¸›ã‚‰ã™ æ”¯å‡ºã‚’ æ¸›ã‚‰ã™ é£Ÿè²»ã‚’ æ¸›ã‚‰ã™ äººå“¡ã‚’ æ¸›ã‚‰ã™ ã‚´ãƒŸã‚’ æ¸›ã‚‰ã™",
    # "ç‰©èª": "åè© + ç‰©èª é•·ç·¨ç‰©èª çŸ­ç·¨ç‰©èª å†’é™ºç‰©èª æ‹æ„›ç‰©èª ä¼èª¬ç‰©èª ã€Œç‰©èªã€ + åè©ãƒ»å¥ ç‰©èª é¢¨æ™¯ ç‰©èª ä¸–ç•Œ ç‰©èª èˆå° ç‰©èª çµæœ« ç‰©èª èªã‚Šæ‰‹ å‹•è© + ç‰©èª ç‰©èªã‚’ ç´¡ã ç‰©èªã‚’ æã ç‰©èªã‚’ èªã‚‹ ç‰©èªãŒ é€²ã‚€ ç‰©èªãŒ å±•é–‹ã™ã‚‹ å½¢å®¹è©ãƒ»ä¿®é£¾èª + ç‰©èª æ„Ÿå‹•çš„ãª ç‰©èª é‡åšãª ç‰©èª ç¾ã—ã„ ç‰©èª å¤å…¸çš„ãª ç‰©èª å¹»æƒ³çš„ãª ç‰©èª ã€Œã€œã®ç‰©èªã€ã€Œç‰©èªã®ã€œã€ äººç”Ÿ ã®ç‰©èª å°‘å¹´ ã®ç‰©èª å®¶æ— ã®ç‰©èª ç‰©èª ã®å§‹ã¾ã‚Š ç‰©èª ã®çµæœ«",
    # "ç¾ã‚Œã‚‹": "åè© + ã€Œç¾ã‚Œã‚‹ã€ å½±ãŒç¾ã‚Œã‚‹ ç‰©éŸ³ãŒç¾ã‚Œã‚‹ ç¾è±¡ãŒç¾ã‚Œã‚‹ ç—‡çŠ¶ãŒç¾ã‚Œã‚‹ æœ¬æ€§ãŒç¾ã‚Œã‚‹ ã€Œã€œã« / ã€œã‹ã‚‰ã€ + ã€Œç¾ã‚Œã‚‹ã€ ï¼ˆå ´æ‰€ãƒ»èµ·ç‚¹ã‚’è¡¨ã™å‰¯è©çš„èªã‚’ä¼´ã†ï¼‰ é—‡ã«ç¾ã‚Œã‚‹ ç©ºã‹ã‚‰ç¾ã‚Œã‚‹ å±±ä¸­ã«ç¾ã‚Œã‚‹ æ£®ã®ä¸­ã‹ã‚‰ç¾ã‚Œã‚‹ èˆå°ã«ç¾ã‚Œã‚‹ å‰¯è© / å‰¯è©å¥ + ã€Œç¾ã‚Œã‚‹ã€ ï¼ˆå‹•ä½œã‚’ä¿®é£¾ã™ã‚‹ï¼‰ çªç„¶ç¾ã‚Œã‚‹ æ€¥ã«ç¾ã‚Œã‚‹ å¾ã€…ã«ç¾ã‚Œã‚‹ ã‚†ã£ãã‚Šç¾ã‚Œã‚‹ ã±ã£ã¨ç¾ã‚Œã‚‹ å½¢å®¹è© / å½¢å®¹è©å¥ + ã€Œç¾ã‚Œã‚‹ã€ ï¼ˆã©ã®ã‚ˆã†ã«è¦‹ãˆã‚‹ã‹ã‚’å¼·èª¿ï¼‰ ã¯ã£ãã‚Šç¾ã‚Œã‚‹ ã¼ã‚“ã‚„ã‚Šç¾ã‚Œã‚‹ å¾®ã‹ã«ç¾ã‚Œã‚‹ é®®æ˜ã«ç¾ã‚Œã‚‹ æ˜ç¢ºã«ç¾ã‚Œã‚‹ å‹•è©å¥ + ã€Œç¾ã‚Œã‚‹ã€ ï¼ˆå‹•è©å¥ã§æ¡ä»¶ã‚„å¤‰åŒ–ã‚’è¡¨ã™ï¼‰ å§¿ã‚’ç¾ã‚Œã‚‹ â†’ å§¿ã‚’ç¾ã™ï¼ˆâ€»ãŸã ã—ã€Œç¾ã™ã€ãŒä½¿ã‚ã‚Œã‚‹æ–¹ãŒä¸€èˆ¬çš„ï¼‰ è¡¨æƒ…ãŒç¾ã‚Œã‚‹ è¡¨æƒ…ã«ç¾ã‚Œã‚‹ æ°—æŒã¡ãŒç¾ã‚Œã‚‹ è‘—ã—ãç¾ã‚Œã‚‹",
    # "ç™ºè¡Œ": "åè© + ã€Œç™ºè¡Œã€ é›‘èªŒã‚’ç™ºè¡Œ åˆŠè¡Œç‰©ã‚’ç™ºè¡Œ å‚µåˆ¸ã‚’ç™ºè¡Œ å…è¨±ã‚’ç™ºè¡Œ è¨¼æ›¸ã‚’ç™ºè¡Œ ã€Œã€œã‚’ï¼ã€œãŒã€ + ã€Œç™ºè¡Œã™ã‚‹ï¼ç™ºè¡Œã•ã‚Œã‚‹ã€ ï¼ˆä»–å‹•ãƒ»è‡ªå‹•ã®æ–‡å‹ã§ä½¿ã‚ã‚Œã‚‹ï¼‰ åˆ¸ã‚’ç™ºè¡Œã™ã‚‹ å…è¨±ãŒç™ºè¡Œã•ã‚Œã‚‹ é€šè²¨ã‚’ç™ºè¡Œã™ã‚‹ è¨¼æ˜æ›¸ã‚’ç™ºè¡Œã™ã‚‹ ãƒã‚±ãƒƒãƒˆã‚’ç™ºè¡Œã™ã‚‹ å‰¯è© / å‰¯è©å¥ + ã€Œç™ºè¡Œã€ ï¼ˆç™ºè¡Œã®æ§˜æ…‹ã‚’è¡¨ã™å‰¯è©çš„èªï¼‰ å®šæœŸçš„ã«ç™ºè¡Œã™ã‚‹ éšæ™‚ç™ºè¡Œã™ã‚‹ æ–°ãŸã«ç™ºè¡Œã™ã‚‹ å¤§é‡ã«ç™ºè¡Œã™ã‚‹ å†åº¦ç™ºè¡Œã™ã‚‹ å½¢å®¹è© / å½¢å®¹è©å¥ + ã€Œç™ºè¡Œã€ ï¼ˆç™ºè¡Œã«é–¢ã™ã‚‹æ€§è³ªã‚’è¡¨ã™èªï¼‰ æ­£å¼ã«ç™ºè¡Œã™ã‚‹ ç„¡æ–™ã§ç™ºè¡Œã™ã‚‹ è¿…é€Ÿã«ç™ºè¡Œã™ã‚‹ é©åˆ‡ã«ç™ºè¡Œã™ã‚‹ æ­£ç¢ºã«ç™ºè¡Œã™ã‚‹ å‹•è©å¥ + ã€Œç™ºè¡Œã™ã‚‹ã€ ï¼ˆå‹•è©å¥ã¨çµã³ã¤ã„ã¦ä½¿ã‚ã‚Œã‚‹è¡¨ç¾ï¼‰ ç™ºè¡Œæ‰‹ç¶šãã‚’ã™ã‚‹ ç™ºè¡Œé–‹å§‹ã™ã‚‹ ç™ºè¡Œåœæ­¢ã™ã‚‹ ç™ºè¡Œç”³è«‹ã™ã‚‹ ç™ºè¡Œè¨±å¯ã™ã‚‹",
    # "çœŸã£èµ¤": "å½¢å®¹è© + åè©ï¼ˆâ€œçœŸã£èµ¤ãª Xâ€ï¼‰ çœŸã£èµ¤ãªèŠ± çœŸã£èµ¤ãªå¤•ç„¼ã‘ çœŸã£èµ¤ãªé¡” çœŸã£èµ¤ãªãƒªãƒ³ã‚´ çœŸã£èµ¤ãªæ—— åè© + ã€ŒçœŸã£èµ¤ã«ãªã‚‹ï¼çœŸã£èµ¤ã«æŸ“ã¾ã‚‹ï¼çœŸã£èµ¤ã«ã™ã‚‹ã€ãªã©ã®å‹•è©å¥ é¡”ãŒçœŸã£èµ¤ã«ãªã‚‹ ç©ºãŒçœŸã£èµ¤ã«æŸ“ã¾ã‚‹ é¡”ã‚’çœŸã£èµ¤ã«ã™ã‚‹ è‘‰ãŒçœŸã£èµ¤ã«æŸ“ã¾ã‚‹ è¡€ãŒçœŸã£èµ¤ã«ãªã‚‹ å‰¯è© / å‰¯è©å¥ + ã€ŒçœŸã£èµ¤ã«ï¼çœŸã£èµ¤ã§ã€ãªã© ã¨ã¦ã‚‚çœŸã£èµ¤ã« ã˜ã‚ã˜ã‚çœŸã£èµ¤ã« ã™ã£ã‹ã‚ŠçœŸã£èµ¤ã« ã‚ã£ãã‚ŠçœŸã£èµ¤ã« å®Œå…¨ã«çœŸã£èµ¤ã§ å‹•è©å¥ + ã€ŒçœŸã£èµ¤ã€ï¼‹ è£œèªï¼ˆçŠ¶æ…‹ã‚’è¡¨ã™èªï¼‰ èµ¤ãçœŸã£èµ¤ï¼ˆã«ãªã‚‹ï¼‰ æŸ“ã‚ã¦çœŸã£èµ¤ ç„¼ã‘ã¦çœŸã£èµ¤ ç‡ƒãˆã¦çœŸã£èµ¤ è½ã¡ã¦çœŸã£èµ¤ æ¯”å–©çš„ç”¨æ³•ï¼è»¢ç”¨è¡¨ç¾ + ã€ŒçœŸã£èµ¤ã€ ï¼ˆæ¯”å–©ãƒ»å¼·èª¿è¡¨ç¾ï¼‰ çœŸã£èµ¤ãªå˜˜ çœŸã£èµ¤ãªã‚¦ã‚½ çœŸã£èµ¤ãªå˜˜ã‚’ã¤ã çœŸã£èµ¤ãªç½ª çœŸã£èµ¤ãªé–“é•ã„",
    # "ç¡çœ ": "åè© + ã€Œã‚’ï¼ãŒã€ + å‹•è©ãƒ‘ã‚¿ãƒ¼ãƒ³ï¼ˆç¡çœ ã¨å‹•è©ã®çµã³ã¤ãï¼‰ ç¡çœ ã‚’å–ã‚‹ ç¡çœ ã‚’ã¨ã‚‹ ç¡çœ ã‚’ã¨ã‚‹ã‚ˆã†ã«ã™ã‚‹ ç¡çœ ã‚’ç¢ºä¿ã™ã‚‹ ç¡çœ ã‚’å¦¨ã’ã‚‹ ç¡çœ ã‚’æ”¹å–„ã™ã‚‹ ç¡çœ ãŒæµ…ã„ ç¡çœ ãŒæ·±ã„ ç¡çœ ãŒä¹±ã‚Œã‚‹ ç¡çœ ãŒè¶³ã‚Šãªã„ å‰¯è© / å‰¯è©å¥ + ã€Œç¡çœ ã€ + å‹•è©ï¼è¿°èªãƒ‘ã‚¿ãƒ¼ãƒ³ ååˆ†ãªç¡çœ ã‚’å–ã‚‹ é©åˆ‡ãªç¡çœ ã‚’ç¢ºä¿ã™ã‚‹ è‰¯è³ªãªç¡çœ ã‚’ã¨ã‚‹ é•·æ™‚é–“ã®ç¡çœ ã‚’ã¨ã‚‹ ååˆ†ãªç¡çœ ãŒå–ã‚Œã¦ã„ãªã„ è³ªã®é«˜ã„ç¡çœ ã‚’å¾—ã‚‹ ã€Œç¡çœ ã€ + åŠ©è© + åè©ï¼å½¢å®¹è©ãƒ‘ã‚¿ãƒ¼ãƒ³ ç¡çœ æ™‚é–“ ç¡çœ ä¸è¶³ ç¡çœ éšœå®³ ç¡çœ çŠ¶æ…‹ ç¡çœ ç¿’æ…£ å½¢å®¹è© + ã€Œç¡çœ ã€ è‰¯è³ªãªç¡çœ  æ·±ã„ç¡çœ  å¿«é©ãªç¡çœ  ååˆ†ãªç¡çœ  ç¡çœ ä¸è¶³ã® å‹•è©å¥ + ã€Œç¡çœ ã€ + è£œèªï¼ä¿®é£¾èªï¼ˆç¡çœ ã®çŠ¶æ…‹ã‚’è¡¨ã™èªã‚’ä¼´ã†è¡¨ç¾ï¼‰ ç¡çœ ãŒå¦¨ã’ã‚‰ã‚Œã‚‹ ç¡çœ ãŒå¦¨å®³ã•ã‚Œã‚‹ ç¡çœ ã®è³ªãŒä½ä¸‹ã™ã‚‹ ç¡çœ ã®è³ªã‚’é«˜ã‚ã‚‹ ç¡çœ ã®ã‚µã‚¤ã‚¯ãƒ«ã‚’æ•´ãˆã‚‹",
    # "ç¢ºã‹ã‚ã‚‹": "åè© + ã€Œã‚’ã€ + ã€Œç¢ºã‹ã‚ã‚‹ã€ èº«åˆ†è¨¼æ˜æ›¸ã‚’ç¢ºã‹ã‚ã‚‹ ä½æ‰€ã‚’ç¢ºã‹ã‚ã‚‹ é›»è©±ç•ªå·ã‚’ç¢ºã‹ã‚ã‚‹ äºˆç´„å†…å®¹ã‚’ç¢ºã‹ã‚ã‚‹ æƒ…å ±ã‚’ç¢ºã‹ã‚ã‚‹ ã€Œã€œãŒã€ + ã€Œç¢ºã‹ã‚ã‚‰ã‚Œã‚‹ï¼ç¢ºã‹ã‚ã‚‹ã€ çœŸå½ãŒç¢ºã‹ã‚ã‚‰ã‚Œã‚‹ äº‹å®ŸãŒç¢ºã‹ã‚ã‚‰ã‚Œã‚‹ å®‰å…¨ãŒç¢ºã‹ã‚ã‚‰ã‚Œã‚‹ æ­£ç¢ºã•ãŒç¢ºã‹ã‚ã‚‰ã‚Œã‚‹ åŠ¹æœãŒç¢ºã‹ã‚ã‚‰ã‚Œã‚‹ å‰¯è© / å‰¯è©å¥ + ã€Œç¢ºã‹ã‚ã‚‹ã€ ãã¡ã‚“ã¨ç¢ºã‹ã‚ã‚‹ å¿µå…¥ã‚Šã«ç¢ºã‹ã‚ã‚‹ å†åº¦ç¢ºã‹ã‚ã‚‹ ã‚‚ã†ä¸€åº¦ç¢ºã‹ã‚ã‚‹ ã¡ã‚ƒã‚“ã¨ç¢ºã‹ã‚ã‚‹ å½¢å®¹è© / å½¢å®¹è©å¥ + ã€Œç¢ºã‹ã‚ã‚‹ã€ æ­£ã—ãç¢ºã‹ã‚ã‚‹ ååˆ†ã«ç¢ºã‹ã‚ã‚‹ å®Œå…¨ã«ç¢ºã‹ã‚ã‚‹ æ­£ç¢ºã«ç¢ºã‹ã‚ã‚‹ ç¢ºå®Ÿã«ç¢ºã‹ã‚ã‚‹ å‹•è©å¥ + ã€Œç¢ºã‹ã‚ã‚‹ã€ ç¢ºã‹ã‚ãŸã†ãˆã§ã€œã™ã‚‹ ç¢ºã‹ã‚ã‚‹æ‰‹ç¶šãã‚’ã™ã‚‹ ç¢ºã‹ã‚ã¦ãŠã ç¢ºã‹ã‚ã«è¡Œã ç¢ºã‹ã‚ç›´ã™",
    # "ç¤¼å„€": "åè© + ã€Œã®ã€ + ã€Œç¤¼å„€ã€ ç¤¼å„€ã®ä½œæ³• ç¤¼å„€ã®æ¬ å¦‚ ç¤¼å„€ã®å•é¡Œ ç¤¼å„€ã®åŸºæœ¬ ç¤¼å„€ã®ç¯„å›² ã€Œç¤¼å„€ã€ + åŠ©è© + åè©ï¼èªå¥ ç¤¼å„€ä½œæ³• ç¤¼å„€æ­£ã—ã• ç¤¼å„€æ­£ã—ã„æ…‹åº¦ ç¤¼å„€æ­£ã—ã„è¨€è‘‰é£ã„ ç¤¼å„€æ­£ã—ã„å¯¾å¿œ å‹•è© + ã€Œç¤¼å„€ã‚’ã€œï¼ç¤¼å„€ã‚’é‡ã‚“ã˜ã‚‹ï¼ç¤¼å„€ã‚’å®ˆã‚‹ã€ãªã© ç¤¼å„€ã‚’é‡ã‚“ã˜ã‚‹ ç¤¼å„€ã‚’å®ˆã‚‹ ç¤¼å„€ã‚’æ¬ ã ç¤¼å„€ã‚’ã‚ãã¾ãˆã‚‹ ç¤¼å„€ã‚’ç¤ºã™ å‰¯è©ï¼å‰¯è©å¥ + ã€Œç¤¼å„€æ­£ã—ãã€œã€ãªã© éå¸¸ã«ç¤¼å„€æ­£ã—ã ãã¡ã‚“ã¨ç¤¼å„€æ­£ã—ã ã‚‚ã£ã¨ç¤¼å„€æ­£ã—ã ã‚ã‚‹ç¨‹åº¦ç¤¼å„€æ­£ã—ã æ¯”è¼ƒçš„ç¤¼å„€æ­£ã—ã",
    # "ç©ã‚„ã‹": "å½¢å®¹è© + åè©ï¼ˆâ€œç©ã‚„ã‹ãª Ã—Ã—â€ï¼‰ ç©ã‚„ã‹ãªå¤©æ°— ç©ã‚„ã‹ãªæµ· ç©ã‚„ã‹ãªè¡¨æƒ… ç©ã‚„ã‹ãªç¬‘é¡” ç©ã‚„ã‹ãªæš®ã‚‰ã— åè© + ã€Œã§ï¼ãªï¼ã«ã€ + ã€Œç©ã‚„ã‹ã€ï¼å¤‰åŒ–è¡¨ç¾ å¿ƒãŒç©ã‚„ã‹ã§ã‚ã‚‹ æ°—æŒã¡ãŒç©ã‚„ã‹ã«ãªã‚‹ ç©ºæ°—ãŒç©ã‚„ã‹ã  æ—¥å·®ã—ãŒç©ã‚„ã‹ã  æ³¢é¢¨ãŒç©ã‚„ã‹ã  å‰¯è© / å‰¯è©å¥ + ã€Œç©ã‚„ã‹ã«ã€ ç©ã‚„ã‹ã«è©±ã™ ç©ã‚„ã‹ã«æš®ã‚‰ã™ ç©ã‚„ã‹ã«è¦‹å®ˆã‚‹ ç©ã‚„ã‹ã«é€²ã‚€ ç©ã‚„ã‹ã«éã”ã™ å‹•è©å¥ + ã€Œç©ã‚„ã‹ã€ï¼ã€Œç©ã‚„ã‹ã§ã‚ã‚‹ã€ ç©ã‚„ã‹ã•ã‚’æ„Ÿã˜ã‚‹ ç©ã‚„ã‹ã•ã‚’ä¿ã¤ ç©ã‚„ã‹ãªé›°å›²æ°—ã‚’é†¸ã™ ï½ã‚’ç©ã‚„ã‹ã«ã™ã‚‹ ï½ã‚’ç©ã‚„ã‹ã«å—ã‘æ­¢ã‚ã‚‹",
    # "ç«‹æ´¾": "å½¢å®¹è© + åè©ï¼ˆâ€œç«‹æ´¾ãª Ã—Ã—â€ï¼‰ ç«‹æ´¾ãªäººç‰© ç«‹æ´¾ãªå»ºç‰© ç«‹æ´¾ãªä»•äº‹ ç«‹æ´¾ãªæˆæœ ç«‹æ´¾ãªæ…‹åº¦ åè© + ã€Œã ï¼ã§ã‚ã‚‹ï¼ã«ã€ + ã€Œç«‹æ´¾ã€ï¼å½¢å®¹å‹•è©åŒ– äººã¯ç«‹æ´¾ã  å»ºç‰©ã¯ç«‹æ´¾ã§ã‚ã‚‹ å½¼ã¯ç«‹æ´¾ãªäººã  ãã‚Œã¯ç«‹æ´¾ãªã“ã¨ã  ç«‹æ´¾ãªã‚‚ã®ã§ã‚ã‚‹ å‹•è© + ã€Œç«‹æ´¾ã«ã€œã™ã‚‹ï¼ç«‹æ´¾ã«ãªã‚‹ã€ãªã© ç«‹æ´¾ã«è‚²ã¤ ç«‹æ´¾ã«æˆé•·ã™ã‚‹ ç«‹æ´¾ã«æŒ¯ã‚‹èˆã† ç«‹æ´¾ã«è¦‹ãˆã‚‹ ç«‹æ´¾ã«å‹™ã‚ã‚‹ å‰¯è© / å‰¯è©å¥ + ã€Œç«‹æ´¾ã«ã€ ã¨ã¦ã‚‚ç«‹æ´¾ã« å®Ÿã«ç«‹æ´¾ã« ãªã‹ãªã‹ç«‹æ´¾ã« ç«‹æ´¾ã«â€¦ã—ã¦ã„ã‚‹ è¦‹äº‹ã«ç«‹æ´¾ã« å‹•è©å¥ + ã€Œç«‹æ´¾ã€ + è£œèªï¼ä¿®é£¾èª ç«‹æ´¾ã•ã‚’æ„Ÿã˜ã‚‹ ç«‹æ´¾ã•ã‚’ç¤ºã™ ç«‹æ´¾ã•ã‚’å‚™ãˆã‚‹ ç«‹æ´¾ã•ã‚’å¤±ã† ç«‹æ´¾ãªæ§˜å­ã‚’å‘ˆã™ã‚‹",
    # "ç­‹è‚‰": "åè© + ã€Œã®ã€ + ã€Œç­‹è‚‰ã€ ç­‹è‚‰ã®åç¸® ç­‹è‚‰ã®é‡ ç­‹è‚‰ã®ç™ºé” ç­‹è‚‰ã®ç–²åŠ´ ç­‹è‚‰ã®å¼·ã• åè© + ã€Œç­‹è‚‰ã€ + åŠ©è© / å‹•è©ãƒ‘ã‚¿ãƒ¼ãƒ³ ç­‹è‚‰ã‚’é›ãˆã‚‹ ç­‹è‚‰ã‚’å¢—ã‚„ã™ ç­‹è‚‰ã‚’ã¤ã‘ã‚‹ ç­‹è‚‰ã‚’è½ã¨ã™ ç­‹è‚‰ãŒã¤ã ç­‹è‚‰ãŒæ¸›ã‚‹ ç­‹è‚‰ãŒã¤ã‚‰ã„ ç­‹è‚‰ãŒç–²ã‚Œã‚‹ å½¢å®¹è© / å½¢å®¹å‹•è© + ã€Œç­‹è‚‰ã€ å¤ªã„ç­‹è‚‰ ç´°ã„ç­‹è‚‰ å¼·ã„ç­‹è‚‰ å¼·é­ãªç­‹è‚‰ å¼±ã„ç­‹è‚‰ å‰¯è© / å‰¯è©å¥ + ã€Œç­‹è‚‰ã€ + å‹•è©ï¼å½¢å®¹èªãƒ‘ã‚¿ãƒ¼ãƒ³ ã—ã£ã‹ã‚Šç­‹è‚‰ã‚’é›ãˆã‚‹ åŠ¹ç‡ã‚ˆãç­‹è‚‰ã‚’ã¤ã‘ã‚‹ ç„¡ç†ã›ãšç­‹è‚‰ã‚’å¢—ã‚„ã™ å¾ã€…ã«ç­‹è‚‰ãŒã¤ã æ€¥æ¿€ã«ç­‹è‚‰ãŒè½ã¡ã‚‹ å‹•è©å¥ + ã€Œç­‹è‚‰ã€ + è£œèªï¼ä¿®é£¾èª ç­‹è‚‰ã®ç™ºé”ã‚’ä¿ƒã™ ç­‹è‚‰ã®ç¶­æŒã‚’å›³ã‚‹ ç­‹è‚‰ã®ä¿®å¾©ã‚’åŠ©ã‘ã‚‹ ç­‹è‚‰ç—›ã‚’æ„Ÿã˜ã‚‹ ç­‹è‚‰ç—›ãŒæ®‹ã‚‹",
    # "çµ„ã‚€": "åè© + ã€Œã‚’ã€ + ã€Œçµ„ã‚€ã€ ãƒãƒ¼ãƒ ã‚’çµ„ã‚€ çµ„ç¹”ã‚’çµ„ã‚€ è¨ˆç”»ã‚’çµ„ã‚€ æ™‚é–“å‰²ã‚’çµ„ã‚€ æ—¥ç¨‹ã‚’çµ„ã‚€ ã€Œã€œã¨ï¼ã€œã§ï¼ã€œã«ã€ + ã€Œçµ„ã‚€ã€ å½¼ã¨çµ„ã‚€ äºŒäººã§çµ„ã‚€ ä»–ç¤¾ã¨çµ„ã‚€ åŒç›Ÿã‚’çµ„ã‚€ é€£æºã‚’çµ„ã‚€ ä½“ã®éƒ¨åˆ† + ã€Œã‚’çµ„ã‚€ã€ è…•ã‚’çµ„ã‚€ è¶³ã‚’çµ„ã‚€ æ‰‹ã‚’çµ„ã‚€ è†ã‚’çµ„ã‚€ æŒ‡ã‚’çµ„ã‚€ å‹•è©å¥ + ã€Œçµ„ã‚€ã€ çµ„ã‚“ã§å”åŠ›ã™ã‚‹ çµ„ã‚“ã§é€²ã‚ã‚‹ çµ„ã‚“ã§å–ã‚Šçµ„ã‚€ çµ„ã‚“ã§å¯¾å¿œã™ã‚‹ çµ„ã‚“ã§è§£æ±ºã™ã‚‹ å½¢å®¹è© / å‰¯è© + ã€Œçµ„ã‚€ã€ å¯†ã«çµ„ã‚€ ç·Šå¯†ã«çµ„ã‚€ ç¶¿å¯†ã«çµ„ã‚€ ç·»å¯†ã«çµ„ã‚€ å¼·å›ºã«çµ„ã‚€",
    # "èŠ¸èƒ½": "åè© + ã€Œã®ã€ + ã€ŒèŠ¸èƒ½ã€ èŠ¸èƒ½ç•Œã® èŠ¸èƒ½æ´»å‹•ã® èŠ¸èƒ½ç•Œã®äº‹æƒ… èŠ¸èƒ½ç•Œã®äººè„ˆ èŠ¸èƒ½ç•Œã®å¸¸è­˜ ã€ŒèŠ¸èƒ½ã€ + åŠ©è© + åè©ï¼èªå¥ èŠ¸èƒ½ç•Œ èŠ¸èƒ½äºº èŠ¸èƒ½æ´»å‹• èŠ¸èƒ½ç•Œå…¥ã‚Š èŠ¸èƒ½ç•Œãƒ‡ãƒ“ãƒ¥ãƒ¼ å‹•è© + ã€ŒèŠ¸èƒ½ã‚’ï¼èŠ¸èƒ½ç•Œã‚’ã€œã€ èŠ¸èƒ½ç•Œã«å…¥ã‚‹ èŠ¸èƒ½æ´»å‹•ã‚’ã™ã‚‹ èŠ¸èƒ½ç•Œã§æ´»èºã™ã‚‹ èŠ¸èƒ½ç•Œã‚’å¼•é€€ã™ã‚‹ èŠ¸èƒ½ç•Œã‚’å¸­å·»ã™ã‚‹ å‰¯è©ï¼å‰¯è©å¥ + ã€ŒèŠ¸èƒ½ç•Œã§ï¼èŠ¸èƒ½æ´»å‹•ã‚’ã€œã€ ç²¾åŠ›çš„ã«èŠ¸èƒ½æ´»å‹•ã‚’ã™ã‚‹ æœ¬æ ¼çš„ã«èŠ¸èƒ½ç•Œå…¥ã‚Šã™ã‚‹ ç©æ¥µçš„ã«èŠ¸èƒ½æ´»å‹•ã‚’è¡Œã† ç€ã€…ã¨èŠ¸èƒ½ç•Œã§ã‚­ãƒ£ãƒªã‚¢ã‚’ç©ã‚€ é–“ã‚‚ãªãèŠ¸èƒ½ç•Œãƒ‡ãƒ“ãƒ¥ãƒ¼ã™ã‚‹",
    # "è¦‹èˆã„": "åè© + ã€Œã®ã€ + ã€Œè¦‹èˆã„ï¼ãŠè¦‹èˆã„ã€ å…¥é™¢ã®è¦‹èˆã„ ç—…æ°—ã®ãŠè¦‹èˆã„ ç«äº‹ã®è¦‹èˆã„ å‡ºç”£ã®ãŠè¦‹èˆã„ è¦‹èˆã„å®¢ å‹•è© + ã€Œè¦‹èˆã„ï¼ãŠè¦‹èˆã„ã‚’ã€œã€ è¦‹èˆã„ã«è¡Œãï¼è¡Œã† ãŠè¦‹èˆã„ã‚’ã™ã‚‹ ãŠè¦‹èˆã„ã‚’æ¸¡ã™ ãŠè¦‹èˆã„ã‚’é€ã‚‹ ãŠè¦‹èˆã„ã‚’å—ã‘ã‚‹ ã€Œã€œã«ã€ + ã€ŒãŠè¦‹èˆã„ã€ ç—…é™¢ã«ãŠè¦‹èˆã„ï¼ˆã«è¡Œãï¼‰ å®¶ã«ãŠè¦‹èˆã„ï¼ˆã«è¡Œãï¼‰ å…¥é™¢å…ˆã«ãŠè¦‹èˆã„ï¼ˆã«è¡Œãï¼‰ ä¼šç¤¾ã«ãŠè¦‹èˆã„ï¼ˆã‚’ã™ã‚‹ï¼‰ å½¢å®¹è© / å‰¯è© + ã€ŒãŠè¦‹èˆã„ã€ï¼ã€Œè¦‹èˆã„ã€ï¼‹ å‹•è© å¿ƒã‹ã‚‰ã®ãŠè¦‹èˆã„ å¿ƒã°ã‹ã‚Šã®ãŠè¦‹èˆã„ ã¡ã‚‡ã£ã¨ã—ãŸãŠè¦‹èˆã„ ãŠè¦‹èˆã„ç”³ã—ä¸Šã’ã¾ã™ æ—©æœŸã®ãŠè¦‹èˆã„ å‹•è©å¥ + ã€Œè¦‹èˆã„ï¼ãŠè¦‹èˆã„ã€ + è£œèªï¼ä¿®é£¾èª è¦‹èˆã„å“ã‚’è´ˆã‚‹ ãŠè¦‹èˆã„çŠ¶ã‚’æ›¸ã ãŠè¦‹èˆã„ã®è¨€è‘‰ã‚’è¿°ã¹ã‚‹ ãŠè¦‹èˆã„é‡‘ã‚’åŒ…ã‚€ è¦‹èˆã„é‡‘ã‚’é€ã‚‹",
    # "è¦šã‚ã‚‹": "åè© + ã€ŒãŒã€ + ã€Œè¦šã‚ã‚‹ï¼è¦šã‚ã¦ã€ å¤¢ãŒè¦šã‚ã‚‹ é…”ã„ãŒè¦šã‚ã‚‹ æ„è­˜ãŒè¦šã‚ã‚‹ å¹»æƒ³ãŒè¦šã‚ã‚‹ æƒ…ç†±ãŒè¦šã‚ã‚‹ å‰¯è© / å‰¯è©å¥ + ã€Œè¦šã‚ã‚‹ï¼è¦šã‚ã¦ã€ ã™ã£ã¨è¦šã‚ã‚‹ ã±ã£ã¨è¦šã‚ã‚‹ ã ã‚“ã ã‚“è¦šã‚ã‚‹ ç›®ãŒè¦šã‚ã‚‹ã‚ˆã†ã«è¦šã‚ã‚‹ ã¯ã£ã¨è¦šã‚ã‚‹ å‹•è©å¥ + ã€Œè¦šã‚ã‚‹ï¼è¦šã‚ã¦ã€ + è£œèªï¼ä¿®é£¾èª è¦šã‚ãŸç›®ã§è¦‹ã‚‹ è¦šã‚ãŸå¿ƒã§å—ã‘æ­¢ã‚ã‚‹ é…”ã„ãŒã™ã£ã‹ã‚Šè¦šã‚ã‚‹ å¤¢ã‹ã‚‰è¦šã‚ã¦ç¾å®Ÿã‚’è¦‹ã‚‹ è¡å‹•ãŒè¦šã‚ã¦å†·é™ã«ãªã‚‹",
    # "è¨ˆç®—": "åè© + ã€Œã‚’ã€ + ã€Œè¨ˆç®—ã™ã‚‹ã€ æ•°å­—ã‚’è¨ˆç®—ã™ã‚‹ é¢ç©ã‚’è¨ˆç®—ã™ã‚‹ é‡‘é¡ã‚’è¨ˆç®—ã™ã‚‹ ç¨é‡‘ã‚’è¨ˆç®—ã™ã‚‹ åˆ©ç›Šã‚’è¨ˆç®—ã™ã‚‹ åè© + ã€Œã®ã€ + ã€Œè¨ˆç®—ã€ é›»å“ã®è¨ˆç®— ãƒ‘ã‚½ã‚³ãƒ³ã®è¨ˆç®— æ•°å­¦ã®è¨ˆç®— æ©Ÿæ¢°ã®è¨ˆç®— è¤‡é›‘ãªè¨ˆç®— å½¢å®¹è© / å‰¯è© + ã€Œè¨ˆç®—ã€ æ­£ç¢ºãªè¨ˆç®— è¤‡é›‘ãªè¨ˆç®— ç°¡å˜ãªè¨ˆç®— å¤§ã¾ã‹ãªè¨ˆç®— ç²¾å¯†ãªè¨ˆç®— ã€Œè¨ˆç®—ã€+ å‹•è©å¥ è¨ˆç®—ãŒåˆã† è¨ˆç®—ãŒåˆã‚ãªã„ è¨ˆç®—ãŒæ—©ã„ è¨ˆç®—ã«å¼·ã„ è¨ˆç®—ã«å¼±ã„ æ…£ç”¨çš„è¡¨ç¾ è¨ˆç®—ã«å…¥ã‚Œã‚‹ è¨ˆç®—ã‹ã‚‰å¤–ã™ è¨ˆç®—ã‚’èª¤ã‚‹ è¨ˆç®—ã‚’é–“é•ãˆã‚‹ è¨ˆç®—é€šã‚Šã«é€²ã‚€",
    # "è¨±å¯": "åè© + ã€Œã‚’ã€ + ã€Œè¨±å¯ã™ã‚‹ã€ å…¥å ´ã‚’è¨±å¯ã™ã‚‹ ä½¿ç”¨ã‚’è¨±å¯ã™ã‚‹ å»ºç¯‰ã‚’è¨±å¯ã™ã‚‹ è¼¸å…¥ã‚’è¨±å¯ã™ã‚‹ é–‹ç™ºã‚’è¨±å¯ã™ã‚‹ åè© + ã€Œã®ã€ + ã€Œè¨±å¯ã€ å–¶æ¥­ã®è¨±å¯ å»ºç¯‰ã®è¨±å¯ ä½¿ç”¨ã®è¨±å¯ å‡ºå›½ã®è¨±å¯ ç«‹ã¡å…¥ã‚Šã®è¨±å¯ å½¢å®¹è© / å‰¯è© + ã€Œè¨±å¯ã€ ç‰¹åˆ¥ã«è¨±å¯ã™ã‚‹ ä¾‹å¤–çš„ã«è¨±å¯ã™ã‚‹ æ­£å¼ã«è¨±å¯ã™ã‚‹ ä¸€æ™‚çš„ã«è¨±å¯ã™ã‚‹ å³æ ¼ã«è¨±å¯ã™ã‚‹ ã€Œè¨±å¯ã€+ å‹•è©å¥ è¨±å¯ãŒä¸‹ã‚Šã‚‹ è¨±å¯ã‚’å¾—ã‚‹ è¨±å¯ã‚’ç”³è«‹ã™ã‚‹ è¨±å¯ã‚’ä¸ãˆã‚‹ è¨±å¯ã‚’å–ã‚Šæ¶ˆã™ æ…£ç”¨çš„è¡¨ç¾ è¨±å¯ãªãç«‹ã¡å…¥ã‚‹ è¨±å¯ãªã—ã§ä½¿ç”¨ã™ã‚‹ è¨±å¯ã®ã‚‚ã¨ã«è¡Œã† è¨±å¯ã‚’å¾…ã¤ è¨±å¯ã‚’å¿…è¦ã¨ã™ã‚‹",
    # "è¨ºå¯Ÿ": "åè© + ã€Œã‚’ã€ + ã€Œè¨ºå¯Ÿã™ã‚‹ã€ æ‚£è€…ã‚’è¨ºå¯Ÿã™ã‚‹ ç—…äººã‚’è¨ºå¯Ÿã™ã‚‹ å­ã©ã‚‚ã‚’è¨ºå¯Ÿã™ã‚‹ å¤–æ¥ã‚’è¨ºå¯Ÿã™ã‚‹ å†…ç§‘ã‚’è¨ºå¯Ÿã™ã‚‹ åè© + ã€Œã®ã€ + ã€Œè¨ºå¯Ÿã€ åŒ»å¸«ã®è¨ºå¯Ÿ ä»Šæ—¥ã®è¨ºå¯Ÿ å¤–æ¥ã®è¨ºå¯Ÿ å®šæœŸã®è¨ºå¯Ÿ æ¬¡å›ã®è¨ºå¯Ÿ ã€Œè¨ºå¯Ÿã€+ å‹•è©å¥ è¨ºå¯Ÿã‚’å—ã‘ã‚‹ è¨ºå¯Ÿã‚’çµ‚ãˆã‚‹ è¨ºå¯Ÿã‚’å¾…ã¤ è¨ºå¯Ÿã‚’å§‹ã‚ã‚‹ è¨ºå¯Ÿã‚’å¸Œæœ›ã™ã‚‹ å½¢å®¹è© / å‰¯è© + ã€Œè¨ºå¯Ÿã€ ä¸å¯§ã«è¨ºå¯Ÿã™ã‚‹ æ…é‡ã«è¨ºå¯Ÿã™ã‚‹ ã—ã£ã‹ã‚Šè¨ºå¯Ÿã™ã‚‹ æ­£ç¢ºã«è¨ºå¯Ÿã™ã‚‹ è©³ã—ãè¨ºå¯Ÿã™ã‚‹ é–¢é€£çš„è¡¨ç¾ è¨ºå¯Ÿåˆ¸ã‚’å‡ºã™ è¨ºå¯Ÿå®¤ã«å…¥ã‚‹ è¨ºå¯Ÿæ™‚é–“ãŒå§‹ã¾ã‚‹ è¨ºå¯Ÿé †ã‚’å¾…ã¤ è¨ºå¯Ÿæ—¥ã‚’æ±ºã‚ã‚‹",
    # "è©•ä¾¡": "åè© + ã€Œã‚’ã€ + ã€Œè©•ä¾¡ã™ã‚‹ã€ æˆç¸¾ã‚’è©•ä¾¡ã™ã‚‹ èƒ½åŠ›ã‚’è©•ä¾¡ã™ã‚‹ æˆæœã‚’è©•ä¾¡ã™ã‚‹ æ¥­ç¸¾ã‚’è©•ä¾¡ã™ã‚‹ å•†å“ã‚’è©•ä¾¡ã™ã‚‹ åè© + ã€Œã®ã€ + ã€Œè©•ä¾¡ã€ ä½œå“ã®è©•ä¾¡ äººã®è©•ä¾¡ å¸‚å ´ã®è©•ä¾¡ ç¤¾ä¼šã®è©•ä¾¡ ç§‘å­¦çš„è©•ä¾¡ å½¢å®¹è© / å‰¯è© + ã€Œè©•ä¾¡ã€ é«˜ã„è©•ä¾¡ æ­£å½“ãªè©•ä¾¡ å®¢è¦³çš„ãªè©•ä¾¡ å³ã—ã„è©•ä¾¡ å…¬å¹³ãªè©•ä¾¡ ã€Œè©•ä¾¡ã€+ å‹•è©å¥ è©•ä¾¡ã‚’å—ã‘ã‚‹ è©•ä¾¡ã‚’ä¸‹ã™ è©•ä¾¡ã‚’å¾—ã‚‹ è©•ä¾¡ã‚’ä¸ãˆã‚‹ è©•ä¾¡ãŒé«˜ã¾ã‚‹ æ…£ç”¨çš„è¡¨ç¾ è©•ä¾¡ã«å€¤ã™ã‚‹ è©•ä¾¡ã®å¯¾è±¡ã¨ãªã‚‹ è©•ä¾¡ã‚’èª¤ã‚‹ è©•ä¾¡ã‚’å·¦å³ã™ã‚‹ è©•ä¾¡ã‚’è¡Œã†",
    # "è©³ã—ã„": "åè© + ã€Œã«è©³ã—ã„ã€ æ­´å²ã«è©³ã—ã„ éŸ³æ¥½ã«è©³ã—ã„ çµŒæ¸ˆã«è©³ã—ã„ æ³•å¾‹ã«è©³ã—ã„ ãƒ‘ã‚½ã‚³ãƒ³ã«è©³ã—ã„ ã€Œè©³ã—ã„ã€+ åè© è©³ã—ã„èª¬æ˜ è©³ã—ã„å†…å®¹ è©³ã—ã„æƒ…å ± è©³ã—ã„åœ°å›³ è©³ã—ã„æ‰‹é † å‰¯è© + ã€Œè©³ã—ã„ã€ ã¨ã¦ã‚‚è©³ã—ã„ ã‹ãªã‚Šè©³ã—ã„ ãŸã„ã¸ã‚“è©³ã—ã„ ã™ã”ãè©³ã—ã„ ã‚ã‚Šã¨è©³ã—ã„ å‹•è©å¥ + ã€Œè©³ã—ã„ã€é–¢é€£ è©³ã—ã„ã“ã¨ã‚’çŸ¥ã‚‹ è©³ã—ã„äººã«èã è©³ã—ã„è³‡æ–™ã‚’èª­ã‚€ è©³ã—ã„è©±ã‚’ã™ã‚‹ è©³ã—ã„çŠ¶æ³ã‚’èª¿ã¹ã‚‹ æ…£ç”¨çš„ãªä½¿ã„æ–¹ è©³ã—ã„ã“ã¨ã¯åˆ†ã‹ã‚‰ãªã„ è©³ã—ã„èª¬æ˜ã‚’æ±‚ã‚ã‚‹ è©³ã—ã„æ‰‹ç¶šãã‚’æ¡ˆå†…ã™ã‚‹ è©³ã—ã„çµŒç·¯ã‚’æ˜ã‚‰ã‹ã«ã™ã‚‹ è©³ã—ã„äº‹æƒ…ã‚’èª¬æ˜ã™ã‚‹",
    # "èª‡ã‚Š": "åè© + ã€Œã‚’ã€ + ã€Œèª‡ã‚Šã«æ€ã†ã€ æ—¥æœ¬ã‚’èª‡ã‚Šã«æ€ã† ä»•äº‹ã‚’èª‡ã‚Šã«æ€ã† å®¶æ—ã‚’èª‡ã‚Šã«æ€ã† æˆç¸¾ã‚’èª‡ã‚Šã«æ€ã† å‡ºèº«ã‚’èª‡ã‚Šã«æ€ã† ã€Œèª‡ã‚Šã€+ åŠ©è© + åè©ï¼å‹•è©å¥ èª‡ã‚Šã‚’æŒã¤ èª‡ã‚Šã‚’æŠ±ã èª‡ã‚Šã‚’å¤±ã† èª‡ã‚Šã‚’å®ˆã‚‹ èª‡ã‚Šã‚’å‚·ã¤ã‘ã‚‹ å½¢å®¹è© / å‰¯è© + ã€Œèª‡ã‚Šã€ å¤§ããªèª‡ã‚Š é«˜ã„èª‡ã‚Š å¼·ã„èª‡ã‚Š æœ€å¤§ã®èª‡ã‚Š ç„¡ä¸Šã®èª‡ã‚Š æ…£ç”¨çš„è¡¨ç¾ èª‡ã‚Šé«˜ã„ç²¾ç¥ èª‡ã‚Šã‚ã‚‹ä¼çµ± èª‡ã‚Šé«˜ãæ°‘æ— èª‡ã‚Šã®è±¡å¾´ èª‡ã‚Šã®æº å‹•è©å¥ + ã€Œèª‡ã‚Šã€ èª‡ã‚Šã«è€ãˆã‚‹ èª‡ã‚Šã‚’ã‹ã‘ã¦æˆ¦ã† èª‡ã‚Šã‚’èƒ¸ã«æŠ±ã èª‡ã‚Šã‚’å–ã‚Šæˆ»ã™ èª‡ã‚Šã‚’ç¤ºã™",
    # "è±Šå¯Œ": "ã€Œè±Šå¯Œãªã€+ åè© è±Šå¯ŒãªçŸ¥è­˜ è±Šå¯ŒãªçµŒé¨“ è±Šå¯Œãªè³‡æº è±Šå¯Œãªç¨®é¡ è±Šå¯Œãªé£Ÿæ åè© + ã€ŒãŒè±Šå¯Œã ã€ æ „é¤ŠãŒè±Šå¯Œã  æ°´ãŒè±Šå¯Œã  è³‡é‡‘ãŒè±Šå¯Œã  äººæãŒè±Šå¯Œã  æƒ…å ±ãŒè±Šå¯Œã  å‰¯è© + ã€Œè±Šå¯Œã«ã€ éå¸¸ã«è±Šå¯Œã« ãã‚ã‚ã¦è±Šå¯Œã« æ¯”è¼ƒçš„è±Šå¯Œã« æœ€ã‚‚è±Šå¯Œã« ã‹ãªã‚Šè±Šå¯Œã« å‹•è©å¥ + ã€Œè±Šå¯Œã€ ãƒãƒªã‚¨ãƒ¼ã‚·ãƒ§ãƒ³ãŒè±Šå¯Œ ãƒ¡ãƒ‹ãƒ¥ãƒ¼ãŒè±Šå¯Œ é¸æŠè‚¢ãŒè±Šå¯Œ å“æƒãˆãŒè±Šå¯Œ çµŒé¨“ãŒè±Šå¯Œ æ…£ç”¨çš„ãªä½¿ã„æ–¹ è±Šå¯Œãªæ‰èƒ½ è±Šå¯Œãªæ„Ÿæ€§ è±Šå¯Œãªè¡¨ç¾åŠ› è±Šå¯Œãªå†…å®¹ è±Šå¯Œãªè‰²å½©",
    # "è³‡æœ¬": "åè© + ã€Œã®ã€ + ã€Œè³‡æœ¬ã€ ä¼æ¥­ã®è³‡æœ¬ å›½å®¶ã®è³‡æœ¬ å¤–å›½ã®è³‡æœ¬ é‡‘èã®è³‡æœ¬ æ ªä¸»ã®è³‡æœ¬ åè© + ã€Œã‚’ã€ + ã€Œè³‡æœ¬ã«ã™ã‚‹ï¼è³‡æœ¬ã¨ã™ã‚‹ã€ çŸ¥è­˜ã‚’è³‡æœ¬ã«ã™ã‚‹ çµŒé¨“ã‚’è³‡æœ¬ã«ã™ã‚‹ æŠ€è¡“ã‚’è³‡æœ¬ã«ã™ã‚‹ äººè„ˆã‚’è³‡æœ¬ã«ã™ã‚‹ åŠ´åƒåŠ›ã‚’è³‡æœ¬ã¨ã™ã‚‹ ã€Œè³‡æœ¬ã€+ åè©ï¼ˆè¤‡åˆèªï¼‰ è³‡æœ¬ä¸»ç¾© è³‡æœ¬æ”¿ç­– è³‡æœ¬ææº è³‡æœ¬é‡‘ è³‡æœ¬å¸‚å ´ å‹•è© + ã€Œè³‡æœ¬ã€ è³‡æœ¬ã‚’æŠ•å…¥ã™ã‚‹ è³‡æœ¬ã‚’é›†ã‚ã‚‹ è³‡æœ¬ã‚’å°å…¥ã™ã‚‹ è³‡æœ¬ã‚’æŒã¤ è³‡æœ¬ã‚’é‹ç”¨ã™ã‚‹ å½¢å®¹è© / å‰¯è© + ã€Œè³‡æœ¬ã€ å¤–å›½è³‡æœ¬ å¤šé¡ã®è³‡æœ¬ å·¨å¤§ãªè³‡æœ¬ è±Šå¯Œãªè³‡æœ¬ ä¸è¶³ã™ã‚‹è³‡æœ¬",
    # "è³¢ã„": "è³¢ã„ã€+ åè© è³¢ã„äºº è³¢ã„å­ã©ã‚‚ è³¢ã„é¸æŠ è³¢ã„æ–¹æ³• è³¢ã„åˆ¤æ–­ åè© + ã€ŒãŒè³¢ã„ã€ çŠ¬ãŒè³¢ã„ å­ã©ã‚‚ãŒè³¢ã„ ç”Ÿå¾’ãŒè³¢ã„ å½¼ãŒè³¢ã„ å½¼å¥³ãŒè³¢ã„ å‰¯è© + ã€Œè³¢ã„ã€ ã¨ã¦ã‚‚è³¢ã„ å®Ÿã«è³¢ã„ ã‹ãªã‚Šè³¢ã„ ã™ã”ãè³¢ã„ æ¯”è¼ƒçš„è³¢ã„ å‹•è©å¥ + ã€Œè³¢ã„ã€é–¢é€£ è³¢ãç”Ÿãã‚‹ è³¢ãä½¿ã† è³¢ãåˆ¤æ–­ã™ã‚‹ è³¢ãé¸ã¶ è³¢ãå¯¾å¿œã™ã‚‹ æ…£ç”¨çš„è¡¨ç¾ è³¢ã„è€ƒãˆ è³¢ã„ç”Ÿãæ–¹ è³¢ã„æ¶ˆè²»è€… è³¢ã„è²·ã„ç‰© è³¢ã„è¨€ã„æ–¹",
    # "è¿ãˆ": "åè© + å‹•è©ï¼ˆè¿ãˆ + å‹•è©ï¼‰ è¿ãˆã«è¡Œã è¿ãˆã«æ¥ã‚‹ è¿ãˆã‚’é ¼ã‚€ è¿ãˆã«å‘ã‹ã† è¿ãˆã«å‡ºã‚‹ å‹•è© + åè©ï¼ˆå‹•è© + è¿ãˆï¼‰ å‡ºè¿ãˆã‚’ã™ã‚‹ å‡ºè¿ãˆã«å‡ºã‚‹ æ­“è¿ã¨è¿ãˆï¼ˆä¾‹ï¼šæ­“è¿ã¨è¿ãˆï¼‰ è¿ãˆå½¹ã‚’å‹™ã‚ã‚‹ è¿ãˆã®è»Š åè© + åè©ï¼ˆè¿ãˆ + åè©ï¼‰ è¿ãˆã®è»Š è¿ãˆã®æ™‚é–“ è¿ãˆã®è¨€è‘‰ è¿ãˆã®æº–å‚™ è¿ãˆã®é€£çµ¡ å½¢å®¹è©/é€£ä½“ä¿®é£¾ + è¿ãˆï¼ˆâ€¦ãªè¿ãˆã€â€¦ã®è¿ãˆï¼‰ æ¸©ã‹ã„è¿ãˆ ç››å¤§ãªè¿ãˆ ç°¡ç´ ãªè¿ãˆ äºˆæœŸã›ã¬è¿ãˆ å¿ƒå¼·ã„è¿ãˆ",
    # "éã”ã™": "åè© + ã‚’ + éã”ã™ï¼ˆåè©ã‚’ç›®çš„èªã«ã¨ã‚‹ãƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ æ™‚é–“ã‚’éã”ã™ æ—¥ã€…ã‚’éã”ã™ ä¸€æ—¥ã‚’éã”ã™ ä¼‘æ—¥ã‚’éã”ã™ å¹³å’Œãªæ—¥ã€…ã‚’éã”ã™ ï½ã‚’ + ã©ã†éã”ã™ï¼ˆç–‘å•æ–‡ãƒ»è£œèªã‚’ä½¿ã£ãŸãƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ ã©ã†éã”ã™ã‹ ï½ã‚’ã©ã†éã”ã™ã‹è€ƒãˆã‚‹ ï½ã‚’æœ‰æ„ç¾©ã«éã”ã™ ï½ã‚’å¿«é©ã«éã”ã™ ï½ã‚’æ¥½ã—ãéã”ã™ å‰¯è©ï¼å½¢å®¹è© + ã« + éã”ã™ï¼ˆæ§˜æ…‹ã‚’è¡¨ã™ï¼‰ ã‚†ã£ãŸã‚Šã¨éã”ã™ ã®ã‚“ã³ã‚Šã¨éã”ã™ é™ã‹ã«éã”ã™ ç©ã‚„ã‹ã«éã”ã™ ç„¡äº‹ã«éã”ã™ ã€œã¨ / ã€œã¨ã¨ã‚‚ã« + éã”ã™ï¼ˆå…±åŒè¡Œå‹•ãƒ»æ™‚é–“ã‚’å…±ã«ã™ã‚‹ï¼‰ å®¶æ—ã¨éã”ã™ å‹äººã¨éã”ã™ æ‹äººã¨éã”ã™ ä¸€ç·’ã«éã”ã™ å¤§åˆ‡ãªäººã¨éã”ã™",
    # "é‡‘é¡": "åè© + ã® + é‡‘é¡ ç·é¡ã®é‡‘é¡ è¦‹ç©é¡ã®é‡‘é¡ æ”¯æ‰•ã†é‡‘é¡ è«‹æ±‚é‡‘é¡ æ§é™¤å¾Œã®é‡‘é¡ é‡‘é¡ + ã‚’ + å‹•è©ï¼ˆå‹•è©ã‚’ä¼´ã†ï¼‰ é‡‘é¡ã‚’ç¢ºèªã™ã‚‹ é‡‘é¡ã‚’æç¤ºã™ã‚‹ é‡‘é¡ã‚’æ”¯æ‰•ã† é‡‘é¡ã‚’è¨˜å…¥ã™ã‚‹ é‡‘é¡ã‚’äº¤æ¸‰ã™ã‚‹ å‰¯è© / å½¢å®¹è© + ã« + é‡‘é¡ï¼ˆæ§˜æ…‹ä¿®é£¾ï¼‰ å¤§ããªé‡‘é¡ å°‘é¡ã®é‡‘é¡ é«˜é¡ãªé‡‘é¡ é©æ­£ãªé‡‘é¡ æƒ³å®šå¤–ã®é‡‘é¡ ï¼ˆï½ã¨ / ï½ã«ã‚ˆã£ã¦ï¼‰ + é‡‘é¡ + ãŒ / ã‚’ + å‹•è©ï¼ˆé–¢ä¿‚ãƒ»æ¯”è¼ƒï¼‰ é‡‘é¡ãŒç•°ãªã‚‹ é‡‘é¡ã«ã‚ˆã‚‹å·® é‡‘é¡ã«å¿œã˜ã¦ é‡‘é¡ãŒè†¨ã‚‰ã‚€ é‡‘é¡ãŒæ±ºã¾ã‚‹",
    # "é‹­ã„": "å½¢å®¹è© + åè©ï¼ˆé‹­ã„ + åè©ï¼‰ é‹­ã„ ç—›ã¿ é‹­ã„ è¦–ç·š é‹­ã„ æ‰¹åˆ¤ é‹­ã„ æ„Ÿè¦š é‹­ã„ åˆƒ åè© + ã‚’ + å½¢å®¹è© ï¼ˆåè©ã‚’ä¿®é£¾ã™ã‚‹å½¢ã«å±•é–‹ï¼‰ åˆ‡ã‚Œå‘³ãŒ é‹­ã„ ãƒŠã‚¤ãƒ• æ‰¹è©•ãŒ é‹­ã„ è¨€è‘‰ æ„Ÿè¦šãŒ é‹­ã„ äºº ç›®ã¤ããŒ é‹­ã„ è¡¨æƒ… æ´å¯ŸåŠ›ãŒ é‹­ã„ åˆ†æ å‹•è© + å½¢å®¹è©åŒ–ï¼å‰¯è©åŒ–ï¼ˆé‹­ã + å‹•è©ãªã©ï¼‰ é‹­ã åˆ‡ã‚Šè¾¼ã‚€ é‹­ã æŒ‡æ‘˜ã™ã‚‹ é‹­ã åå¿œã™ã‚‹ é‹­ã è¿«ã‚‹ é‹­ã æ„Ÿã˜ã‚‹ å‰¯è©ãƒ»é€£ä½“ä¿®é£¾èª + ã« + é‹­ã„ éå¸¸ã« é‹­ã„ çªç„¶ é‹­ã„ ä¸€æ®µã¨ é‹­ã„ ã•ã‚‰ã« é‹­ã„ ç•°æ§˜ã« é‹­ã„",
    # "é™¤ã": "åè© + ã‚’ + é™¤ãï¼ˆå¯¾è±¡ã‚’ç›®çš„èªã«ã¨ã‚‹ãƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ ä¸è‰¯å“ã‚’é™¤ã é›‘è‰ã‚’é™¤ã éšœå®³ã‚’é™¤ã ä¾‹å¤–ã‚’é™¤ã ã”ã¿ã‚’é™¤ã åè© + ã‚’é™¤ã„ã¦ï¼ˆç¯„å›²ã‚’é™å®šã™ã‚‹ï¼é™¤å¤–ã‚’ç¤ºã™è¡¨ç¾ï¼‰ã€æ–‡æ³•è¡¨ç¾ã€‘ Bunpro +2 JSMORI +2 åœŸæ—¥ã‚’é™¤ã„ã¦ ç¨é‡‘ã‚’é™¤ã„ã¦ ã‚¢ãƒ«ãƒã‚¤ãƒˆã‚’é™¤ã„ã¦ å½¼ã‚’é™¤ã„ã¦ å ´æ‰€ã‚’é™¤ã„ã¦ å‹•è© + åè©ï¼ˆâ€¦ã‚’é™¤ã + åè©ã€ä¿®é£¾è¡¨ç¾ã¨ã—ã¦ä½¿ã†ï¼‰ é™¤ã„ãŸéƒ¨åˆ† é™¤ãå¯¾è±¡ é™¤ãã¹ãã‚‚ã® é™¤ã„ã¦ã‚ˆã„è¦ç´  é™¤ã„ãŸã‚ã¨ã®çŠ¶æ…‹ å‰¯è©ãƒ»å½¢å®¹è© + ã« + é™¤ãï¼ˆä¿®é£¾å½¢ï¼‰ ä¾‹å¤–çš„ã«é™¤ã æ˜ç¤ºçš„ã«é™¤ã é©åˆ‡ã«é™¤ã å®Œå…¨ã«é™¤ã ä¸€éƒ¨ã‚’é™¤ã",
    # "éš ã™": "åè© + ã‚’ + éš ã™ï¼ˆç›®çš„èªã‚’ã¨ã‚‹ãƒ‘ã‚¿ãƒ¼ãƒ³ï¼‰ çœŸå®Ÿã‚’éš ã™ ç§˜å¯†ã‚’éš ã™ æ¬ ç‚¹ã‚’éš ã™ æœ¬éŸ³ã‚’éš ã™ è²¡ç”£ã‚’éš ã™ ï½ã‚’ + éš ã—ã¦ï¼ˆæ‰‹æ®µãƒ»æ–¹æ³•ã‚’è¡¨ã™å½¢ï¼‰ å½±ã§éš ã—ã¦ å£°ã‚’éš ã—ã¦ èº«ã‚’éš ã—ã¦ é¡”ã‚’éš ã—ã¦ æœ¬å½“ã®æ„å›³ã‚’éš ã—ã¦ å‹•è© + éš ã™ï¼ˆå‹•è©ã¨çµ„ã‚€è¡¨ç¾ï¼‰ æ„å›³ã‚’éš ã™ æ„Ÿæƒ…ã‚’éš ã™ äº‹å®Ÿã‚’éš ã™ ç—›ã¿ã‚’éš ã™ å¼±ã¿ã‚’éš ã™ å‰¯è©ï¼å½¢å®¹è© + ã« + éš ã™ï¼ˆæ§˜å­ãƒ»ç¨‹åº¦ã‚’è¡¨ã™ä¿®é£¾èªï¼‰ ã“ã£ãã‚Šéš ã™ å®Œå…¨ã«éš ã™ ã†ã¾ãéš ã™ å·§å¦™ã«éš ã™ å¯†ã‹ã«éš ã™",
    # "é›†ã¾ã‚Š": "åè© + å‹•è© ï¼ˆåè© â€œé›†ã¾ã‚Šâ€ ã‚’ç›®çš„èªãƒ»ä¸»èªã«ã¨ã‚‹å‹•è©ã¨ã®çµã³ã¤ãï¼‰ é›†ã¾ã‚Š ãŒã‚ã‚‹ é›†ã¾ã‚Š ãŒæ‚ªã„ é›†ã¾ã‚Š ã«å‚åŠ ã™ã‚‹ é›†ã¾ã‚Š ã‚’é–‹ã é›†ã¾ã‚Š ã‚’è¨­ã‘ã‚‹ å½¢å®¹è© + åè© ï¼ˆâ€œã©ã®ã‚ˆã†ãªé›†ã¾ã‚Šã‹â€ ã‚’è¡¨ã™ï¼‰ å°‘ãªã„é›†ã¾ã‚Š å¤šã„é›†ã¾ã‚Š å°ã•ãªé›†ã¾ã‚Š è¦ªã—ã„é›†ã¾ã‚Š éå…¬å¼ãªé›†ã¾ã‚Š åè© + ã® + åè© ï¼ˆâ€œï½ã®é›†ã¾ã‚Šâ€ ã®å½¢ã§ä½¿ã‚ã‚Œã‚‹ï¼‰ å‹é”ã®é›†ã¾ã‚Š å®¶æ—ã®é›†ã¾ã‚Š åŒçª“ä¼šã®é›†ã¾ã‚Š è¶£å‘³ã®é›†ã¾ã‚Š åœ°åŸŸã®é›†ã¾ã‚Š å‹•è© + ã® + åè© ï¼ˆâ€œï½ã®é›†ã¾ã‚Šâ€ ã‚’å‹•è©åŒ–ã—ãŸè¡¨ç¾ï¼‰ ä¼šã† ã®é›†ã¾ã‚Šï¼ˆå£èªçš„ã«ä½¿ã‚ã‚Œã‚‹ã“ã¨ã‚‚ï¼‰ é›†ã¾ã‚‹ ã®é›†ã¾ã‚Š é–‹ã ã®é›†ã¾ã‚Š æ•´ãˆã‚‹ ã®é›†ã¾ã‚Šï¼ˆä¾‹ï¼šæº–å‚™ã‚’æ•´ãˆã‚‹ã®é›†ã¾ã‚Šã€ã¯å°‘ã—ä¸è‡ªç„¶ã ãŒè¦‹ã‹ã‘ã‚‹ã“ã¨ã‚‚ï¼‰ æ±ºã‚ã‚‹ ã®é›†ã¾ã‚Š å‹•è©å¥ + åè© ï¼ˆä»–ã®å‹•è©å¥ã¨çµ„ã¿åˆã‚ã›ã¦ä½¿ã‚ã‚Œã‚‹è¡¨ç¾ï¼‰ äºˆå®šã•ã‚ŒãŸé›†ã¾ã‚Š çªç™ºçš„ãªé›†ã¾ã‚Š å®šæœŸçš„ãªé›†ã¾ã‚Š æ˜¼é–“ã®é›†ã¾ã‚Š å¤œã®é›†ã¾ã‚Š",
    # "éœ‡ãˆã‚‹": "ä¸»èª + éœ‡ãˆã‚‹ ï¼ˆä½•ãŒéœ‡ãˆã‚‹ã®ã‹ã‚’è¡¨ã™ï¼‰ æ‰‹ãŒéœ‡ãˆã‚‹ è¶³ãŒéœ‡ãˆã‚‹ å£°ãŒéœ‡ãˆã‚‹ å…¨èº«ãŒéœ‡ãˆã‚‹ è‚©ãŒéœ‡ãˆã‚‹ éœ‡ãˆã‚‹ + åŠ©è© + åè© / éœ‡ãˆã‚‹ + åè© ï¼ˆéœ‡ãˆã‚‹å¯¾è±¡ã‚„åŸå› ãªã©ã‚’ç¶šã‘ã‚‹è¡¨ç¾ï¼‰ éœ‡ãˆã‚‹ ã»ã© éœ‡ãˆã‚‹ éŸ³ éœ‡ãˆã‚‹ æ‰‹ã§ éœ‡ãˆã‚‹ å£°ã§ éœ‡ãˆã‚‹ ç¬‘é¡”ï¼ˆæ¯”å–©çš„è¡¨ç¾ï¼‰ éœ‡ãˆã‚‹ + åŠ©è© + å‹•è© ï¼ˆéœ‡ãˆãªãŒã‚‰ã€œã™ã‚‹ã€ãªã©ã®è¡¨ç¾ï¼‰ éœ‡ãˆãªãŒã‚‰è©±ã™ éœ‡ãˆãªãŒã‚‰ç«‹ã¤ éœ‡ãˆãªãŒã‚‰ç¬‘ã† éœ‡ãˆãªãŒã‚‰ç¥ˆã‚‹ éœ‡ãˆãªãŒã‚‰æ­Œã† åŸå› ãƒ»ç†ç”±ã‚’è¡¨ã™èª + ã§ / ã« + éœ‡ãˆã‚‹ ï¼ˆãªãœéœ‡ãˆã‚‹ã‹ã‚’è¡¨ã™æ§‹é€ ï¼‰ å¯’ã•ã§éœ‡ãˆã‚‹ ç·Šå¼µã§éœ‡ãˆã‚‹ ææ€–ã§éœ‡ãˆã‚‹ æ„Ÿå‹•ã§éœ‡ãˆã‚‹ é©šãã«éœ‡ãˆã‚‹ éœ‡ãˆã‚‹ + è£œåŠ©è¡¨ç¾ï¼ˆé€£ç”¨å½¢ã‚„å‰¯è©å¥ï¼‰ ã¶ã‚‹ã¶ã‚‹éœ‡ãˆã‚‹ ããããéœ‡ãˆã‚‹ ãµã‚‹ãµã‚‹éœ‡ãˆã‚‹ å¾®ã‹ã«éœ‡ãˆã‚‹ æ¿€ã—ãéœ‡ãˆã‚‹",
    # "éå¸¸": "å‰¯è© + å½¢å®¹è©ï¼å‹•è© ï¼ˆéå¸¸ã«ï½ã™ã‚‹ï¼éå¸¸ã«ï½ã ï¼‰ éå¸¸ã« é‡è¦ã  éå¸¸ã« é›£ã—ã„ éå¸¸ã« åŠ©ã‹ã‚‹ éå¸¸ã« é©šã éå¸¸ã« å¹¸é‹ã  å½¢å®¹å‹•è© + åè© ï¼ˆã€Œéå¸¸ãªï½ã€ã®å½¢ã§ä½¿ã‚ã‚Œã‚‹ã“ã¨ãŒå¤šã„ï¼‰ éå¸¸ãª äº‹æ…‹ éå¸¸ãª çŠ¶æ…‹ éå¸¸ãª å±é™º éå¸¸ãª äº‹æ…‹ éå¸¸ãª åŠªåŠ› åè© + ã® + éå¸¸ ï¼ˆã€Œï½ã®éå¸¸ï¼ˆï½ã®ç·Šæ€¥æ€§ï¼‰ã€ã®æ„å‘³åˆã„ã§ä½¿ã‚ã‚Œã‚‹ï¼‰ å‘½ã®éå¸¸ ç¤¾ä¼šã®éå¸¸ å›½ã®éå¸¸ çŠ¶æ³ã®éå¸¸ åœ°åŸŸã®éå¸¸ ï¼ˆãŸã ã—ã“ã®ãƒ‘ã‚¿ãƒ¼ãƒ³ã¯ã‚„ã‚„æ–‡èªçš„ãƒ»é™å®šçš„ãªç”¨æ³•ï¼‰ éå¸¸ + ã€å‹•è©å¥ã€‘ ï¼ˆã€Œéå¸¸ï½ã™ã‚‹ï¼éå¸¸ï½ã•ã›ã‚‹ã€ãªã©ã®æ§‹é€ ï¼‰ éå¸¸ã‚’ è¦ã™ã‚‹ éå¸¸ã‚’ å‘¼ã³ã‹ã‘ã‚‹ éå¸¸ã‚’ å®£è¨€ã™ã‚‹ éå¸¸ã‚’ èªã‚ã‚‹ éå¸¸ã‚’ ç¶­æŒã™ã‚‹",
    # "é ‚ã": "ä¸»èª + é ‚ã ï¼ˆç›®ä¸Šã®äººã‹ã‚‰ä½•ã‹ã‚’ã€Œé ‚ãã€ï¼‰ è³ã‚’é ‚ã æ‰¿èªã‚’é ‚ã ã”æ”¯æ´ã‚’é ‚ã ã”é€£çµ¡ã‚’é ‚ã ã”æ„è¦‹ã‚’é ‚ã ï½ã—ã¦ + é ‚ã ï¼ˆã€Œï½ã—ã¦ã‚‚ã‚‰ã†ã€ã®è¬™è­²è¡¨ç¾ï¼šã€Œï½ã—ã¦é ‚ãï¼é ‚ããŸã„ã€ãªã©ï¼‰ ã”æ•™ç¤ºã—ã¦é ‚ã ã”å”åŠ›ã—ã¦é ‚ã ã”é€£çµ¡ã—ã¦é ‚ã ã”èª¬æ˜ã—ã¦é ‚ã ã”æ¤œè¨ã—ã¦é ‚ã ï½ã•ã›ã¦ + é ‚ã ï¼ˆè‡ªåˆ†ã®è¡Œç‚ºã«ã¤ã„ã¦ã€ç›¸æ‰‹ã®è¨±å¯ã‚’å¾—ã¦è¡Œã†è¡¨ç¾ï¼‰ ã”æ¡ˆå†…ã•ã›ã¦é ‚ã ã”èª¬æ˜ã•ã›ã¦é ‚ã ã”ç´¹ä»‹ã•ã›ã¦é ‚ã ã”ç›¸è«‡ã•ã›ã¦é ‚ã ã”æ±ºå®šã•ã›ã¦é ‚ã ãŠï¼ã” + åè© + é ‚ã ï¼ˆåè©ã«ã€ŒãŠï¼ã”ã€ã‚’ä»˜ã‘ã¦æ•¬æ„ã‚’è¾¼ã‚ã€ã€Œé ‚ãã€ã‚’ä½¿ã†ï¼‰ ã”åˆ©ç”¨é ‚ã ã”ç†è§£é ‚ã ã”å¯¾å¿œé ‚ã ã”é€£çµ¡é ‚ã ã”è©•ä¾¡é ‚ã ï½ã¦ + é ‚ãï¼ˆè£œåŠ©å‹•è©çš„ç”¨æ³•ï¼‰ ï¼ˆã€Œï½ã¦ã‚‚ã‚‰ã†ã€ã®ä¸å¯§ãªè¨€ã„æ–¹ï¼‰ ã”æ‰¿çŸ¥ã„ãŸã ã ã”è¦§ã„ãŸã ã ãŠèãã„ãŸã ã ãŠè©±ã—ã„ãŸã ã ã”ç¢ºèªã„ãŸã ã",
    # "é ­ç—›": "ä¸»èª + å‹•è©ï¼ˆã€œãŒã™ã‚‹ï¼ã€œãŒèµ·ã“ã‚‹ ãªã©ï¼‰ é ­ç—›ãŒã™ã‚‹ é ­ç—›ãŒèµ·ã“ã‚‹ é ­ç—›ãŒç¶šã é ­ç—›ãŒã²ã©ã„ é ­ç—›ãŒæ²»ã‚‹ åè© + ã® + åè© åé ­ç—›ã® ç™ºä½œ ç·Šå¼µå‹é ­ç—›ã® ç—‡çŠ¶ ç‰‡é ­ç—›ã® ç™ºä½œ é ­ç—›ã® åŸå›  é ­ç—›ã® ç¨®é¡ å½¢å®¹è©ï¼å‰¯è© + åè© æ¿€ã—ã„é ­ç—› è»½ã„é ­ç—› ã²ã©ã„é ­ç—› é ­ç—›ã®ã²ã©ã• é ­ç—›ã®é »åº¦ åŸå› ãƒ»ç†ç”± + ã§ï¼ã« + é ­ç—› ã‚¹ãƒˆãƒ¬ã‚¹ã§é ­ç—› ç–²ã‚Œã§é ­ç—› å¯ä¸è¶³ã§é ­ç—› é¨’éŸ³ã«é ­ç—› å…‰ã«é ­ç—› è£œåŠ©èªå¥ + é ­ç—› é ­ç—›æŒã¡ï¼ˆï¼é ­ç—›ãŒèµ·ã“ã‚Šã‚„ã™ã„äººï¼‰ é ­ç—›è–¬ é ­ç—›å¤–æ¥ é ­ç—›æ—¥è¨˜ é ­ç—›äºˆé˜²",
    # "é£½ãã¾ã§": "é£½ãã¾ã§ + åè© é£½ãã¾ã§ è‡ªèª¬ é£½ãã¾ã§ æ–¹é‡ é£½ãã¾ã§ åå¯¾ é£½ãã¾ã§ å¯èƒ½æ€§ é£½ãã¾ã§ æ„è¦‹ é£½ãã¾ã§ + å‹•è©ï¼å¥ï¼ˆï½ã™ã‚‹ï¼ï½ã—ã‚ˆã†ã¨ã™ã‚‹ etc.ï¼‰ é£½ãã¾ã§ ä¸»å¼µã™ã‚‹ é£½ãã¾ã§ è²«ã é£½ãã¾ã§ åå¯¾ã™ã‚‹ é£½ãã¾ã§ æ±‚ã‚ã‚‹ é£½ãã¾ã§ å®ˆã‚ã†ã¨ã™ã‚‹ å‰¯è©ï¼‹é£½ãã¾ã§ ï¼ˆã€Œé£½ãã¾ã§ã‚‚ã€ã®å½¢ã§å¼·èª¿ã‚’åŠ ãˆã‚‹ï¼‰ ã‚ãã¾ã§ã‚‚ æ„è¦‹ã¨ã—ã¦ ã‚ãã¾ã§ã‚‚ ç§è¦‹ã ãŒ ã‚ãã¾ã§ã‚‚ ãã®ç¯„å›²ã§ ã‚ãã¾ã§ã‚‚ å½¢å¼çš„ã« ã‚ãã¾ã§ã‚‚ åç›®ä¸Š",
    # "é«ªã®æ¯›": "ä¸»èª + å‹•è©ï¼ˆï½ãŒï½ã™ã‚‹ï¼ï½ã‚’ï½ã™ã‚‹ ç­‰ï¼‰ é«ªã®æ¯› ãŒæŠœã‘ã‚‹ é«ªã®æ¯› ãŒä¼¸ã³ã‚‹ é«ªã®æ¯› ã‚’åˆ‡ã‚‹ é«ªã®æ¯› ã‚’æŸ“ã‚ã‚‹ é«ªã®æ¯› ã‚’æ´—ã† å½¢å®¹è©ï¼å‰¯è© + é«ªã®æ¯› é•·ã„é«ªã®æ¯› çŸ­ã„é«ªã®æ¯› æŸ”ã‚‰ã‹ã„é«ªã®æ¯› ãµã•ãµã•ã®é«ªã®æ¯› å‚·ã‚“ã é«ªã®æ¯› é«ªã®æ¯› + ã® + åè© é«ªã®æ¯› ã®é‡ é«ªã®æ¯› ã®çŠ¶æ…‹ é«ªã®æ¯› ã®è³ª é«ªã®æ¯› ã®è‰² é«ªã®æ¯› ã®æ‰‹å…¥ã‚Œ åŸå› ãƒ»ç†ç”±ï¼æ¡ä»¶ + ã§ï¼ã« + é«ªã®æ¯› ã‚¹ãƒˆãƒ¬ã‚¹ ã§ é«ªã®æ¯›ãŒæŠœã‘ã‚‹ åŠ é½¢ ã« ã‚ˆã‚‹é«ªã®æ¯›ã®å¤‰åŒ– æ „é¤Šä¸è¶³ ã§ é«ªã®æ¯›ãŒç´°ããªã‚‹ ç´«å¤–ç·š ã« ã‚ˆã£ã¦é«ªã®æ¯›ãŒå‚·ã‚€ ã‚·ãƒ£ãƒ³ãƒ—ãƒ¼ ã§ é«ªã®æ¯›ã‚’æ´—ã† è£œåŠ©èªï¼‹é«ªã®æ¯› é«ªã®æ¯› å…ˆï¼ˆä¾‹ï¼šæ¯›å…ˆã€å…ˆç«¯ï¼‰ é«ªã®æ¯› æŸï¼ˆä¾‹ï¼šä¸€æŸï¼äºŒæŸã®é«ªã®æ¯›ï¼‰ é«ªã®æ¯› é‡ï¼ˆå¤šã„ï¼å°‘ãªã„ï¼‰ é«ªã®æ¯› è³ªï¼ˆç¡¬ã•ãƒ»æŸ”ã‚‰ã‹ã•ãªã©ï¼‰ é«ªã®æ¯› è‰²ï¼ˆé»’é«ªï¼ç™½é«ªï¼èŒ¶é«ªãªã©ï¼‰",
    # "é»’æ¿": "ä¸»èª + å‹•è©ï¼ˆï½ã™ã‚‹ï¼ï½ãŒï½ã™ã‚‹ ãªã©ï¼‰ é»’æ¿ ã«æ›¸ã é»’æ¿ ã‚’æ¶ˆã™ é»’æ¿ ã‚’æ‹­ã é»’æ¿ ã‚’ä½¿ã† é»’æ¿ ãŒè¦‹ãˆã‚‹ å½¢å®¹è©ï¼å‰¯è© + é»’æ¿ å¤ã„é»’æ¿ å¤§ããªé»’æ¿ å£ã®é»’æ¿ æ²ç¤ºç”¨é»’æ¿ ç§»å‹•å¼é»’æ¿ é»’æ¿ + ã® + åè© é»’æ¿ æ–‡å­— é»’æ¿ æ¶ˆã—ã‚´ãƒ ï¼ˆé»’æ¿æ¶ˆã—ï¼‰ é»’æ¿ é¢ é»’æ¿ æ¿æ›¸ é»’æ¿ è¡¨é¢ ç”¨é€”ãƒ»æ©Ÿèƒ½ã‚’è¡¨ã™èª + é»’æ¿ æˆæ¥­ç”¨é»’æ¿ ä¼šè­°ç”¨é»’æ¿ é›»å­é»’æ¿ï¼ˆãƒ‡ã‚¸ã‚¿ãƒ«é»’æ¿ï¼‰ ä¼æ¥­ç”¨ãƒ‡ã‚¸ã‚¿ãƒ«ãƒ›ãƒ¯ã‚¤ãƒˆãƒœãƒ¼ãƒ‰ ã€ãƒŸãƒ©ã‚¤ã‚¿ãƒƒãƒBizã€ ç§»å‹•é»’æ¿ ä¸¡é¢é»’æ¿ åŸå› ãƒ»æ‰‹æ®µ + ã§ï¼ã« + é»’æ¿ ãƒãƒ§ãƒ¼ã‚¯ã§ é»’æ¿ ã«æ›¸ã Wikipedia æŒ‡ã§ é»’æ¿ ã‚’è§¦ã‚‹ æ¶ˆã—ã‚´ãƒ ã§ é»’æ¿ ã‚’æ¶ˆã™ é›¨ã§ é»’æ¿ ãŒæ±šã‚Œã‚‹ é¢¨ã§ é»’æ¿ ã«ç²‰ãŒèˆã†",
    # "é›†ã¾ã‚‹": "å‹•è© + å‰¯è©ï¼é€£ä½“èªï¼ˆä¿®é£¾èªï¼‰ ã²ãã‹ã«é›†ã‚‹ æ¬¡ã€…ã¨é›†ã‚‹ è‡ªç„¶ã«é›†ã‚‹ ã¤ã„é›†ã‚‹ ç„¡æ„è­˜ã«é›†ã‚‹ ä¸»èªï¼ˆåè©ï¼‰ + ãŒ + å‹•è©ï¼šã€Œï½ãŒé›†ã‚‹ã€ è™«ãŒé›†ã‚‹ è¦‹ç‰©äººãŒé›†ã‚‹ ç¾¤è¡†ãŒé›†ã‚‹ ã²ã¨ãŒé›†ã‚‹ ä¸è‰¯ã‚°ãƒ«ãƒ¼ãƒ—ãŒé›†ã‚‹ äºº + ã« + å‹•è©ï¼šã€Œï½ã«é›†ã‚‹ã€ ï¼ˆã€Œé‡‘å“ã‚’ã›ã³ã‚‹ãƒ»ãŠã”ã‚‰ã›ã‚‹ã€æ„å‘³ã§ä½¿ã‚ã‚Œã‚‹ç”¨æ³•ï¼‰ å…ˆè¼©ã«é›†ã‚‹ å‹äººã«é›†ã‚‹ è¦ªã«é›†ã‚‹ å®¢ã«é›†ã‚‹ ä»²é–“ã«é›†ã‚‹ å‹•è© + ç›®çš„èªï¼ˆï½ã‚’ + å‹•è©ï¼‰ ï¼ˆãŸã ã—ã€Œé›†ã‚‹ã€ã¯è‡ªå‹•è©çš„ã«ä½¿ã‚ã‚Œã‚‹ã“ã¨ãŒå¤šã„ãŸã‚ã€ç›®çš„èªã‚’å–ã‚‹ä¾‹ã¯é™å®šçš„ï¼‰ é£Ÿäº‹ã‚’é›†ã‚‹ é…’ä»£ã‚’é›†ã‚‹ é…’ã‚’ä¸€æ¯é›†ã‚‹ ç†Ÿèªçš„ï¼æ…£ç”¨è¡¨ç¾ãƒ»åè©å¥ã¨ã®çµã³ã¤ã é»’å±±ï¼ˆã®ã‚ˆã†ã«ï¼‰äººãŒé›†ã‚‹ äººã ã‹ã‚Šï¼ˆãŒã§ãã¦ï¼‰é›†ã‚‹ é»’å±±ã®äººé›†ã‚Š è¦‹ç‰©äººã®é›†ã‚Š å›£ä½“ã«é›†ã‚‹",
    # "ç”Ÿãã‚‹": "å‰¯è©ï¼é€£ä½“èª + ç”Ÿãã‚‹ï¼ˆä¿®é£¾èª + å‹•è©ï¼‰ ç²¾ä¸€æ¯ç”Ÿãã‚‹ ãŸãã¾ã—ãç”Ÿãã‚‹ ã—ã¶ã¨ãç”Ÿãã‚‹ æ­»ã¬ã¾ã§ç”Ÿãã‚‹ æ„ç¾©ã‚’ã‚‚ã£ã¦ç”Ÿãã‚‹ ä¸»èª + ãŒ + ç”Ÿãã‚‹ï¼ˆï½ãŒç”Ÿãã‚‹ï¼‰ äººãŒç”Ÿãã‚‹ å‹•ç‰©ãŒç”Ÿãã‚‹ æ¤ç‰©ãŒç”Ÿãã‚‹ ç¤¾ä¼šãŒç”Ÿãã‚‹ ç´°èƒãŒç”Ÿãã‚‹ ï½ã« + ç”Ÿãã‚‹ï¼ˆï½ã«ç”Ÿãã‚‹ï¼‰ï¼ï½ã§ç”Ÿãã‚‹ å¿ƒã«ç”Ÿãã‚‹ å¤¢ã«ç”Ÿãã‚‹ å¸Œæœ›ã«ç”Ÿãã‚‹ ä¿¡å¿µã§ç”Ÿãã‚‹ çœŸå®Ÿã§ç”Ÿãã‚‹ ï½ã‚’ + ç”Ÿãã‚‹ï¼ˆç›®çš„èªçš„ãªèªã‚’ä¼´ã†è¡¨ç¾ï¼‰ æ˜æ—¥ã‚’ç”Ÿãã‚‹ ä»Šã‚’ç”Ÿãã‚‹ éå»ã‚’ç”Ÿãã‚‹ ç¾å®Ÿã‚’ç”Ÿãã‚‹ è‡ªåˆ†ã‚’ç”Ÿãã‚‹ ç†Ÿèªçš„ï¼æ…£ç”¨è¡¨ç¾ã¨ã®çµã³ã¤ã ç”Ÿãã‚‹ä¾¡å€¤ ç”Ÿãã‚‹æ„å‘³ ç”Ÿãã‚‹åŠ› ç”Ÿãã‚‹è¡“ ç”Ÿãã‚‹è¨¼",
    # "æ¤ãˆã‚‹": "ä¸»èª + ãŒ + æ¤ãˆã‚‹ï¼ˆï½ãŒæ¤ãˆã‚‹ï¼‰ è¾²å®¶ãŒæ¤ãˆã‚‹ å¸‚æ°‘ãŒæ¤ãˆã‚‹ å­¦æ ¡ãŒæ¤ãˆã‚‹ å›£ä½“ãŒæ¤ãˆã‚‹ äººãŒæ¤ãˆã‚‹ ï½ã‚’ + æ¤ãˆã‚‹ï¼ˆç›®çš„èª + å‹•è©ï¼‰ æœ¨ã‚’æ¤ãˆã‚‹ èŠ±ã‚’æ¤ãˆã‚‹ é‡èœã‚’æ¤ãˆã‚‹ ç¨®ã‚’æ¤ãˆã‚‹ è‰ã‚’æ¤ãˆã‚‹ å ´æ‰€ + ã« + æ¤ãˆã‚‹ï¼ˆå ´æ‰€ã‚’ç›®çš„èªã«ã™ã‚‹è¨€ã„æ–¹ï¼‰ åº­ã«æ¤ãˆã‚‹ ç•‘ã«æ¤ãˆã‚‹ å…¬åœ’ã«æ¤ãˆã‚‹ æ ¡åº­ã«æ¤ãˆã‚‹ é‰¢ï¼ˆã«/ã§ï¼‰æ¤ãˆã‚‹ å‰¯è©ï¼ä¿®é£¾èª + æ¤ãˆã‚‹ï¼ˆä¿®é£¾èª + å‹•è©ï¼‰ ã¾ãšæ¤ãˆã‚‹ ã¡ã‚ƒã‚“ã¨æ¤ãˆã‚‹ æ·±ãæ¤ãˆã‚‹ å¤šãæ¤ãˆã‚‹ å°‘ã—ãšã¤æ¤ãˆã‚‹ ç†Ÿèªçš„ï¼æ…£ç”¨è¡¨ç¾ã¨ã®çµã³ã¤ã æ¤ãˆã‚‹ä¾¡å€¤ æ¤ãˆã‚‹æ„ç¾© å¿ƒã«æ¤ãˆã‚‹ï¼ˆæ¯”å–©ç”¨æ³•ï¼‰ æ ¹ã‚’æ¤ãˆã‚‹ï¼ˆæ¯”å–©çš„è¡¨ç¾ï¼‰ æ€æƒ³ã‚’æ¤ãˆã‚‹ï¼ˆæ¯”å–©çš„ã«ï¼‰",
    # "ãŠåœŸç”£": "å½¢å®¹è© + ãŠåœŸç”£ï¼ãŠåœŸç”£ + ã®ï¼ˆä¿®é£¾èª + åè©ï¼‰ åç‰©ã®ãŠåœŸç”£ ä»£è¡¨çš„ãªãŠåœŸç”£ æ‰‹ä½œã‚Šã®ãŠåœŸç”£ åœ°å…ƒã®ãŠåœŸç”£ ã¡ã‚‡ã£ã¨ã—ãŸãŠåœŸç”£ ãŠåœŸç”£ä»£ ãŠåœŸç”£å±‹ ãŠåœŸç”£ç‰© ãŠåœŸç”£åº— ãŠåœŸç”£é¸ã³ å‹•è© + ãŠåœŸç”£ï¼ˆå‹•è© + åè©ï¼‰ è²·ã†ãŠåœŸç”£ æ¸¡ã™ãŠåœŸç”£ é€ã‚‹ãŠåœŸç”£ é¸ã¶ãŠåœŸç”£ ç”¨æ„ã™ã‚‹ãŠåœŸç”£ ãŠåœŸç”£ã‚’è²·ã† ãŠåœŸç”£ã‚’æ¸¡ã™ ãŠåœŸç”£ã‚’æŒã£ã¦å¸°ã‚‹ ãŠåœŸç”£ã‚’è´ˆã‚‹ ãŠåœŸç”£ã‚’åŒ…ã‚€ ï½ã«ï¼ã‹ã‚‰ï¼ã¸ + ãŠåœŸç”£ï¼ˆæ ¼åŠ©è© + åè©è¡¨ç¾ï¼‰ å®¶æ—ã«ãŠåœŸç”£ å‹é”ã«ãŠåœŸç”£ è·å ´ã«ãŠåœŸç”£ åœ°å…ƒã‹ã‚‰ã®ãŠåœŸç”£ æµ·å¤–ã‹ã‚‰ã®ãŠåœŸç”£ æ•°é‡è©ï¼‹ãŠåœŸç”£ï¼é‡ã‚’ç¤ºã™è¡¨ç¾ ä¸€ã¤ã®ãŠåœŸç”£ äºŒã¤ã®ãŠåœŸç”£ ãŸãã•ã‚“ã®ãŠåœŸç”£ å°‘ã—ã®ãŠåœŸç”£ å¤§ããªãŠåœŸç”£",
    # "æ‚²ã—ã„": "å‰¯è©ï¼ä¿®é£¾èª + æ‚²ã—ã„ï¼ˆä¿®é£¾èª + å½¢å®¹è©ï¼‰ ã¨ã¦ã‚‚æ‚²ã—ã„ æœ¬å½“ã«æ‚²ã—ã„ ã²ã©ãæ‚²ã—ã„ ãªã‚“ã¨ãªãæ‚²ã—ã„ å¿ƒã‹ã‚‰æ‚²ã—ã„ ãšã£ã¨æ‚²ã—ã„ é™ã‹ã«æ‚²ã—ã„ æ‚²ã—ã„ + åè©ï¼ˆå½¢å®¹è© + åè©ï¼‰ æ‚²ã—ã„æ°—æŒã¡ æ‚²ã—ã„æ€ã„ æ‚²ã—ã„éå» æ‚²ã—ã„è¨˜æ†¶ æ‚²ã—ã„çµæœ« æ‚²ã—ã„çŸ¥ã‚‰ã› æ‚²ã—ã„ç¾å®Ÿ æ‚²ã—ã„åˆ¥ã‚Œ ï½ãŒ + æ‚²ã—ã„ï¼ˆä¸»èª + è¿°èªï¼‰ å½¼å¥³ãŒæ‚²ã—ã„ å›ãŒæ‚²ã—ã„ å­ã©ã‚‚ãŒæ‚²ã—ã„ ã“ã®äº‹å®ŸãŒæ‚²ã—ã„ çµæœãŒæ‚²ã—ã„ çœŸå®ŸãŒæ‚²ã—ã„ ï½ã‚’ + æ‚²ã—ã„ï¼ˆç›®çš„èªï¼å¯¾è±¡ + å½¢å®¹è©è¡¨ç¾ï¼‰ ã“ã®å½¢ã¯ã‚„ã‚„é™å®šçš„ã§ã™ãŒã€æ¯”å–©çš„ãƒ»æ„Ÿè¦šçš„è¡¨ç¾ã§ä½¿ã‚ã‚Œã‚‹ã“ã¨ãŒã‚ã‚Šã¾ã™ï¼š æ¶™ã‚’æ‚²ã—ã„ã«å¤‰ãˆã‚‹ï¼ˆæ¯”å–©ï¼‰ éŸ³ã‚’æ‚²ã—ã„ã«ã™ã‚‹ï¼ˆéŸ³æ¥½ãªã©ã§ï¼‰ ç†Ÿèªçš„è¡¨ç¾ãƒ»æ…£ç”¨è¡¨ç¾ã¨ã®çµã³ã¤ã æ‚²ã—ã„ç¾å®Ÿ æ‚²ã—ã„è¨˜æ†¶ æ‚²ã—ã„åˆ¥ã‚Œ æ‚²ã—ã„çŸ¥ã‚‰ã› æ‚²ã—ã„çµæœ« æ‚²ã—ã¿ã«æš®ã‚Œã‚‹ æ‚²ã—ã¿ã‚’èƒŒè² ã† æ‚²ã—ã¿ã‚’ç™’ã™ æ‚²ã—ã•ã‚’æŠ‘ãˆã‚‹ æ‚²ã—ã¿ãŒã“ã¿ä¸Šã’ã‚‹",
    # "æ±ºã‚ã‚‹": "å‰¯è©ï¼ä¿®é£¾èª + æ±ºã‚ã‚‹ï¼ˆä¿®é£¾èª + å‹•è©ï¼‰ ã¯ã£ãã‚Šæ±ºã‚ã‚‹ ãã£ã±ã‚Šæ±ºã‚ã‚‹ æ—©ã‚ã«æ±ºã‚ã‚‹ ã‚ã‚‰ã‹ã˜ã‚æ±ºã‚ã‚‹ è‡ªåˆ†ã§æ±ºã‚ã‚‹ æœ€çµ‚çš„ã«æ±ºã‚ã‚‹ ï½ã‚’ + æ±ºã‚ã‚‹ï¼ˆç›®çš„èª + å‹•è©ï¼‰ é€²è·¯ã‚’æ±ºã‚ã‚‹ æ—¥ç¨‹ã‚’æ±ºã‚ã‚‹ ç›®æ¨™ã‚’æ±ºã‚ã‚‹ è¦å‰‡ã‚’æ±ºã‚ã‚‹ å½¹å‰²ã‚’æ±ºã‚ã‚‹ å€¤æ®µã‚’æ±ºã‚ã‚‹ æ–¹æ³•ã‚’æ±ºã‚ã‚‹ æ–¹å‘ã‚’æ±ºã‚ã‚‹ ï½ã« + æ±ºã‚ã‚‹ï¼ˆé¸æŠå¯¾è±¡ + å‹•è©ï¼‰ ã€œã‚’ A ã«æ±ºã‚ã‚‹ ä¾‹ï¼šéƒ¨å±‹ã‚’ A ã«æ±ºã‚ã‚‹ è¡Œãå…ˆã«æ±ºã‚ã‚‹ åº—ã«æ±ºã‚ã‚‹ ç›¸æ‰‹ã«æ±ºã‚ã‚‹ æ¡ˆã«æ±ºã‚ã‚‹ ä¾‹ï¼š ã“ã®åº—ã«æ±ºã‚ã‚‹ ãƒ—ãƒ©ãƒ³Aã«æ±ºã‚ã‚‹ æ—¥æ™‚ã«æ±ºã‚ã‚‹ ä¸»èª + ãŒ + æ±ºã‚ã‚‹ï¼ˆï½ãŒæ±ºã‚ã‚‹ï¼‰ ãƒªãƒ¼ãƒ€ãƒ¼ãŒæ±ºã‚ã‚‹ ãƒãƒ¼ãƒ ãŒæ±ºã‚ã‚‹ ç¤¾é•·ãŒæ±ºã‚ã‚‹ è¦ªãŒæ±ºã‚ã‚‹ ä¼šè­°ãŒæ±ºã‚ã‚‹ é–¢ä¿‚è€…ãŒæ±ºã‚ã‚‹ ç†Ÿèªçš„ï¼æ…£ç”¨è¡¨ç¾ã¨ã®çµã³ã¤ã æ±ºã‚æ‰‹ã«ãªã‚‹ æ±ºã‚ã”ã¨ æ±ºã‚å°è© æ±ºã‚ãƒãƒ¼ã‚º æ±ºã‚äº‹ æ±ºã‚è¾¼ã‚€ï¼ˆæ´¾ç”Ÿï¼‰ æ±ºã‚ãŸã‚‚ã‚“ã  æ±ºã‚ã‚ãã­ã‚‹",
    # "äº¤é€š": "å½¢å®¹è© + äº¤é€š ï¼ äº¤é€š + ã®ï¼ˆä¿®é£¾èª + åè©ï¼‰ å…¬å…±äº¤é€š è‡ªå‹•è»Šäº¤é€š äº¤é€šå®‰å…¨ äº¤é€šæ©Ÿé–¢ äº¤é€šé‡ äº¤é€šæ¸‹æ» äº¤é€šè¦åˆ¶ äº¤é€šæ··é›‘ äº¤é€šäº‹æƒ… äº¤é€šç¶² å‹•è© + äº¤é€š äº¤é€šã‚’åˆ©ç”¨ã™ã‚‹ äº¤é€šãŒæ··ã‚€ äº¤é€šãŒç™ºé”ã™ã‚‹ äº¤é€šãŒé®æ–­ã•ã‚Œã‚‹ äº¤é€šã‚’æ•´å‚™ã™ã‚‹ äº¤é€šã‚’åˆ¶é™ã™ã‚‹ äº¤é€šã‚’ç¢ºä¿ã™ã‚‹ äº¤é€šã‚’é®ã‚‹ ï½ã‚’ + äº¤é€šï¼ˆç›®çš„èªä»˜ãè¡¨ç¾ï¼‰ äº¤é€šäº‹æ•…ã‚’èµ·ã“ã™ äº¤é€šé•åã‚’ã™ã‚‹ äº¤é€šæ‰‹æ®µã‚’é¸ã¶ äº¤é€šè²»ã‚’æ”¯æ‰•ã† äº¤é€šä¿¡å·ã‚’å®ˆã‚‹ ï½ã« + äº¤é€šï¼ˆæ ¼ï¼‹åè©å¥ï¼‰ éƒ½å¸‚ã«äº¤é€šãŒç™ºé”ã—ã¦ã„ã‚‹ åœ°æ–¹ã«äº¤é€šã®ä¾¿ãŒæ‚ªã„ å¸‚å†…ã«äº¤é€šç¶²ãŒæ•´ã£ã¦ã„ã‚‹ é«˜é€Ÿé“è·¯ã«äº¤é€šè¦åˆ¶ãŒã‚ã‚‹ é§…ã«äº¤é€šã‚¢ã‚¯ã‚»ã‚¹ãŒã„ã„ ç†Ÿèªãƒ»æ…£ç”¨è¡¨ç¾ äº¤é€šäº‹æ•… äº¤é€šãƒ«ãƒ¼ãƒ« äº¤é€šå®‰å…¨ äº¤é€šæ©Ÿé–¢ äº¤é€šç¶² äº¤é€šæ‰‹æ®µ äº¤é€šæ”¿ç­– äº¤é€šã‚¤ãƒ³ãƒ•ãƒ© äº¤é€šæ¸‹æ» äº¤é€šäº‹æƒ…",
    words_dict = {}

    for word, raw_collocations in words_dict.items():
        w = JPWord(word, 3, raw_collocations=raw_collocations, ai_init=True)

        # w = JPWord.load_from_json(word)
        w.tts()
        w.pptx_generation()
