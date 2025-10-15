import ast
import json
import os
from typing import Annotated, Any, List

from pydantic import BaseModel, ConfigDict, Field
from pydub import AudioSegment

from src import LANGUAGES_ABBR, get_openai_client, get_translator_client
from src.logger_module import get_logger
from src.word.JPWord import extract_kanji, query_jisho, query_kanji

logger = get_logger("JVD")


class KanjiDetail(BaseModel):
    model_config = ConfigDict(extra="forbid")
    kanji: str = Field(examples=["学", "校"], max_length=1, additionalProperties=False)  # type: ignore
    onyomi: list[str] = Field(
        min_items=0, max_items=2, examples=[["がく"], ["こう", "きょう"]], additionalProperties=False
    )  # type: ignore
    kunyomi: list[str] = Field(
        min_items=0, max_items=2, examples=[["まな.ぶ"], ["つか.う", "つか.える"]], additionalProperties=False
    )  # type: ignore
    meanings_english: list[str] = Field(
        min_items=1, max_items=3, examples=[["study", "learning"], ["school", "exam"]], additionalProperties=False
    )  # type: ignore
    common_words: list[str] = Field(
        min_items=1,
        max_items=2,
        examples=[
            ["土曜日 (どようび): Saturday", "土地 (とち): land, plot"],
            ["出産 (しゅっさん): childbirth", "産業 (さんぎょう): industry"],
        ],
        additionalProperties=False,
    )  # type: ignore


NuanceList = Annotated[
    List[Annotated[str, Field(max_length=20, description="single english word")]], Field(min_items=1, max_items=5)  # type: ignore
]


class JapaneseText(BaseModel):
    model_config = ConfigDict(extra="forbid")
    kanji: str = Field(examples=["本を読む"], max_length=100, additionalProperties=False)  # type: ignore
    furigana: str = Field(examples=["本(ほん)を読(よ)む"], max_length=100, additionalProperties=False)  # type: ignore
    translations: dict[str, str] = Field(default_factory=dict)  # COMMENT OUT THIS LINE BEFORE RUNNING BATCH


class JPWordInfo(BaseModel):
    model_config = ConfigDict(extra="forbid")
    kanji: str = Field(examples=["学校", "猫", "走る"], max_length=4, additionalProperties=False)  # type: ignore
    reading: str = Field(examples=["がっこう", "ねこ", "はしる"], max_length=8, additionalProperties=False)  # type: ignore
    introduction_japanese: str = Field(max_length=250, additionalProperties=False)  # type: ignore
    introduction_english: str = Field(max_length=250, additionalProperties=False)  # type: ignore
    meanings: List[NuanceList] = Field(
        min_items=1,
        max_items=3,
        examples=[[["degree", "level", "amount"], ["balance", "moderation"]]],
        additionalProperties=False,
    )  # type: ignore
    meaning_explanation_japanese: str = Field(max_length=250, additionalProperties=False)  # type: ignore
    meaning_explanation_english: str = Field(max_length=250, additionalProperties=False)  # type: ignore
    youtube_description: str = Field(max_length=250, additionalProperties=False)  # type: ignore
    kanji_details: list[KanjiDetail] = Field(min_items=1, max_items=4)  # type: ignore
    kanji_explanation: str = Field(max_length=500, additionalProperties=False)  # type: ignore
    synonyms: list[Annotated[str, Field(max_length=50)]] = Field(
        min_items=0, max_items=2, examples=[["土産 : みやげ : souvenir", "贈り物 : おくりもの : gift, present"]]
    )  # type: ignore
    synonym_explanation: str = Field(max_length=250, additionalProperties=False)  # type: ignore
    antonyms: list[Annotated[str, Field(max_length=50)]] = Field(
        min_items=0, max_items=2, examples=[["暑い : あつい : hot", "高い : たかい : tall, high"]]
    )  # type: ignore
    antonym_explanation: str = Field(max_length=250, additionalProperties=False)  # type: ignore
    collocations: list[JapaneseText] = Field(min_items=5, max_items=8)  # type: ignore
    example_sentences: list[JapaneseText] = Field(min_items=4, max_items=5)  # type: ignore
    meanings_translations: list[dict[str, Any]] = Field(
        default_factory=list
    )  # COMMENT OUT THIS LINE BEFORE RUNNING BATCH

    def pptx_generation(self, word: str, jlpt_level: int, num_examples: int | None = 4) -> None:
        """Generate PowerPoint presentation for the word"""
        import os

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import (
            MSO_AUTO_SIZE,
            MSO_VERTICAL_ANCHOR,
            PP_PARAGRAPH_ALIGNMENT,
        )
        from pptx.util import Inches, Pt

        logger.info(f"🟦 Generating PPTX for word: {word}")
        file_name = f"./Output/{word}/{word} JLPT N{jlpt_level} Vocabulary.pptx"
        if os.path.exists(file_name):
            return

        prs = Presentation("resources/pptx_templates/template.pptx")

        # Title slide
        first_slide = prs.slides.add_slide(prs.slide_layouts[0])
        presentation_title = first_slide.shapes.title
        if presentation_title:
            presentation_title.text = self.kanji
            presentation_title.text_frame.paragraphs[0].font.size = Pt(160)
            presentation_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 95, 154)

        presentation_subtitle = first_slide.placeholders[1]
        presentation_subtitle.text = self.reading  # type: ignore
        presentation_subtitle.text_frame.paragraphs[0].font.size = Pt(100)  # type: ignore
        presentation_subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(192, 79, 21)  # type: ignore

        first_slide.shapes.add_movie(
            f"./output/{word}/audio/0_introduction.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Definitions slide
        definitions_slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = definitions_slide.shapes.add_textbox(Inches(5), Inches(0), Inches(40 / 3 - 5), Inches(7.5))
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        shape.text_frame.word_wrap = True
        for translations in self.meanings_translations:
            # Extract the EN list which contains the primary nuance words
            en_meanings = translations.get("EN", [])
            if not en_meanings:
                continue

            # First word is the key (primary nuance)
            en = en_meanings[0] if isinstance(en_meanings, list) else en_meanings

            p = shape.text_frame.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            run.text = str(en) + "\n"
            run.font.size = Pt(115)
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(33, 95, 154)

            # Show additional English nuances if available
            if isinstance(en_meanings, list) and len(en_meanings) > 1:
                run = p.add_run()
                run.text = ", ".join(en_meanings[1:]) + "\n"
                run.font.size = Pt(32)
                run.font.color.rgb = RGBColor(192, 79, 21)
                run.font.name = "Berlin Sans FB"
                p.space_after = Pt(0)
                p.line_spacing = 0.9

            for code, translation in translations.items():
                if code.lower() == "en":
                    continue
                run_code = p.add_run()
                run_code.text = f"     {code}"
                run_code.font.size = Pt(16)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)

                # Add ordinary run for the translation
                run_translation = p.add_run()
                run_translation.text = f"{translation}"
                run_translation.font.size = Pt(20)
                run_translation.font.name = "Berlin Sans FB"
                run_translation.font.color.rgb = RGBColor(0, 0, 0)

        definitions_slide.shapes.add_movie(
            f"./output/{word}/audio/1_definition.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Add a slide for Kanji breakdown
        kanji_slide = prs.slides.add_slide(prs.slide_layouts[6])
        for i, k in enumerate(self.kanji_details):
            kanji_detail: KanjiDetail = k
            kanji_shape = kanji_slide.shapes.add_textbox(
                Inches(0), Inches(i * 7.5 / 2), Inches(40 / 3 * 0.3), Inches(7.5 / 2)
            )
            p = kanji_shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = kanji_detail.kanji
            run.font.size = Pt(250)
            run.font.bold = True
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(33, 95, 154)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.space_after = Pt(0)
            p.line_spacing = 0.9

            kanji_explanation_shape = kanji_slide.shapes.add_textbox(
                Inches(40 / 3 * 0.3),
                Inches(i * 7.5 / 2),
                Inches(40 / 3 * 0.7),
                Inches(7.5 / 2),
            )
            kanji_explanation_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            kanji_explanation_shape.text_frame.word_wrap = True
            p = kanji_explanation_shape.text_frame.paragraphs[-1]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            run = p.add_run()
            run.text = "Meanings:"
            run.font.size = Pt(24)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(127, 127, 127)
            run = p.add_run()
            run.text = f" {', '.join(kanji_detail.meanings_english)}\n"
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = "Readings:"
            run.font.size = Pt(24)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(127, 127, 127)
            run = p.add_run()
            run.text = f" {', '.join(kanji_detail.kunyomi + kanji_detail.onyomi)}\n"
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            for vocab in kanji_detail.common_words:
                run.text += f"\n{vocab}"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(192, 79, 21)
            p.space_after = Pt(0)
            p.line_spacing = 0.9

        kanji_slide.shapes.add_movie(
            f"./output/{word}/audio/2_kanji_explanation.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Add a slide for examples
        for i, example in enumerate(self.example_sentences[:num_examples]):
            example_slide = prs.slides.add_slide(prs.slide_layouts[6])

            top_shape = example_slide.shapes.add_textbox(Inches(0), Inches(0), Inches(40 / 3), Inches(3))
            top_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            top_shape.text_frame.word_wrap = True
            top_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = top_shape.text_frame.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            run.text = example.kanji
            run.font.size = Pt(70)
            run.font.bold = True
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(33, 95, 154)
            run = p.add_run()
            run.text = f"\n{example.furigana}"  # type: ignore
            run.font.size = Pt(33)
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(192, 79, 21)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.space_after = Pt(0)
            p.line_spacing = 0.9

            left_languages = [
                "EN",
                "ID",
                "ES",
                "VI",
                "FR",
                "NE",
                "AR",
            ]

            left_shape = example_slide.shapes.add_textbox(
                Inches(20 / 3 * 0.025), Inches(3), Inches(20 / 3 * 0.95), Inches(4.5)
            )
            left_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            left_shape.text_frame.word_wrap = True
            left_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for lang in left_languages:
                p = left_shape.text_frame.paragraphs[-1]
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                run_code = p.add_run()
                run_code.text = f"{lang}: "
                run_code.font.size = Pt(20)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)
                run_translation = p.add_run()
                run_translation.text = f"{example.translations[lang]}\n"
                run_translation.font.size = Pt(25)
                run_translation.font.name = "Berlin Sans FB"
            right_shape = example_slide.shapes.add_textbox(
                Inches(20 / 3 * 1.025), Inches(3), Inches(20 / 3 * 0.95), Inches(4.5)
            )
            right_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            right_shape.text_frame.word_wrap = True
            right_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            right_languages = [
                "BN",
                "ZH",
                "KO",
                "TL",
                "MY",
                "HI",
                "FA",
            ]
            for lang in right_languages:
                p = right_shape.text_frame.paragraphs[-1]
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                run_code = p.add_run()
                run_code.text = f"{lang}: "
                run_code.font.size = Pt(20)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)
                run_translation = p.add_run()
                run_translation.text = f"{example.translations[lang]}\n"  # type: ignore
                run_translation.font.size = Pt(25)
                run_translation.font.bold = False
                run_translation.font.name = "Berlin Sans FB"

            example_slide.shapes.add_movie(
                f"./output/{word}/audio/{3 + i}_example.wav",
                left=Pt(0),
                top=Pt(-50),
                width=Pt(50),
                height=Pt(50),
                mime_type="audio/x-wav",
            )

        # Add a slide for synonyms and antonyms
        synonyms_slide = prs.slides.add_slide(prs.slide_layouts[6])

        synonyms_shape = synonyms_slide.shapes.add_textbox(Inches(0), Inches(0), Inches(20 / 3), Inches(7.5))
        synonyms_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        synonyms_shape.text_frame.word_wrap = True
        p = synonyms_shape.text_frame.paragraphs[-1]
        run = p.add_run()
        run.text = "\nSYNONYMS\n\n"
        run.font.size = Pt(40)
        run.font.name = "Berlin Sans FB Demi"
        run.font.color.rgb = RGBColor(33, 95, 154)
        for synonym in self.synonyms:
            kanji, reading, meaning = synonym.split(":")
            run = p.add_run()
            run.text = f"{kanji}\n"
            run.font.size = Pt(54)
            run.font.bold = True  # type: ignore
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(33, 95, 154)
            run = p.add_run()
            run.text = f"{reading}\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = f"{meaning}\n\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(0)
            p.line_spacing = 0.9
        synonyms_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        antonyms_shape = synonyms_slide.shapes.add_textbox(Inches(20 / 3), Inches(0), Inches(20 / 3), Inches(7.5))
        antonyms_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        antonyms_shape.text_frame.word_wrap = True
        p = antonyms_shape.text_frame.paragraphs[-1]
        run = p.add_run()
        run.text = "\nANTONYMS\n\n"
        run.font.size = Pt(40)
        run.font.name = "Berlin Sans FB Demi"
        run.font.color.rgb = RGBColor(192, 79, 21)
        for antonym in self.antonyms:
            kanji, reading, meaning = antonym.split(":")
            run = p.add_run()
            run.text = f"{kanji}\n"
            run.font.size = Pt(54)
            run.font.bold = True  # type: ignore
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(192, 79, 21)
            run = p.add_run()
            run.text = f"{reading}\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = f"{meaning} \n\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(0)
            p.line_spacing = 0.9
        antonyms_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        synonyms_slide.shapes.add_movie(
            f"./output/{word}/audio/100_synonyms_antonyms.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        prs.save(file_name)

    def add_translations_to_examples(self):
        """
        Add translations to all example sentences by translating the kanji text
        to all supported languages. Modifies the example_sentences in place.
        """
        logger.info(f"Adding translations to example sentences for word: {self.kanji}")
        for example in self.example_sentences:
            # Get the kanji text to translate
            text_to_translate = example.kanji
            logger.debug(f"Translating example: {text_to_translate}")

            # Translate to all languages
            translations = translate_to_all_languages(text_to_translate, source_language="ja")

            # Set translations field
            example.translations = translations

        logger.info(f"Successfully added translations to {len(self.example_sentences)} examples")
        return self

    def add_translations_to_collocations(self):
        """
        Add translations to all collocations by translating the kanji text
        to all supported languages. Modifies the collocations in place.
        """
        logger.info(f"Adding translations to collocations for word: {self.kanji}")
        for collocation in self.collocations:
            # Get the kanji text to translate
            text_to_translate = collocation.kanji
            logger.debug(f"Translating collocation: {text_to_translate}")

            # Translate to all languages
            translations = translate_to_all_languages(text_to_translate, source_language="ja")

            # Set translations field
            collocation.translations = translations

        logger.info(f"Successfully added translations to {len(self.collocations)} collocations")
        return self

    def add_translations_to_meanings(self):
        """
        Add translations to all meanings by translating each nuance word
        to all supported languages. Modifies the meanings in place.
        """
        logger.info(f"Adding translations to meanings for word: {self.kanji}")
        for nuance_list in self.meanings:
            translations = translate_to_all_languages(nuance_list[0], source_language="en")
            translations["EN"] = nuance_list
            self.meanings_translations.append(translations)

    def add_all_translations(self):
        """
        Add translations to both example sentences and collocations.
        Convenience method to translate all text content.
        """
        logger.info(f"Adding translations to all text content for word: {self.kanji}")
        self.add_translations_to_examples()
        self.add_translations_to_collocations()
        self.add_translations_to_meanings()
        logger.info("Successfully added all translations")
        return self

    def tts(self, word: str, num_examples: int | None = 4) -> None:
        """
        Generate text-to-speech audio files for the word, including:
        - Introduction (English and Japanese)
        - Meaning explanations (English and Japanese)
        - Kanji explanation
        - Example sentences
        - Synonyms and antonyms explanation

        Args:
            word: The Japanese word (kanji)
            num_examples: Number of example sentences to generate audio for (default: 4)
        """
        os.makedirs(f"./Output/{word}/audio", exist_ok=True)

        # Generate audio for introduction
        introduction_en_audio_path = f"./output/{word}/audio/introduction_en.mp3"
        if not os.path.exists(introduction_en_audio_path):
            logger.info("🔊 Generating audio for english introduction")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.introduction_english,
                instructions=f"English mixed with Japanese. calmly and gently. Correct pronunciation: {word} = {self.reading}",
                response_format="mp3",
                speed=0.95,
            )
            with open(introduction_en_audio_path, "wb") as audio_file:
                audio_file.write(response.content)

        introduction_jp_audio_path = f"./output/{word}/audio/introduction_jp.mp3"
        if not os.path.exists(introduction_jp_audio_path):
            logger.info("🔊 Generating audio for japanese introduction")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.introduction_japanese,
                instructions=f"Japanese. calmly and gently. Correct pronunciation: {word} = {self.reading}",
                response_format="mp3",
                speed=0.95,
            )
            with open(introduction_jp_audio_path, "wb") as audio_file:
                audio_file.write(response.content)

        introduction_audio = AudioSegment.silent(duration=100)
        introduction_audio += AudioSegment.from_mp3(introduction_en_audio_path).apply_gain(12)
        introduction_audio += AudioSegment.silent(duration=500)
        introduction_audio += AudioSegment.from_mp3(introduction_jp_audio_path).apply_gain(12)
        with open(f"./output/{word}/audio/0_introduction.wav", "wb") as title_file:
            introduction_audio.export(title_file, format="wav")

        # Generate audio for meaning explanation
        definition_en_audio_path = f"./output/{word}/audio/definition_en.mp3"
        if not os.path.exists(definition_en_audio_path):
            logger.info("🔊 Generating audio for english definition")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.meaning_explanation_english,
                instructions=f"English mixed with Japanese. Calmly and gently. Correct pronunciation: {word}:{self.reading}",
                response_format="mp3",
                speed=0.95,
            )
            with open(definition_en_audio_path, "wb") as audio_file:
                audio_file.write(response.content)

        definition_jp_audio_path = f"./output/{word}/audio/definition_jp.mp3"
        if not os.path.exists(definition_jp_audio_path):
            logger.info("🔊 Generating audio for japanese definition")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.meaning_explanation_japanese,
                instructions=f"Japanese. calmly and gently. Correct pronunciation: {word} = {self.reading}",
                response_format="mp3",
                speed=0.95,
            )
            with open(definition_jp_audio_path, "wb") as audio_file:
                audio_file.write(response.content)

        definition_audio = AudioSegment.silent(duration=100)
        definition_audio += AudioSegment.from_mp3(definition_en_audio_path).apply_gain(12)
        definition_audio += AudioSegment.silent(duration=500)
        definition_audio += AudioSegment.from_mp3(definition_jp_audio_path).apply_gain(12)
        with open(f"./output/{word}/audio/1_definition.wav", "wb") as word_file:
            definition_audio.export(word_file, format="wav")

        # Generate audio for kanji explanations
        kanji_explanation_audio_path = f"./output/{word}/audio/kanji_explanation.mp3"
        if not os.path.exists(kanji_explanation_audio_path):
            logger.info("🔊 Generating audio for kanji explanation")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.kanji_explanation,
                instructions=f"English. calmly and gently. Correct pronunciation: {word} = {self.reading}",
                response_format="mp3",
                speed=0.95,
            )
            with open(kanji_explanation_audio_path, "wb") as audio_file:
                audio_file.write(response.content)

        kanji_explanation_audio = AudioSegment.from_mp3(kanji_explanation_audio_path).apply_gain(12)
        with open(f"./output/{word}/audio/2_kanji_explanation.wav", "wb") as kanji_file:
            kanji_explanation_audio.export(kanji_file, format="wav")

        # Generate audio for examples
        for i, example in enumerate(self.example_sentences[:num_examples]):
            example_jp_audio_path = f"./output/{word}/audio/example_{i + 1}_jp.mp3"
            if not os.path.exists(example_jp_audio_path):
                logger.info(f"🔊 Generating audio for example {i + 1}")
                response = get_openai_client().audio.speech.create(
                    model="gpt-4o-mini-tts",
                    voice="coral",
                    input=example.kanji,
                    instructions=f"Japanese. calmly and gently. Correct pronunciation: {word} = {self.reading}",
                    response_format="mp3",
                    speed=0.95,
                )
                with open(example_jp_audio_path, "wb") as audio_file:
                    audio_file.write(response.content)

            example_en_audio_path = f"./output/{word}/audio/example_{i + 1}_en.mp3"
            if not os.path.exists(example_en_audio_path):
                logger.info(f"🔊 Generating audio for example {i + 1} translation")
                response = get_openai_client().audio.speech.create(
                    model="gpt-4o-mini-tts",
                    voice="coral",
                    input=example.translations.get("EN", ""),
                    instructions=f"English. calmly and gently. Correct pronunciation: {word} = {self.reading}",
                    response_format="mp3",
                    speed=0.95,
                )
                with open(example_en_audio_path, "wb") as audio_file:
                    audio_file.write(response.content)

            example_jp_audio = AudioSegment.from_mp3(example_jp_audio_path).apply_gain(12)
            example_en_audio = AudioSegment.from_mp3(example_en_audio_path).apply_gain(12)

            example_audio = AudioSegment.silent(duration=100)
            example_audio += example_en_audio
            example_audio += AudioSegment.silent(duration=500)
            example_audio += example_jp_audio
            example_audio += AudioSegment.silent(duration=500)
            example_audio += example_jp_audio
            with open(f"./output/{word}/audio/{3 + i}_example.wav", "wb") as example_file:
                example_audio.export(example_file, format="wav")

        # Generate audio for synonyms and antonyms explanations
        synonyms_explanation_audio_path = f"./output/{word}/audio/synonyms_explanation.mp3"
        if not os.path.exists(synonyms_explanation_audio_path):
            logger.info("🔊 Generating audio for synonyms explanation")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.synonym_explanation + " " + self.antonym_explanation,
                instructions=f"English mixed with Japanese. calmly and gently. Correct pronunciation: {word} = {self.reading}",
                response_format="mp3",
                speed=0.95,
            )
            with open(synonyms_explanation_audio_path, "wb") as audio_file:
                audio_file.write(response.content)

        synonyms_explanation_audio = AudioSegment.from_mp3(synonyms_explanation_audio_path).apply_gain(12)
        with open(f"./output/{word}/audio/100_synonyms_antonyms.wav", "wb") as syn_ant_file:
            synonyms_explanation_audio.export(syn_ant_file, format="wav")


prompt_template = """You are a friendly teacher who explains Japanese vocabulary to beginners. Use a clear, concise, spoken style (as if to a friend). Keep every section brief but complete.

Target word: {{word}}

Output the sections below using **exactly** these headings and this order—no extra commentary.

## introduction_japanese
In Japanese only. Without giving the meaning or reading, name typical situations/contexts where this word is used. Start with the word itself. 1-2 short spoken sentences suitable for elementary learners.

## introduction_english
English translation of **introduction_japanese**. Write the word in kana. Start with: “The [adjective/noun/verb …] [word] …”

## youtube_description
A short English YouTube description for a video explaining the word meaning and use.

## meanings
List **all** meanings grouped by nuance. Each nuance is a list of single-word English glosses. Return a nested list, e.g.:
[[degree,level,amount],[balance,moderation]]

## meaning_explanation_japanese
A short, complete spoken explanation (Japanese) of the literal meanings based on the previous meanings section. Do **not** use the target word itself—use synonyms or antonyms.

## meaning_explanation_english
A short spoken explanation (English) of the literal meanings  based on the previous meanings section. Include the word in kana.

## kanji_details
For **each kanji** in the word: give 1-2 common words (excluding the target word). For each, provide: kanji word, reading, and meaning.

## kanji_explanation_english
For **each kanji** (in order), write one paragraph of 3-4 short sentences in a teacher spoken voice. Start with “The [first/second/…] kanji means …”. Mention 1-2 example vocab items (not the target word) **written in hiragana only**. No bullet points, parentheses, line breaks, titles, or kanji inside the example vocab.

## synonyms
List 1 (max 2) common synonyms **excluding the target word**. Format exactly:
kanji : reading : meaning

## synonyms_explanation
A very short English explanation of the synonyms nuances and how they overlap with the target word. Start with: “The most common synonym[s] of the word [are/is] …”. Write any Japanese vocab **in hiragana only** (no kanji).

## antonyms
List 1 (max 2) common antonyms **excluding the target word**. Format exactly:
kanji : reading : meaning

## antonyms_explanation
A very short English explanation of the antonyms nuances and how they differ from the target word. Start with: “The most common antonym[s] of the word [are/is] …”. Write any Japanese vocab **in hiragana only** (no kanji).

## collocations
Collocation refers to a group of two or more words that usually go together.
List simple, common collocations based on each of the following patterns with the word ({{word}}). 
1) Noun collocation (Det/Num + Adj + N; N + Adj; N + N; Poss + N; N + case/PP)
2) Verb collocation (S + V + O; V + Adv; V + Obj + PP; Aux + V; serial V if normal)
3) Adjective collocation (Adv + Adj; Adj + PP; basic comparatives/superlatives)
4) Adverbial collocation (Adv + Adv; Adv + PP; common time/place adverbials)
For example {kanji:鋭い痛み, furigana:鋭(するど)い痛(いた)み} of Adj+N, {kanji:営業を開始する, furigana:営業(えいぎょう)を開始(かいし)する} of N+V, {kanji:週末の営業, furigana:週末(しゅうまつ)の営業(えいぎょう)} of N+N. Use other common patterns above as needed.

## Examples
Provide 5-7 short, simple and commonly used sentences using the target word ({{word}}) in different random contexts. Each sentence should be suitable for elementary learners. Use a variety of sentence patterns and structures.

For each collocation and example, give:
- Kanji sentence
- Furigana sentence, placing the reading in parentheses **immediately after each kanji** (if no kanji, write the sentence once).
Keep everything beginner-friendly.
"""


def translate_text(text: str, target_language: str, source_language: str | None = "ja") -> str:
    """Translate text from source language to target language using Google Translate"""
    logger.debug(f"🟢 Translating text: {text} from {source_language} to {target_language}")
    client = get_translator_client()

    if source_language is None:
        result = client.detect_language(text)
        source_language = result["language"]

    try:
        # normalize language codes to lowercase for the translator API
        tgt = target_language.lower() if isinstance(target_language, str) else target_language
        src = source_language.lower() if isinstance(source_language, str) else source_language
        result = client.translate(text, target_language=tgt, source_language=src, format_="text")
        logger.debug(f"🟢 Translation result: {result}")
        return result["translatedText"]
    except Exception as e:
        logger.error(f"🟢 Translation failed: {e}")
        return f"Error: {str(e)}"


def translate_to_all_languages(text: str, source_language: str | None = "ja") -> dict:
    """Translate text to all supported languages"""
    logger.debug(f"🟩 Translating text to all languages: {text}")
    translations = {}
    for lang_code in LANGUAGES_ABBR.values():
        # skip translating to the same language as source (guard if source_language is None)
        if source_language and lang_code.lower() == source_language.lower():
            continue
        translated_text = translate_text(text, lang_code, source_language)
        translations[lang_code] = translated_text
    order_list = ["EN", "ID", "ES", "VI", "FR", "NE", "BN", "ZH", "KO", "TL", "MY", "HI", "AR", "FA"]
    sorted_x = {key: translations[key] for key in order_list if key in translations}
    return sorted_x


def get_schema():
    # Make sure to comment out the translations in JPWordInfo before running this function
    # also comment out meanings_translations in JPWordInfo and translations in JapaneseText

    openai_schema = {
        "name": "jp_word_info",
        "type": "json_schema",
        "strict": True,
        "schema": {**JPWordInfo.model_json_schema(), "additionalProperties": False},
    }

    openai_schema_string = str(openai_schema).replace("'", '"').replace("True", "true").replace("False", "false")
    return openai_schema


def generate_word_requests(words: list[str]):
    # fmt: off
    # fmt: on
    openai_schema = get_schema()
    with open("batch_words.jsonl", "w", encoding="utf-8") as f:
        for word in words:
            print(f"Processing word: {word}")
            request = {
                "custom_id": word,
                "method": "POST",
                "url": "/v1/responses",
                "body": {
                    "model": "gpt-5",
                    "input": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "input_text",
                                    "text": prompt_template.replace("{{word}}", word),
                                }
                            ],
                        }
                    ],
                    "text": {"format": openai_schema, "verbosity": "high"},
                    "reasoning": {"effort": "medium", "summary": None},
                    "tools": [],
                    "store": False,
                    "include": ["reasoning.encrypted_content", "web_search_call.action.sources"],
                },
            }
            f.write(
                request.__str__()
                .replace("'", '"')
                .replace("True", "true")
                .replace("False", "false")
                .replace("None", "null")
                + "\n"
            )


def read_batch_results(filepath: str, jlpt_level: int):
    outputs = []
    with open(filepath, "r", encoding="utf-8") as f:
        for line in f:
            result = json.loads(line)
            word = result["custom_id"]
            dir_path = os.path.join("output", word)
            if not os.path.exists(dir_path):
                os.makedirs(dir_path)
            file_name = f"./Output/{word}/{word}.json"
            with open(file_name, "w", encoding="utf-8") as wf:
                data_text = result["response"]["body"]["output"][1]["content"][0]["text"]
                try:
                    data = json.loads(data_text)
                except json.JSONDecodeError:
                    try:
                        data = ast.literal_eval(data_text)
                    except Exception as e:
                        logger.error(f"Failed to parse AI response for {word}: {e}")
                        # skip this item if parsing fails
                        continue
                jp_w = JPWordInfo.model_validate(data)
                jp_w.add_all_translations()
                outputs.append((word, jp_w))
                d = {
                    "version": "0.3.0",
                    "word": word,
                    "jlpt_level": jlpt_level,
                    "youtube_link": "",
                    "in_db": False,
                    **jp_w.model_dump(),
                    "jisho_data": query_jisho(word),
                    "kanji_list": extract_kanji(word),
                    "kanji_data": {k: query_kanji(k) for k in extract_kanji(word)},
                }
                json.dump(d, wf, ensure_ascii=False, indent=2)
            jp_w.tts(word, num_examples=4)
            jp_w.pptx_generation(word, jlpt_level)
    return outputs


ws = [
    "与える",
    "穴",
    "あんなに",
    "市場",
    "従兄弟",
    "伺う",
    "運転",
    "描く",
    "贈る",
    "香り",
    "財産",
    "幸せ",
    "自身",
    "実行",
    "失望",
    "芝居",
    "示す",
    "症状",
    "審判",
    "速度",
    "袖",
    "他人",
    "注目",
    "著者",
    "掴む",
    "包み",
    "提出",
    "適用",
    "手伝い",
    "日光",
    "人間",
    "年代",
    "別に",
    "迷惑",
    "優勝",
    "輸出",
    "夜明け",
    "要点",
    "止す",
    "余分",
    "宜しい",
    "利口",
    "流行",
    "話題",
    "湾",
]

if __name__ == "__main__":
    # get_schema()
    # generate_word_requests(ws)
    w = read_batch_results(r"C:\Users\eskan\Downloads\batch_68ebaf6956e0819090f86166d6593f31_output.jsonl", 3)
    # print(w)
