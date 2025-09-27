import json
import os
import re

import requests
from pydub import AudioSegment

from src import LANGUAGES_ABBR, get_openai_client, get_translator_client
from src.logger_module import get_logger

logger = get_logger("JVD")


def translate_text(text: str, target_language: str, source_language: str | None = "ja") -> str:
    logger.debug(f"🟢 Translating text: {text} from {source_language} to {target_language}")
    client = get_translator_client()

    if source_language is None:
        result = client.detect_language(text)
        source_language = result["language"]

    try:
        result = client.translate(
            text, target_language=target_language, source_language=source_language, format_="text"
        )
        logger.debug(f"🟢 Translation result: {result}")
        return result["translatedText"]
    except Exception as e:
        logger.error(f"🟢 Translation failed: {e}")
        return f"Error: {str(e)}"


def translate_to_all_languages(text: str, source_language: str | None = "ja") -> dict:
    logger.debug(f"🟩 Translating text to all languages: {text}")
    translations = {}
    for lang_code in LANGUAGES_ABBR.values():
        if lang_code.lower() == source_language.lower():
            continue
        translated_text = translate_text(text, lang_code, source_language)
        translations[lang_code] = translated_text
    order_list = ["EN", "ID", "ES", "VI", "FR", "NE", "BN", "ZH", "KO", "TL", "MY", "HI", "AR", "FA"]
    sorted_x = {key: translations[key] for key in order_list if key in translations}
    return sorted_x


def query_jisho(word: str) -> dict | None:
    logger.debug(f"🔴 Querying Jisho API for word: {word}")
    url = f"https://jisho.org/api/v1/search/words?keyword={word}"
    response = requests.get(url)
    if response.status_code == 200:
        res = response.json()
        data = res.get("data", [])
        if data:
            logger.debug(f"🔴 Jisho API response: {data[0]}")
            return data[0]  # Return the first matching entry
        else:
            logger.debug("🔴 No data found in Jisho API response")
            return None
    else:
        logger.error(f"🔴Jish o API request failed: {response.status_code}")
        return None


def query_kanji(kanji: str) -> dict | None:
    logger.debug(f"⚪️ Querying Kanji API for kanji: {kanji}")
    url = f"https://kanjiapi.dev/v1/kanji/{kanji}"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        if data:
            logger.debug(f"⚪️ Kanji API response: {data}")
            return data
        else:
            logger.debug("⚪️ No data found in Kanji API response")
            return None
    else:
        logger.error(f"⚪️ Kanji API request failed: {response.status_code}")
        return None


def extract_kanji(text) -> list[str]:
    return re.findall(r"[\u4E00-\u9FFF]", text)


class JPWord2:
    def __init__(self, initial_word: str, initial_jlpt_level: int, ai_init: bool = False, raw_collocations: str = ""):
        logger.debug(f"📖 Initializing for word: {initial_word}")
        self.word = initial_word
        self.jlpt_level = initial_jlpt_level
        self.youtube_link = ""
        self.in_db = False
        self.jisho_data = {}
        self.kanji_list = []
        self.kanji_data = {}
        self.collocations = []
        self.ai_explanation = {}
        self.meaning_translations = {}
        if not ai_init:
            return
        self.jisho_data = query_jisho(self.word)
        if self.jisho_data:
            self.kanji_list = extract_kanji(self.word)
            for kanji in self.kanji_list:
                kanji_info = query_kanji(kanji)
                if kanji_info:
                    self.kanji_data[kanji] = kanji_info
        else:
            self.kanji_list = []
            self.kanji_data = {}
        self.collocations = self.get_collocations(raw_collocations)
        self.ai_explanation = self.get_ai_explanation()
        self.translate()
        self.save_to_json()

    def get_collocations(self, raw_collocations: str = "") -> dict:
        logger.info(f"🖥️ Getting collocations for word: {self.word}")
        if raw_collocations == "":
            response = get_openai_client().responses.create(
                model="gpt-4o-mini",
                input=[
                    {
                        "role": "system",
                        "content": [
                            {
                                "type": "input_text",
                                "text": "\n".join(
                                    [
                                        "Search on the web for finding the most common and useful collocations of the given Japanese word.",
                                        "A collocation is a natural combination (pattern) of words that are often used together (some usual patterns are: noun + verb, adjective + noun, verb + adverb, adverb + adjective, verb + noun, noun + noun, etc.).",
                                        "For each pattern provide 1 to 5 common collocations that are useful for language learners. Avoid rare or archaic collocations.",
                                        "Consider the collocation only (in kanji) and avoid extra explanations/translations.",
                                        "First write the pattern, then under it report the collocations with bullet point.",
                                    ]
                                ),
                            }
                        ],
                    },
                    {"role": "user", "content": [{"type": "input_text", "text": self.word}]},
                ],
                text={"format": {"type": "text"}},
                reasoning={},
                tools=[
                    {
                        "type": "web_search_preview",
                        "filters": None,
                        "search_context_size": "medium",
                        "user_location": {
                            "type": "approximate",
                            "city": None,
                            "country": None,
                            "region": None,
                            "timezone": None,
                        },
                    }  # type: ignore
                ],
                tool_choice={"type": "web_search_preview"},
                temperature=1,
                max_output_tokens=2024,
                top_p=1,
                store=False,
                include=["web_search_call.action.sources"],  # type: ignore
            )
            raw_collocations = response.output_text
            logger.debug(f"🖥️ raw collocations: {raw_collocations}")
        response = get_openai_client().responses.create(
            model="gpt-4o-mini",
            input=[
                {
                    "role": "system",
                    "content": [
                        {
                            "type": "input_text",
                            "text": 'you have to make the JSON format of the given text about collocations of a Japanese vocabulary.\nOutput JSON:\n{ "collocations": list of strings}',
                        }
                    ],
                },
                {"role": "user", "content": [{"type": "input_text", "text": raw_collocations}]},
            ],
            text={"format": {"type": "json_object"}},
            reasoning={},
            tools=[],
            tool_choice="none",
            temperature=1,
            max_output_tokens=2375,
            top_p=1,
            store=False,
            include=["web_search_call.action.sources"],  # type: ignore
        )
        return eval(response.output_text)["collocations"]

    def get_ai_explanation(self) -> dict:
        logger.info(f"🖥️ Getting AI explanation for word: {self.word}")
        response = get_openai_client().responses.create(
            model="gpt-5-mini",
            input=[
                {
                    "role": "developer",
                    "content": [
                        {
                            "type": "input_text",
                            "text": "\n".join(
                                [
                                    "You are a helpful assistant that explains Japanese words in detail with respect to the given data collected from Jisho.org and KanjiAPI.dev.",
                                    "Provide clear, concise, comprehensive and short explanations suitable for language learners.",
                                    "The explanation should be in spoken form, as if you are explaining to a friend.",
                                    "Make sure to strictly follow the output format specified in the user prompt. No extra text outside the JSON object.",
                                ]
                            ),
                        }
                    ],
                },
                {
                    "role": "user",
                    "content": "\n".join(
                        [
                            f"The target word is '{self.word}'. Use the following data:\n\nJisho Data: {self.jisho_data}\n\nKanji Data: {self.kanji_data if self.kanji_data else 'No kanji data available.'}\n\nCollocations: {self.collocations}\n\nProvide a comprehensive explanation suitable for a language learner.",
                            "## introduction_japanese: ",
                            "Without mentioning the meaning and reading of the word, name the situations, contexts, and occasions in which this word is used. The explanation should be in Japanese, short and in spoken form suitable for elementary level language learners. Start with the word itself.",
                            "Example: 「はげしい」は、感情、天気、強調が必要な状況についての会話でよく見かける形容詞です。",
                            "## introduction_english: ",
                            "Translation of introduction_japanese in English. Write the word in its hiragana/katakana form. Start with 'The [adjective/noun/verb ...] [word] ...",
                            "Example: 'はげしい is an adjective commonly seen in conversations about emotions, weather, and situations that require emphasis.'",
                            "## meanings: ",
                            "List of all meanings grouped by nuance. Each nuance should have its own list of meanings. If there is no nuance, put all meanings in a single list.",
                            "## meaning_explanation_japanese: ",
                            "Short but complete explanation of literal meanings of the word in Japanese in spoken form suitable for elementary level language learners.",
                            "Example: 「はげしい」は、強い感情や激しい動作を表す言葉です。例えば、激しい雨や激しい議論など、何かが非常に強いことを示します。",
                            "## meaning_explanation_english: ",
                            "Shortly provide the literal meanings of the word. Write the word in its hiragana/katakana form. The explanation should be spoken form and suitable for elementary level language learners.",
                            "Example: The word はげしい means 'intense', 'fierce', or 'violent'. This word can describe actions, emotions, or conditions that are of great intensity or force.",
                            "## kanji_details: ",
                            "For each kanji in the word, provide its meanings, readings (on'yomi and kun'yomi), and 1 or 2 (max) common words using this kanji except the main word. For each common word, provide its reading and meaning.",
                            "## kanji_explanation_english: ",
                            "Write in a natural, conversational transcript of a teacher explaining the kanji and its meanings. for each kanji (in the order it appears), compose one 3-4 short sentence paragraph that: 1. Describes the kanji's core meaning. 2. except the original word, Presents the 1-2 vocabularies that use that kanji, and how this kanji gives meaning in this vocabulary. Constrains: * Explanation field is the transcription of a speech. Don't use bullet points, parenthesis, new lines, titles, or anything similar. * Do not include the original word in the vocabs. * In the explanation field only insert the hiragana for of Japanese vocabs. No kanjis. * Explanation starts with English phrases such as:  The [first/second/...] kanji means ...",
                            "## synonyms: ",
                            "List the 1 (maximum 2) most commonly used synonym for the provided Japanese vocabulary word (no readings or any other extra text, perferebly in kanji). Excluding the original word.",
                            "## synonyms_explanation: ",
                            'provide the English transcription of a very short explanation about the synonyms listed, including their nuances and meanings. If there were no synonyms in the list, say "No common synonyms found." Constraints 1. Only insert the hiragana for of Japanese vocabs. No kanjis. 2. Explanation starts with English phrases such as:  The most common synonym[s] of the [word] [are/is] ... 3. Very shortly explain the nuances of each synonym and antonym listed, and how they look like the original word.',
                            "## antonyms: ",
                            "List the 1 (maximum 2) most commonly used antonym for the provided Japanese vocabulary word (no readings or any other extra text, perferebly in kanji). Excluding the original word.",
                            "## antonyms_explanation: ",
                            'provide the English transcription of a very short explanation about the antonyms listed, including their nuances and meanings. If there were no antonyms in the list, say "No common antonyms found." Constraints 1. Only insert the hiragana for of Japanese vocabs. No kanjis. 2. Explanation starts with English phrases such as:  The most common antonym[s] of the [word] [are/is] ... 3. Very shortly explain the nuances of each synonym and antonym listed, and how they differ from the original word.',
                            "## Examples: ",
                            "With respect to the provided collocations, provide 5 - 7 example short and simple sentences using the target word in different contexts useful for language learners.",
                            "Each example should be in kanji and its furigana form. In the furigana form, provide the reading of kanjis in parenthesis right after the kanji. If there are no kanjis in the sentence, just write the sentence as is.",
                            "For example for the word 'はげしい', an example sentence in kanji could be '彼ははげしい感情を持っている。' and its furigana would be '彼(かれ)ははげしい感情(かんじょう)を持(も)っている。'.",
                            "make sure the sentences are short, simple and easy to understand for elementary level language learners.",
                            "make sure you follow the rules for furigana.",
                            "## output_format: ",
                            "The output should be in JSON format:",
                            """
{
  "kanji": the word's kanji (if the kanji of the word is not commonly used, write the word itself - no kanji),
  "reading": the word's reading,
  "introduction_japanese": string,
  "introduction_english": string,
  "meanings": [[list of all meanings in nuance1], [list of all meanings in nuance2], ...],
  "meaning_explanation_japanese": string,
  "meaning_explanation_english": string,
  "kanji_explanation_english": string
  "kanji_details": [
    {
      "kanji": "the kanji character",
      "common_words":  # 1 or 2 common words using this kanji except the main word
        ["word1 in kanji (word1 in hiragana): meaning", "word2 in kanji (word2 in hiragana): meaning"]
    },
    ...
  ],
  "synonyms": ["synonym1 in kanji : synonym1 in hiragana : meaning", ...],
  "synonyms_explanation": string,
  "antonyms": ["antonym1 in kanji : antonym1 in hiragana : meaning", ...],
  "antonyms_explanation": string,
  "examples": [{"kanji": string, "furigana": string}, ...]
}
""",
                        ]
                    ),
                },
            ],
            text={"format": {"type": "json_object"}, "verbosity": "medium"},
            reasoning={
                "effort": "medium",
                "summary": None,
            },
            tools=[],
            tool_choice="none",
            store=False,
            include=["reasoning.encrypted_content", "web_search_call.action.sources"],  # type: ignore
        )
        logger.debug(f"🟩 AI explanation response: {response.output_text}")
        return eval(response.output_text)

    def translate(self):
        logger.info(f"🟩 Translating AI explanation for word: {self.word}")
        self.ai_explanation["meaning_translations"] = {}
        for nuance in self.ai_explanation.get("meanings", []):
            nuance_translations = translate_to_all_languages(nuance[0], source_language="en")
            self.ai_explanation["meaning_translations"][nuance[0]] = nuance_translations
            self.ai_explanation["meaning_translations"][nuance[0]]["EN"] = nuance[1:]
        for example in self.ai_explanation.get("examples", []):
            example["translations"] = translate_to_all_languages(example["kanji"])

    def to_dict(self) -> dict:
        return {
            "version": "0.2.0",
            "word": self.word,
            "jlpt_level": self.jlpt_level,
            "youtube_link": self.youtube_link,
            "in_db": self.in_db,
            "ai_explanation": self.ai_explanation,
            "jisho_data": self.jisho_data,
            "kanji_list": self.kanji_list,
            "kanji_data": self.kanji_data,
            "collocations": self.collocations,
        }

    def save_to_json(self) -> None:
        dir_path = os.path.join("output", self.word)
        if not os.path.exists(dir_path):
            os.makedirs(dir_path)
        file_path = os.path.join(dir_path, f"{self.word}.json")
        with open(file_path, "w", encoding="utf-8") as f:
            json.dump(self.to_dict(), f, ensure_ascii=False, indent=4)

    @staticmethod
    def load_from_json(word: str) -> "JPWord2":
        # file_path = os.path.join("output", word, f"{word}.json")
        file_path = os.path.join("resources", "words", f"{word}.json")
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        obj = JPWord2(data["word"], data["jlpt_level"], ai_init=False)
        obj.youtube_link = data["youtube_link"]
        obj.in_db = data["in_db"]
        obj.ai_explanation = data["ai_explanation"]
        obj.jisho_data = data["jisho_data"]
        obj.kanji_list = data["kanji_list"]
        obj.kanji_data = data["kanji_data"]
        obj.collocations = data["collocations"]
        return obj

    def __str__(self) -> str:
        from pprint import pformat

        return pformat(
            self.to_dict(),
            indent=1,
            width=120,
            compact=False,
            sort_dicts=False,
        )

    def pptx_generation(self, num_examples: int | None = 4) -> None:
        logger.info(f"🟦 Generating PPTX for word: {self.word}")
        file_name = f"./Output/{self.word}/{self.word} JLPT N{self.jlpt_level} Vocabulary.pptx"
        if os.path.exists(file_name):
            return

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import (
            MSO_AUTO_SIZE,
            MSO_VERTICAL_ANCHOR,
            PP_PARAGRAPH_ALIGNMENT,
        )
        from pptx.util import Inches, Pt

        prs = Presentation("resources/pptx_templates/template.pptx")

        # Title slide
        first_slide = prs.slides.add_slide(prs.slide_layouts[0])
        presentation_title = first_slide.shapes.title
        if presentation_title:
            presentation_title.text = self.ai_explanation.get("kanji", self.word)
            presentation_title.text_frame.paragraphs[0].font.size = Pt(160)
            presentation_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 95, 154)

        presentation_subtitle = first_slide.placeholders[1]
        presentation_subtitle.text = self.ai_explanation.get("reading", self.word)  # type: ignore
        presentation_subtitle.text_frame.paragraphs[0].font.size = Pt(100)  # type: ignore
        presentation_subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(192, 79, 21)  # type: ignore

        first_slide.shapes.add_movie(
            f"./output/{self.word}/audio/0_introduction.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Definitions slide
        definitions_slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = definitions_slide.shapes.add_textbox(Inches(5), Inches(0), Inches(40 / 3 - 5), Inches(7.5))
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        shape.text_frame.word_wrap = True
        for en, translations in self.ai_explanation["meaning_translations"].items():
            p = shape.text_frame.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            run.text = en + "\n"
            run.font.size = Pt(115)
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(33, 95, 154)
            if len(translations["EN"]):
                run = p.add_run()
                run.text = ", ".join(translations["EN"]) + "\n"
                run.font.size = Pt(32)
                run.font.color.rgb = RGBColor(192, 79, 21)
                run.font.name = "Berlin Sans FB"
                p.space_after = Pt(0)
                p.line_spacing = 0.9
            for code, translation in translations.items():
                if code.lower() == "en":
                    continue
                run_code = p.add_run()
                run_code.text = f"     {code}"
                run_code.font.size = Pt(16)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)

                # Add ordinary run for the translation
                run_translation = p.add_run()
                run_translation.text = f"{translation}"
                run_translation.font.size = Pt(20)
                run_translation.font.name = "Berlin Sans FB"
                run_translation.font.color.rgb = RGBColor(0, 0, 0)

        definitions_slide.shapes.add_movie(
            f"./output/{self.word}/audio/1_definition.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Add a slide for Kanji breakdown
        kanji_slide = prs.slides.add_slide(prs.slide_layouts[6])
        for i, kanji in enumerate(self.kanji_list):
            kanji_shape = kanji_slide.shapes.add_textbox(
                Inches(0), Inches(i * 7.5 / 2), Inches(40 / 3 * 0.3), Inches(7.5 / 2)
            )
            p = kanji_shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = kanji
            run.font.size = Pt(250)
            run.font.bold = True
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(33, 95, 154)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.space_after = Pt(0)
            p.line_spacing = 0.9

            kanji_explanation_shape = kanji_slide.shapes.add_textbox(
                Inches(40 / 3 * 0.3),
                Inches(i * 7.5 / 2),
                Inches(40 / 3 * 0.7),
                Inches(7.5 / 2),
            )
            kanji_explanation_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            kanji_explanation_shape.text_frame.word_wrap = True
            p = kanji_explanation_shape.text_frame.paragraphs[-1]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            run = p.add_run()
            run.text = "Meanings:"
            run.font.size = Pt(24)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(127, 127, 127)
            run = p.add_run()
            run.text = f" {', '.join(self.kanji_data[kanji]['meanings'][:4])}\n"
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = "Readings:"
            run.font.size = Pt(24)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(127, 127, 127)
            run = p.add_run()
            run.text = f" {', '.join(self.kanji_data[kanji]['kun_readings'] + self.kanji_data[kanji]['on_readings'])}\n"
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            for vocab in self.ai_explanation["kanji_details"][i]["common_words"]:
                run.text += f"\n{vocab}"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(192, 79, 21)
            p.space_after = Pt(0)
            p.line_spacing = 0.9

        kanji_slide.shapes.add_movie(
            f"./output/{self.word}/audio/2_kanji_explanation.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        # Add a slide for examples
        for i, example in enumerate(self.ai_explanation["examples"][:num_examples]):
            example_slide = prs.slides.add_slide(prs.slide_layouts[6])

            top_shape = example_slide.shapes.add_textbox(Inches(0), Inches(0), Inches(40 / 3), Inches(3))
            top_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            top_shape.text_frame.word_wrap = True
            top_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = top_shape.text_frame.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            run.text = example["kanji"]
            run.font.size = Pt(70)
            run.font.bold = True
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(33, 95, 154)
            run = p.add_run()
            run.text = f"\n{example['furigana']}"  # type: ignore
            run.font.size = Pt(33)
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(192, 79, 21)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.space_after = Pt(0)
            p.line_spacing = 0.9

            left_languages = [
                "EN",
                "ID",
                "ES",
                "VI",
                "FR",
                "NE",
                "AR",
            ]

            left_shape = example_slide.shapes.add_textbox(
                Inches(20 / 3 * 0.025), Inches(3), Inches(20 / 3 * 0.95), Inches(4.5)
            )
            left_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            left_shape.text_frame.word_wrap = True
            left_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for lang in left_languages:
                p = left_shape.text_frame.paragraphs[-1]
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                run_code = p.add_run()
                run_code.text = f"{lang}: "
                run_code.font.size = Pt(20)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)
                run_translation = p.add_run()
                run_translation.text = f"{example['translations'][lang]}\n"
                run_translation.font.size = Pt(25)
                run_translation.font.name = "Berlin Sans FB"
            right_shape = example_slide.shapes.add_textbox(
                Inches(20 / 3 * 1.025), Inches(3), Inches(20 / 3 * 0.95), Inches(4.5)
            )
            right_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            right_shape.text_frame.word_wrap = True
            right_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            right_languages = [
                "BN",
                "ZH",
                "KO",
                "TL",
                "MY",
                "HI",
                "FA",
            ]
            for lang in right_languages:
                p = right_shape.text_frame.paragraphs[-1]
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                run_code = p.add_run()
                run_code.text = f"{lang}: "
                run_code.font.size = Pt(20)
                run_code.font.name = "Berlin Sans FB"
                run_code.font.color.rgb = RGBColor(127, 127, 127)
                run_translation = p.add_run()
                run_translation.text = f"{example['translations'][lang]}\n"  # type: ignore
                run_translation.font.size = Pt(25)
                run_translation.font.bold = False
                run_translation.font.name = "Berlin Sans FB"

            example_slide.shapes.add_movie(
                f"./output/{self.word}/audio/{3 + i}_example.wav",
                left=Pt(0),
                top=Pt(-50),
                width=Pt(50),
                height=Pt(50),
                mime_type="audio/x-wav",
            )

        # Add a slide for synonyms and antonyms
        synonyms_slide = prs.slides.add_slide(prs.slide_layouts[6])

        synonyms_shape = synonyms_slide.shapes.add_textbox(Inches(0), Inches(0), Inches(20 / 3), Inches(7.5))
        synonyms_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        synonyms_shape.text_frame.word_wrap = True
        p = synonyms_shape.text_frame.paragraphs[-1]
        run = p.add_run()
        run.text = "\nSYNONYMS\n\n"
        run.font.size = Pt(40)
        run.font.name = "Berlin Sans FB Demi"
        run.font.color.rgb = RGBColor(33, 95, 154)
        for synonym in self.ai_explanation.get("synonyms", []):
            kanji, reading, meaning = synonym.split(":")
            run = p.add_run()
            run.text = f"{kanji}\n"
            run.font.size = Pt(54)
            run.font.bold = True  # type: ignore
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(33, 95, 154)
            run = p.add_run()
            run.text = f"{reading}\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = f"{meaning}\n\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(0)
            p.line_spacing = 0.9
        synonyms_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        antonyms_shape = synonyms_slide.shapes.add_textbox(Inches(20 / 3), Inches(0), Inches(20 / 3), Inches(7.5))
        antonyms_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        antonyms_shape.text_frame.word_wrap = True
        p = antonyms_shape.text_frame.paragraphs[-1]
        run = p.add_run()
        run.text = "\nANTONYMS\n\n"
        run.font.size = Pt(40)
        run.font.name = "Berlin Sans FB Demi"
        run.font.color.rgb = RGBColor(192, 79, 21)
        for antonym in self.ai_explanation.get("antonyms", []):
            kanji, reading, meaning = antonym.split(":")
            run = p.add_run()
            run.text = f"{kanji}\n"
            run.font.size = Pt(54)
            run.font.bold = True  # type: ignore
            run.font.name = "Berlin Sans FB Demi"
            run.font.color.rgb = RGBColor(192, 79, 21)
            run = p.add_run()
            run.text = f"{reading}\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            run = p.add_run()
            run.text = f"{meaning} \n\n"  # type: ignore
            run.font.size = Pt(32)
            run.font.name = "Berlin Sans FB"
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(0)
            p.line_spacing = 0.9
        antonyms_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        synonyms_slide.shapes.add_movie(
            f"./output/{self.word}/audio/100_synonyms_antonyms.wav",
            left=Pt(0),
            top=Pt(-50),
            width=Pt(50),
            height=Pt(50),
            mime_type="audio/x-wav",
        )

        prs.save(file_name)

    def tts(self, num_examples: int | None = 4) -> None:
        os.makedirs(f"./Output/{self.word}/audio", exist_ok=True)

        # Generate audio for introduction
        introduction_en_audio_path = f"./output/{self.word}/audio/introduction_en.mp3"
        if not os.path.exists(introduction_en_audio_path):
            logger.info("🔊 Generating audio for english introduction")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['introduction_english']}",
                instructions=f"English mixed with Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(introduction_en_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        introduction_jp_audio_path = f"./output/{self.word}/audio/introduction_jp.mp3"
        if not os.path.exists(introduction_jp_audio_path):
            logger.info("🔊 Generating audio for japanese introduction")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['introduction_japanese']}",
                instructions=f"Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(introduction_jp_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        introduction_audio = AudioSegment.silent(duration=100)
        introduction_audio = AudioSegment.from_mp3(introduction_en_audio_path).apply_gain(12)
        introduction_audio += AudioSegment.silent(duration=500)
        introduction_audio += AudioSegment.from_mp3(introduction_jp_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/0_introduction.wav", "wb") as title_file:
            introduction_audio.export(title_file, format="wav")

        # Generate audio for meaning explanation
        definition_en_audio_path = f"./output/{self.word}/audio/definition_en.mp3"
        if not os.path.exists(definition_en_audio_path):
            logger.info("🔊 Generating audio for english definition")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['meaning_explanation_english']}",
                instructions=f"English mixed with Japanese. Calmly and gently. Correct pronunciation: {self.word}:{self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(definition_en_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        definition_jp_audio_path = f"./output/{self.word}/audio/definition_jp.mp3"
        if not os.path.exists(definition_jp_audio_path):
            logger.info("🔊 Generating audio for japanese definition")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=f"{self.ai_explanation['meaning_explanation_japanese']}",
                instructions=f"Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(definition_jp_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        definition_audio = AudioSegment.silent(duration=100)
        definition_audio = AudioSegment.from_mp3(definition_en_audio_path).apply_gain(12)
        definition_audio += AudioSegment.silent(duration=500)
        definition_audio += AudioSegment.from_mp3(definition_jp_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/1_definition.wav", "wb") as word_file:
            definition_audio.export(word_file, format="wav")

        # Generate audio for kanji explanations
        kanji_explanation_audio_path = f"./output/{self.word}/audio/kanji_explanation.mp3"
        if not os.path.exists(kanji_explanation_audio_path):
            logger.info("🔊 Generating audio for kanji explanation")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.ai_explanation["kanji_explanation_english"],
                instructions=f"English. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(kanji_explanation_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        kanji_explanation_audio = AudioSegment.from_mp3(kanji_explanation_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/2_kanji_explanation.wav", "wb") as kanji_file:
            kanji_explanation_audio.export(kanji_file, format="wav")

        # Generate audio for examples
        for i, example in enumerate(self.ai_explanation["examples"][:num_examples]):
            example_jp_audio_path = f"./output/{self.word}/audio/example_{i + 1}_jp.mp3"
            if not os.path.exists(example_jp_audio_path):
                logger.info(f"🔊 Generating audio for example {i + 1}")
                response = get_openai_client().audio.speech.create(
                    model="gpt-4o-mini-tts",
                    voice="coral",
                    input=example["kanji"],
                    instructions=f"Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                    response_format="mp3",
                    speed=0.95,
                )
                with open(example_jp_audio_path, "wb") as audio_file:
                    audio_file.write(response.content)
            example_en_audio_path = f"./output/{self.word}/audio/example_{i + 1}_en.mp3"
            if not os.path.exists(example_en_audio_path):
                logger.info(f"🔊 Generating audio for example {i + 1} translation")
                response = get_openai_client().audio.speech.create(
                    model="gpt-4o-mini-tts",
                    voice="coral",
                    input=example["translations"]["EN"],
                    instructions=f"English. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                    response_format="mp3",
                    speed=0.95,
                )
                with open(example_en_audio_path, "wb") as audio_file:
                    audio_file.write(response.content)
            example_jp_audio = AudioSegment.from_mp3(example_jp_audio_path).apply_gain(12)
            example_en_audio = AudioSegment.from_mp3(example_en_audio_path).apply_gain(12)

            example_audio = AudioSegment.silent(duration=100)
            example_audio += example_en_audio
            example_audio += AudioSegment.silent(duration=500)
            example_audio += example_jp_audio
            example_audio += AudioSegment.silent(duration=500)
            example_audio += example_jp_audio
            with open(f"./output/{self.word}/audio/{3 + i}_example.wav", "wb") as example_file:
                example_audio.export(example_file, format="wav")

        # Generate audio for synonyms and antonyms explanations
        synonyms_explanation_audio_path = f"./output/{self.word}/audio/synonyms_explanation.mp3"
        if not os.path.exists(synonyms_explanation_audio_path):
            logger.info("🔊 Generating audio for synonyms explanation")
            response = get_openai_client().audio.speech.create(
                model="gpt-4o-mini-tts",
                voice="coral",
                input=self.ai_explanation["synonyms_explanation"] + " " + self.ai_explanation["antonyms_explanation"],
                instructions=f"English mixed with Japanese. calmly and gently. Correct pronunciation: {self.word} = {self.ai_explanation['reading']}",
                response_format="mp3",
                speed=0.95,
            )
            with open(synonyms_explanation_audio_path, "wb") as audio_file:
                audio_file.write(response.content)
        synonyms_explanation_audio = AudioSegment.from_mp3(synonyms_explanation_audio_path).apply_gain(12)
        with open(f"./output/{self.word}/audio/100_synonyms_antonyms.wav", "wb") as syn_ant_file:
            synonyms_explanation_audio.export(syn_ant_file, format="wav")

    def show_in_streamlit(self, st, auth: dict | None = None) -> None:
        if self.word == self.ai_explanation.get("reading", ""):
            st.markdown(
                f"# {self.word} :orange-badge[N{self.jlpt_level}] :green-badge[{'Common' if self.jisho_data['is_common'] else 'Uncommon'}]"
            )
        else:
            st.markdown(
                f"# {self.word} ({self.ai_explanation['reading']}) :orange-badge[N{self.jlpt_level}] :green-badge[{'Common' if self.jisho_data['is_common'] else 'Uncommon'}]"
            )
        if self.youtube_link:
            st.video(self.youtube_link)

        st.markdown(self.ai_explanation["introduction_english"])
        st.markdown(self.ai_explanation["introduction_japanese"])
        st.markdown("### Meanings")
        st.markdown(self.ai_explanation["meaning_explanation_english"])
        st.markdown(self.ai_explanation["meaning_explanation_japanese"])
        for m, tr in self.ai_explanation["meaning_translations"].items():
            tr["EN"].insert(0, m)
            tr["EN"] = ", ".join(tr["EN"])
            with st.container(border=1, horizontal=True):
                for k, v in tr.items():
                    if auth:
                        user_langs = [LANGUAGES_ABBR[lang] for lang in auth.get("preferred_languages", [])]
                        if k not in user_langs:
                            continue
                    st.markdown(f":gray-badge[{k}] {v}")

        col1, col2 = st.columns(2)
        with col1:
            st.markdown("### Synonyms")
            with st.container(border=1):
                for s in self.ai_explanation["synonyms"]:
                    st.markdown(s)
        with col2:
            st.markdown("### Antonyms")
            with st.container(border=1):
                for a in self.ai_explanation["antonyms"]:
                    st.markdown(a)

        st.markdown("### Kanji")
        for i, k in enumerate(self.kanji_list):
            st.markdown(f"#### {k}")
            cont = st.container(border=1)
            col1, col2 = cont.columns([2, 4])
            with col1:
                st.image(
                    f"https://raw.githubusercontent.com/KanjiVG/kanjivg/refs/heads/master/kanji/0{self.kanji_data[k]['unicode'].lower()}.svg",
                    width="stretch",
                )
            with col2:
                st.markdown(f":gray-badge[Meaning] {', '.join(self.kanji_data[k]['meanings'])}")
                st.markdown(f":gray-badge[On-yomi] {', '.join(self.kanji_data[k]['on_readings'])}")
                st.markdown(f":gray-badge[Kun-yomi] {', '.join(self.kanji_data[k]['kun_readings'])}")
                st.markdown("---")
                for common_word in self.ai_explanation["kanji_details"][i]["common_words"]:
                    st.markdown(common_word)

        st.markdown("### Examples")
        for ex in self.ai_explanation["examples"]:
            from src.utils import create_html_with_ruby

            with st.expander(ex["kanji"]):
                ruby = create_html_with_ruby(ex["furigana"])
                st.markdown(ruby, unsafe_allow_html=True)
                for key, value in ex["translations"].items():
                    if auth:
                        if key not in auth.get("preferred_languages", []):
                            continue
                    st.markdown(f":gray-badge[{key}] {value}")


if __name__ == "__main__":
    # "上達": "名詞 + 動詞 *（「〜が／を上達する／させる／望む」など） *技術が上達する *能力が上達する *実力が上達する *練習して上達する *習い事を上達させる *副詞 + 上達する／上達させる *（程度や様子を表す副詞と組む） *飛躍的に上達する *着実に上達する *急速に上達する *少しずつ上達する *目覚ましく上達する *上達 + の + 名詞 *（上達した結果・性質を表す名詞） *上達の過程 *上達の目安 *上達の速度 *上達の程度 *上達の跡 *上達 + が + 形容詞 *（「上達が早い／遅い／著しい／順調／目立つ」など） *上達が早い *上達が遅い *上達が著しい *上達が順調だ *上達が目立つ *動詞 + 上達 *（他の動詞が「上達」にかかるパターン） *実感して上達を感じる *努力して上達を図る *様子を見て上達を判断する *継続して上達を目指す *比較して上達を見る",
    # "うまい": "副詞 + うまい／上手い／旨い *（様子・程度を強めたり緩めたり） *とても／すごくうまい *すごく上手い *かなりうまい *ずいぶんうまい *案外うまい *うまい／上手い／旨い + 名詞 *（その“うまさ”を対象化する） *うまい料理 *上手い表現 *うまい話 *上手い使い方 *上手い演技 *動詞 + うまく／上手に／旨く + 動詞 *（行為・動作に対して「うまく〜する」など） *うまくやる *上手に話す *うまく伝える *上手に使う *うまくまとめる *〜が + うまい／上手い／旨い *（主語の能力・技巧を表す） *彼は上手い（だ） *君がうまい *あの人がうまい *誰が上手いか *それがうまい（使われる） *うまい + （形容詞・句） *（「うまい」が修飾語句と組まれる形） *うまい具合に *うまい加減に *うまい調子で *うまい話に乗る *うまいタイミングで",
    # "不安": "形容詞 + 不安 * 強い不安 * 大きな不安 * 重大な不安 * 根強い不安 * 常に不安 * 不安 + を + 動詞 * 不安を感じる * 不安を抱く * 不安を覚える * 不安を煽る * 不安を和らげる * 不安 + が + 形容詞／動詞 * 不安がある * 不安が消える * 不安が募る * 不安が高まる * 不安が和らぐ * 副詞 + 不安 * 非常に不安 * やや不安 * ただ不安 * 何となく不安 * 少し不安 * 不安 + の + 名詞 * 不安の種 * 不安の声 * 不安の原因 * 不安の念 * 不安の兆し",
    # "丸い": "形容動詞 + 名詞 * 丸い形 * 丸い頭 * 丸い目 * 丸い顔 * 丸い机 * 丸い + 名詞 + を + 動詞 * 丸い石を投げる * 丸いコインを手に取る * 丸い皿を洗う * 丸い風船を膨らませる * 丸いボールを打つ * 副詞・修飾語 + 丸い * 非常に丸い * 比較的丸い * 丸くて光る * きれいに丸い * ほぼ丸い * 名詞 + が + 丸い／丸く * 月が丸い * 顔が丸い * 目が丸い * お腹が丸い * 魚が丸い * 慣用表現・比喩的表現 * 丸く収める * 丸くなる * 丸く治める * 丸く言う * 丸く収まる",
    # "乗客": "名詞 + 動詞／動詞＋名詞 乗客を乗せる 乗客を降ろす 乗客が乗る 乗客が下車する 乗客を案内する 〜に + 乗客 （乗客を目的語・対象場面で取る表現） 列車に乗客を乗せる 飛行機に乗客を乗せる バスに乗客を載せる 電車に乗客が殺到する 駅に乗客があふれる 乗客 + の + 名詞 乗客数 乗客案内 乗客サービス 乗客層 乗客運送 乗客 + が + 動詞／形容詞 乗客が多い 乗客が減る 乗客が増える 乗客が混雑する 乗客が苦情を言う 副詞／修飾語 + 乗客 多くの乗客 少数の乗客 主な乗客 観光客乗客 巨大な乗客数",
    # "予期": "名詞 + 動詞（「〜を予期する」型） 事態を予期する 結果を予期する 変化を予期する 反応を予期する 事前に予期する 副詞 + 動詞（予期とともに使われる副詞） 十分に予期する おおよそ予期する 全く予期しない まったく予期しなかった 予期せず 名詞 + に + 動詞（「〜に予期される／〜に予期されない」など） 予期に反する 予期に沿う 予期に達する 予期に届く 予期に違わず 名詞＋形容詞（名詞「予期」と組む形容詞修飾） 予期せぬ出来事 予期外の反応 予期しがたい結果 予期どおりの展開 予期以上の成果 名詞 + 名詞（「〜の予期」などの構造） 将来の予期 期待との予期 予期の範囲 予期の違い 予期の変化",
    # "予防": "名詞 + を + 動詞（「〜を予防する／〜を予防できる」型） 感染を予防する 病気を予防する 事故を予防する 災害を予防する 虫歯を予防する 名詞 + の + 名詞（「〜の予防」「〜予防」型） 病気の予防 感染予防 火災予防 災害予防 健康予防（※あまり日常語ではないが使われることあり） 名詞 + に + 動詞／形容動詞（「〜に予防策を講じる」「〜に予防的な」など） 感染に対して予防策を講じる 病気に対して予防手段をとる 災害に備えて予防的措置をとる 被害に対する予防的対応 リスクに対する予防的処置 形容詞 + 名詞（「〜予防」「〜的予防」など修飾語） 予防接種 予防措置 予防的 予防効果 予防医療 副詞／慣用句 + 動詞（予防を表す表現補助語） 十分に予防する 事前に予防する 適切に予防する 徹底して予防する 予防できる／予防可能",
    # "事件": "名詞 + を + 動詞（「〜を事件とみなす／〜を扱う」など） 事件を扱う 事件を立件する 事件を捜査する 事件を報道する 事件を解決する 名詞 + の + 名詞（「〜事件」「〜の事件」型） 殺人事件 交通事故事件 窃盗事件 重大事件 未解決事件 名詞 + に + 動詞／形容詞 事件に発展する 事件に巻き込まれる 事件に関与する 事件に至る 事件に見舞われる 形容詞 + 名詞（「〜事件」「〜な事件」型） 衝撃的な事件 異常な事件 悲惨な事件 惨劇的な事件 意外な事件 副詞／接続語 + 名詞／動詞（「〜事件が起こる」「〜により事件が〜」など） 突如として事件が起こる 社会問題化する事件 次々と事件が発覚する 事件が相次ぐ 事件が表面化する",
    # "交換": "名詞 + を + 動詞（〜を交換する／〜を交換できる 等） 情報を交換する 名刺を交換する 連絡先を交換する 意見を交換する 部品を交換する 名詞 + の + 名詞（〜の交換／交換〜 等） 交換条件 交換手数料 交換所 交換期 交換制度 名詞 + と + 動詞（〜と交換する／交換し合う 等） 物と交換する 意見と交換し合う 書籍と交換する 情報と交換する 贈り物と交換する 形容詞 + 名詞（交換を修飾する形容語との組み合わせ） 相互交換 無償交換 等価交換 即時交換 自由交換 動詞 + 名詞（交換する対象を表す） 交換操作 交換業務 交換取引 交換対応 交換処理 副詞／補助語 + 動詞（交換の仕方・頻度を表す） 定期的に交換する 相互に交換する 直接交換する 部分的に交換する 速やかに交換する",
    # "代金": "名詞 + を + 動詞（「〜代金を〜する」型） 代金を支払う 代金を受け取る 代金を請求する 代金を立て替える 代金を回収する 名詞 + の + 名詞 支払代金 販売代金 商品代金 請求代金 後払い代金 名詞 + と + 動詞（交換・引換など） 商品と代金を交換する 代金と引き換えに渡す 物と代金を換える 代金と相殺する 代金と対象を対応させる 形容詞 + 名詞 前払い代金 後払い代金 一括代金 分割代金 預り代金 副詞／補助語 + 名詞／動詞 あらかじめ代金 全額代金 一部代金 正式に代金を支払う ただちに代金を振込む",
    # "仲間": "名詞 + を + 動詞（「〜を仲間に加える／〜を仲間とする」型） 仲間を増やす 仲間を募る 仲間を集める 仲間をつくる 仲間を信頼する 名詞 + と + 動詞（「〜と仲間になる／仲間と共に〜する」型） 友と仲間になる 仲間と協力する 仲間と過ごす 仲間と連携する 仲間と分かち合う 名詞 + の + 名詞（「〜の仲間」「仲間の〜」型） 仲間の一人 仲間意識 仲間関係 仲間外れ 仲間入り 形容詞 + 名詞（仲間を修飾する表現） 親しい仲間 本当の仲間 心強い仲間 多様な仲間 新しい仲間 副詞／補助語 + 動詞・名詞（仲間と含める表現） 気軽に仲間と 一緒に仲間と もはや仲間ではない 仲間として認める 仲間である",
    # "任せる": "名詞 + を + 名詞 + に + 任せる（〜を〜に任せる型） 仕事を先輩に任せる 責任を部下に任せる 判断を他人に任せる 決定を幹部に任せる 管理を専門業者に任せる 名詞 + を + 動詞句 + に + 任せる（〜を〜することを任せる型） 掃除を全面的に任せる 交渉を丸ごと任せる 企画をすべて任せる 片付けを後任者に任せる 手配を現地スタッフに任せる 名詞 + に + 任せる（主体を示す表現） 部下に任せる 専門家に任せる 家族に任せる 友人に任せる プロに任せる 補助表現・副詞 + 任せる（任せるの仕方・度合いを表す） 完全に任せる 思い切って任せる 安心して任せる すべて任せる ほとんど任せる 慣用形・定型表現 任せてください 任せっきりにする／任せっきりだ 任せるしかない ～に任せておく ～に任される",
    # "休暇": "名詞 + を + 動詞（「〜を取る／〜を申請する」型） 休暇を取る 休暇を申請する 休暇を取得する 休暇を延長する 休暇を与える 名詞 + 中／間（状態・期間を表す） 休暇中 休暇期間中 休暇の間 長期休暇中 労働休暇中 名詞 + の + 名詞（休暇を修飾する語との組み合わせ） 有給休暇 年次休暇 病気休暇 産休・育休（産前・育児休暇） 特別休暇 名詞 + に + 動詞（目的・理由を導く構文など） 休暇にあてる 休暇に当てる 休暇に出る 休暇に充てる 休暇に入る 形容詞／修飾語 + 名詞（休暇を修飾する語） 長期休暇 短期休暇 連休休暇 無給休暇 完全休暇",
    # "作業": "名詞 + を + 動詞（〜を作業する／〜を進める 等） 作業を進める 作業を行う 作業を開始する 作業を完了する 作業を中断する 名詞 + 中／間（進行中・期間を表す） 作業中 作業の途中 作業期間中 作業の合間 夜間作業中 名詞 + の + 名詞（作業を修飾する語との組み合わせ） 作業工程 作業手順 作業効率 作業環境 作業負荷 副詞／修飾語 + 動詞（作業するのしかた・程度を表す語） 丁寧に作業する 急いで作業する 正確に作業する 手早く作業する 徐々に作業を進める 名詞 + に + 動詞（作業に関与する・作業対象を表す表現） 作業に携わる 作業に従事する 作業に参加する 作業に取りかかる 作業に着手する",
    # "作物": "名詞 + を + 動詞（～を作物とする／～を育てる 等）  作物を育てる  作物を収穫する  作物を栽培する  作物を植える  作物を保護する  名詞 + の + 名詞（～の作物／作物の～）  農作物（の作物）  作物の収量  作物の品質  作物の出荷  作物の生育  形容詞／修飾語 + 名詞（作物を修飾する語）  食用作物  工芸作物  農業作物  主作物  副作物  名詞 + に + 動詞（作物に関わる動詞句）  作物に被害を与える  作物に影響を及ぼす  作物に適する  作物に対応する  作物に栄養を与える  副詞／補助語 + 動詞・名詞（頻度・程度・状態を表す）  多くの作物  よい作物  十分な作物  豊作の作物  作物全体",
    # "信用": "名詞 + を + 動詞（「〜を信用する／〜を得る」型）  人を信用する  相手を信用する  自分を信用する  報告を信用する  証言を信用する  名詞 + の + 名詞（信用を修飾する語、信用の種類など）  信用度  信用調査  信用格付け  信用取引  信用保証  名詞 + と + 動詞（信用を前提とする構文）  信用と裏切り  信用と信頼を結ぶ  信用と期待を寄せる  信用と実績を重ねる  信用と義務を果たす  形容詞／修飾語 + 名詞（信用を強調・条件づける語）  高い信用  強い信用  絶対の信用  社会的信用  長期信用  副詞／補助語 + 動詞（信用するのしかた・程度を表す語）  安易に信用する  安心して信用する  完全に信用する  半ば信用する  十分に信用する  動詞句 + 名詞（信用を対象とする動詞との組み合わせ）  信用を失う  信用を回復する  信用を築く  信用を損なう  信用を高める",
    # "修正": "名詞 + を + 動詞（〜を修正する型）  誤りを修正する  文章を修正する  計画を修正する  デザインを修正する  スケジュールを修正する  名詞 + の + 名詞（修正を修飾／修正対象を示す語）  修正案  修正点  修正後  修正前  修正作業  補助語／副詞 + 動詞（修正するのしかた・度合いを表す語）  少し修正する  軽く修正する  大幅に修正する  再度修正する  最終修正する  名詞 + に + 動詞（修正を加える対象・目的を示す構文）  ～に修正を加える  仕様に修正を加える  条文に修正を加える  表現に修正を加える  内容に修正を加える  慣用表現・複合語  軌道修正  修正バージョン  修正対応  修正フィードバック  修正可能",
    # "値": "名詞 + を + 動詞（「〜を〜する／〜を〜とする」型）  値を付ける  値を下げる  値を上げる  値を見直す  値を設定する  名詞 + の + 名詞（値を修飾する語・値の種類など）  値段  値幅  値上げ  値下げ  最低値／最高値  名詞 + が + 動詞（値が〜する・値が〜だという表現）  値が上がる  値が下がる  値が変動する  値が安定する  値が高い／低い  動詞 + 名詞（動詞と組む名詞で値を表す表現）  値を表示する  値を比較する  値を記録する  値を報告する  値を予測する  補助語・副詞 + 動詞（値を扱う動詞表現を修飾する語）  適正に値を付ける  大幅に値を下げる  急激に値が上がる  わずかに値を変える  比較的高い値",
    # "僅か": "形容動詞 + 名詞 “わずかな + 名詞” の形で、量・程度・時間などの少なさを表す 僅かな時間 僅かな差 僅かな利益 僅かな期待 僅かな可能性 副詞（僅かに） + 動詞・形容詞 “わずかに + ～する／～である” の形で、ほどよく／わずかに変化・状態を示す 僅かに見える 僅かにずれる 僅かに違う 僅かに傾く 僅かに聞こえる 名詞 + 格助詞 + 僅か 名詞を主語や対象とし、「〜が僅か（である／しかない）」のような表現 残りが僅か 差が僅か 余裕が僅か 在庫が僅か 変化が僅か 数詞 + 僅か 数を限定して「たった…だけ」の意味合いを強める句 僅か1人 僅か数分 僅か数日 僅か1回 僅か数個 僅か + 助詞（〜に／〜で／〜しか） “僅かに” や “僅かで／しか” の形で修飾語的に用いる 僅かに変化する 僅かに違う 僅かである 僅かでしかない",
    # "典型": "名詞 + の + 形容詞／名詞 （「典型の～」という表現） 典型の例 典型のパターン 典型の事例 典型の特徴 典型の姿 ～として + 典型（的） + 名詞 （何かを典型例として位置づける表現） ～として典型例 ～として典型像 ～として典型モデル ～として典型的事例 ～として典型的パターン 形容詞 + 典型（例・姿・像など） （典型を修飾する言葉付き表現） 最も典型的な例 古典的で典型的なパターン わかりやすい典型像 標準的な典型例 一般的な典型像 ～は + 典型 + だ / ～が + 典型 + だ （述語として使う表現） これは典型だ 彼／彼女は典型だ その意見は典型だ この現象が典型だ そのケースが典型だ 動詞 + 典型（例・姿など） （動作とともに「典型」を使う表現） 示す典型例 見せる典型像 成る／なる典型例 採る典型例 形成する典型像",
    # "冒険": "名詞 + の + 名詞 （「冒険の～」という表現） 冒険の旅 冒険の始まり 冒険の物語 冒険の舞台 冒険の記録 名詞 + を + 動詞 （「冒険を～する」など） 冒険をする 冒険を始める 冒険を求める 冒険を重ねる 冒険を恐れる ～に + 冒険（する・出る） （目的地や場面を示す） ～に冒険する ～に冒険に出る ～に冒険心を抱く ～に冒険を挑む ～に冒険を持ちかける 形容詞 + 冒険／冒険的 （冒険を修飾する表現） 刺激的な冒険 危険な冒険 未知の冒険 壮大な冒険 冒険的な計画 冒険 + 名詞 （「冒険～」の複合語的な使い方） 冒険小説 冒険譚 凍険者（冒険者） 冒険旅行 冒険心",
    # "出席": "名詞 + の + 出席 正式の出席 全員の出席 来賓の出席 必要な出席 臨時の出席 出席 + を + 動詞 出席を求める 出席を取る 出席を確認する 出席を促す 出席を保障する ～に + 出席（する／できる／見送る） 会議に出席する 授業に出席する 式典に出席できる 総会に出席する 授業を出席する（少し重複だが慣用） 形容詞 + 出席 欠席と出席 全員出席 多数出席 正式出席 名誉出席 出席 + 名詞 出席率 出席者 出席名簿 出席記録 出席簿",
    # "利益": "名詞 + の + 名詞 利益の追求 利益の確保 利益の配分 利益の最大化 利益の還元 利益 + を + 動詞 利益を得る／あげる 利益を追求する 利益を確保する 利益を分配する 利益を再投資する ～に + 利益（する／得る／還元する） 顧客に利益を還元する 市場に利益をもたらす 社会に利益をもたらす 株主に利益を還元する 企業に利益をもたらす 形容詞 + 利益／利益的 純粋な利益 実質的な利益 莫大な利益 潜在的な利益 短期的な利益 利益 + 名詞 利益率 利益額 利益拡大 利益構造 利益分配",
    # "制度": "名詞 + の + 名詞 制度の導入 制度の改正 制度の確立 制度の維持 制度の見直し 制度 + を + 動詞 制度を設ける 制度を導入する 制度を改める／改正する 制度を適用する 制度を撤廃する ～に + 制度（を／が） 国に制度を導入する 同様に制度が存在する 市場に制度を設ける 社会に制度を根付かせる 組織に制度を整備する 形容詞 + 制度／制度的 公的な制度 法的な制度 新しい制度 既存の制度 柔軟な制度 制度 + 名詞 制度改革 制度設計 制度運用 制度改正 制度基準",
    # "加減": "動詞 + 加減（を／する） 加減をする 加減を調整する 加減を見る 加減を考える 加減をつける 名詞 + の + 加減 塩の加減 火の加減 照明の加減 湯の加減 焼き加減 ～具合 + 加減／～かげん 具合の加減 体の加減（例：「お加減はいかがですか」） 痛みの加減 気分の加減 調子の加減 接尾辞的用法（～加減） （動詞・状態を表す語に付いて、「～の具合／やや～な様子」を表す） 焼き加減 焼け加減 書き加減 歩き加減 うつむき加減 形容詞・副詞 + 加減 いい加減 適切な加減 程よい加減 非常に加減（あまり使われないが文脈で） やや加減",
    # "化粧": "名詞 + の + 名詞 化粧の仕方 化粧の種類 化粧の効果 化粧の道具 化粧の技術 化粧 + を + 動詞 化粧をする 化粧を落とす 化粧を直す 化粧を濃くする 化粧を薄くする ～に + 化粧（する／される） 顔に化粧する 頬に化粧される 目元に化粧を施す 人に化粧をする 拡張：壁に化粧を施す 形容詞 + 化粧 / 化粧的 濃い化粧 薄い化粧 きれいな化粧 自然な化粧 派手な化粧 化粧 + 名詞 化粧品 化粧道具 化粧水 化粧室 化粧筆",
    # "医療": "名詞 + の + 名詞 医療の現場 医療の制度 医療の費用 医療の質 医療の提供 医療 + を + 動詞 医療を受ける 医療を提供する 医療を改善する 医療を支える 医療を改革する ～に + 医療（を／が） 地域に医療を提供する 国に医療制度がある 社会に医療の格差がある 高齢者に医療が必要だ 災害時に医療が不足する 形容詞 + 医療／医療的 高度な医療 先進的な医療 公的な医療 緊急医療 必要不可欠な医療 医療 + 名詞 医療機関 医療従事者 医療費 医療技術 医療サービス",
    # "厄介": "形容詞 + 名詞 厄介な問題 厄介な相手 厄介な話 厄介な状況 厄介な事態 名詞 + を + 動詞 厄介を引き受ける 厄介をかける 厄介をみる 厄介をさせる 厄介を避ける ～に + 厄介（なる／巻き込まれる／なる） 親の厄介になる 人間関係に厄介になる 問題に厄介になる トラブルに巻き込まれて厄介になる ある事柄に厄介になる 副詞・接続語 + 厄介（に／で） 非常に厄介になる ちょっと厄介である とても厄介に思う やっかいにしてしまう 厄介でたまらない",
    # "収める": "動詞 + 収める／納める 成功を収める 利益を収める 学費を納める 税金を納める 紛争を収める 名詞 + を + 収める／納める 鋼材を倉庫に収める 書物を目録に収める 注文の品を納める 売上を市に納める 本を本棚に収める 収める + 名詞 収める成果 収める勝利 納める金額 納める義務 収める利益 ～に／～へ + 収める／納める 顧客に利益を収める 国へ税金を納める 書庫に資料を収める 工場に製品を収める 渋滞を収める（暴動などを） 形容詞・副詞 + 収める／納める 無事に収める きちんと納める 早く収める 適切に納める 確実に収める",
    # "叱る": "名詞 + を + 叱る 子どもを叱る 部下を叱る 生徒を叱る 犬を叱る 親を叱る （ただし使われる場面は少なめ） ～に + 叱る／叱られる 子どもに叱る 誤りをした人に叱る 教師に叱られる 親に叱られる 先生に叱られる 叱る + 名詞・補語 叱る声 叱る理由 叱る態度 叱る口調 叱る必要性 副詞・修飾語 + 叱る 厳しく叱る きつく叱る 強く叱る 優しく叱る 時折叱る 叱る + 動詞・句 叱ってしまう 叱ってもらう 叱ってやる 叱られて反省する 叱られて気づく",
    # "合図": "名詞 + の + 合図 出発の合図 開始の合図 終了の合図 合図のタイミング 合図の意味 合図 + を + 動詞 合図を出す 合図を送る 合図を忘れる 合図をする 合図を見落とす ～の + 合図 + に（なる／なるもの） それを合図にして始める その言葉が合図になった 笛が合図になって一斉に動く 目配せが合図になった 手を挙げるのが合図になる 形容詞・修飾語 + 合図 小さな合図 微妙な合図 明確な合図 無言の合図 見えにくい合図",
    # "周囲": "名詞 + の + 名詞 周囲の人々 周囲の環境 周囲の状況 周囲の目 周囲の人たち ～の + 周囲 + に／で／へ／から 周囲に注意する 周囲で観察する 周囲から見える 周囲へ声をかける 周囲に溶け込む 周囲 + を + 動詞 周囲を見渡す 周囲を気にする 周囲を観察する 周囲を騒がせる 周囲を取り巻く 形容詞 + 周囲／周囲の 広い周囲 狭い周囲 近い周囲 快適な周囲 静かな周囲",
    # "下ろす": "名詞 + 下ろす／降ろす／卸す *（何を「下ろす／降ろす／卸す」か） *荷物を下ろす *お金を下ろす *ブラインドを下ろす *看板を下ろす／降ろす *商品を卸す *下ろす／降ろす + 名詞 *（「下ろす」したものを指す） *下ろした荷物 *下ろしたブラインド *下ろした看板 *卸した商品 *下ろしたお金 *副詞 + 下ろす／降ろす／卸す *（動作の様子・程度を表す） *ゆっくり下ろす *真っすぐ下ろす *一気に下ろす *こっそり下ろす *大量に卸す *下ろす／降ろす + 助詞句 *（目的・場所などを表す語との組合せ） *～から下ろす（例：棚から下ろす） *～に下ろす（例：床に下ろす） *～から降ろす（例：バスから降ろす） *～へ卸す（例：卸先に卸す） *～まで下ろす *慣用表現・慣用句的用法 *（比喩的・慣用的な使い方） *肩の荷を下ろす *胸をなで下ろす *人をこき下ろす *根を下ろす *髪を下ろす",
    # "付近": "名詞 + の + 名詞（～の付近） 駅の付近 公園の付近 学校の付近 商店街の付近 駐車場の付近 名詞 + に + 動詞（付近に～がある／付近に～が見える 等） 付近にある 付近に見える 付近に位置する 付近に点在する 付近に存在する 名詞 + で + 動詞（～付近で～する） 駅付近で待ち合わせする 公園付近で散歩する 商店街付近で買い物する 交差点付近で渋滞する 道路付近で立ち止まる 副詞 + 名詞／動詞（多少修飾語とともに使われる） すぐ付近 近い付近 その付近 周辺付近 付近一帯 形容詞 + 名詞（修飾語＋付近） 周辺付近 近辺付近 直近付近 最寄り付近 住宅付近",
    # "命令": "名詞 + 動詞 「命令 + 動詞」の形で、「命令を（～する／出す／受ける／下す…）」など 命令を出す 命令を受ける 命令を下す 命令を伝える 命令を守る 動詞 + 命令 動詞と「命令」が結びつく形。「～命令」「命令～」「～を命令する」など 強い命令を（出す／下す） 上からの命令 不当な命令 違法な命令 命令権を持つ 名詞 + の + 命令 所有・修飾の形で「～の命令」「命令の～」 上司の命令 親の命令 規則の命令 軍の命令 組織の命令 命令 + による／によって 「手段・原因・根拠」を示す形で使う 命令による処分 命令によって実行される 命令による混乱 命令による変更 命令による影響 命令 + を + 名詞 命令が対象・内容を示す形 命令を遵守（じゅんしゅ）する 命令を改訂する 命令を取消す 命令を強化する 命令を遂行する",
    # "回復": "名詞 + 助詞（回復 + を／が／に）＋動詞 （対象＋回復する／回復を～する／回復が～する など） 健康が回復する 体力が回復する 備蓄を回復する 信用を回復する 名誉を回復する ～が + 回復 + する （主語が “〜が回復する” の形） 天候が回復する 景気が回復する ダイヤ（運行）が回復する 気力が回復する 意識が回復する 回復 + の + 名詞 （回復を修飾する語句） 回復力 回復期 回復過程 回復傾向 回復速度 回復 + を + 動詞 （回復を行為として表す語句） 回復を待つ 回復を早める 回復を図る 回復を促す 回復を保証する 他動詞 + 回復 + する （～が回復を “する” / “させる” の構文） 回復させる 回復する 回復させてくれる 回復できる 回復させようとする",
    # "困難": "名詞 + 助詞／動詞（困難 + を／が／に + 動詞） 困難を乗り越える 困難を克服する 困難を抱える 困難に直面する 困難が生じる 形容詞／連体修飾 + 困難 非常に困難 極めて困難 極度の困難 大きな困難 深刻な困難 困難 + の + 名詞 困難さ 困難度 困難性 困難な状況 困難な課題 困難 + を + 動詞 困難を感じる 困難を予想する 困難を克服する 困難を引き受ける 困難を理解する",
    # "基本": "形容詞 + 基本 （「基本」を修飾する語） 最も基本 ごく基本 基本的 実質的な基本 根本的な基本 基本 + の + 名詞 （基本を修飾語とする形） 基本概念 基本原則 基本構造 基本方針 基本姿勢 名詞 + を + 基本 + とする （～を基本とする） 安全を基本とする 信頼を基本とする 利益を基本とする 継続を基本とする 効率を基本とする 基本 + に + 動詞／形容詞 （基本に○○する／基本に○○である） 基本に忠実である 基本に従う 基本に立ち戻る 基本に照らす 基本に基づく 動詞 + 基本 + を + 名詞／動詞 （～を基本とする、～を基本に～する） ～を基本に据える ～を基本に考える ～を基本として採用する ～を基本に設計する ～を基本に実行する",
    # "女優": "形容詞／連体修飾 + 女優 女優を修飾する言葉 若手女優 ベテラン女優 名女優 人気女優 実力派女優 女優 + の + 名詞 女優が所有・所属・性質を示す語 女優の卵 女優の卵（意味：これから女優になる人） 女優の素質 女優の道 女優の経歴 名詞 + を + 女優 対象を女優にする・扱う語 〜を女優に抜擢する 〜を女優として起用する 〜を女優として売り出す 〜を女優として認められる 女優 + を + 動詞 女優に関する動作語 女優を目指す 女優を演じる（比喩的に） 女優を育てる 女優を支える ～と / ～として + 女優 立場・比較・役割を示す語 映画女優 舞台女優 トップ女優 国際的女優 女優として成長する",
    # "完全": "形容詞／連体修飾 + 完全 完全を修飾する言葉 まったく完全 真の完全 ほぼ完全 完全無欠な 完全な 完全 + の + 名詞 完全を修飾語とする名詞 完全性 完全版 完全犯罪 完全一致 完全燃焼 名詞 + を + 完全 + に + 動詞 （～を完全に ○○する） 障害を完全に克服する 任務を完全に遂行する 合意を完全に形成する 整備を完全に行う 制度を完全に整える 完全 + に + 動詞／形容詞 （“完全に～する／～である”の構文） 完全に満足する 完全に理解する 完全に消失する 完全に解決する 完全に支配する ～が + 完全 + に + ～する （主語 + 完全に + 動詞） システムが完全に稼働する 権力が完全に掌握される 支配が完全に行き届く 依存が完全に解消する 検査が完全に終了する",
    # "宗教": "形容詞／連体修飾 + 宗教 宗教を修飾する語 世界宗教 宗教的 伝統宗教 宗教間 国家宗教 宗教 + の + 名詞 宗教を修飾語とする語句 宗教法人 宗教観 宗教信仰 宗教儀式 宗教団体 名詞 + を + 宗教 （～を宗教にする／扱う構文など） キリスト教を宗教とする 宗教を問題視する 宗教を信じる 宗教 + を + 動詞 宗教に関する動作語 宗教を信仰する 宗教を広める 宗教を否定する 宗教を尊重する 宗教を改革する ～と / ～として + 宗教 立場・比較・役割を示す語 宗教として 宗教法人として 宗教色 宗教施設 宗教活動",
    # "幅": "名詞 + 動詞 幅が広がる 幅が狭まる 幅を持つ 幅を広げる 幅を制限する 形容詞 + 名詞 広い幅 狭い幅 十分な幅 最大の幅 適切な幅 名詞 + 名詞 幅方向 幅員（はばいん） 幅広さ 幅制限 幅寄せ 動詞 + 名詞 増す幅 減る幅 変化の幅 許容の幅 影響の幅",
    # "延期": "名詞 + 動詞 延期する 延期になる 延期を決定する 延期を発表する 延期を要請する 形容詞 + 名詞 急な延期 不可避な延期 大幅な延期 臨時の延期 長期の延期 名詞 + 名詞 延期決定 延期通知 延期措置 延期期間 延期理由 動詞 + 名詞 行事の延期 会議の延期 出発の延期 発表の延期 締切の延期",
    # "建設": "名詞 ＋ 動詞（＝「～を建設する」など） 建設する 建設を行う 建設を進める 建設に着手する 建設が進む 形容詞／修飾語 ＋ 名詞（「～な建設」） 大規模建設 急速な建設 安全な建設 共同建設 持続可能な建設 名詞 + 名詞（「～の建設」「建設～」など） 建設工事 建設費用 建設計画 建設業者 建設現場 動詞句／句動詞 ＋ 名詞（動詞＋目的語的な組み合わせ） 土地を建設する（※ただし「土地を建設する」はやや不自然。「宅地を造成して建設する」など文脈次第） 橋を建設する 道路を建設する ダムを建設する 施設を建設する",
    # "引退": "名詞 + 動詞（＝「〜を引退する／〜が引退する」など） 引退する 引退を表明する 引退を決断する 引退が囁かれる 引退を迎える 形容詞／修飾語 + 名詞（＝「〜な引退」） 突然の引退 不本意な引退 待望の引退 円満な引退 早期の引退 名詞 + 名詞（＝「〜の引退」「引退〜」など） 引退表明 引退試合 引退会見 引退セレモニー 引退後 動詞句／句動詞 + 名詞（動詞＋目的語・補語としての組み合わせ） 現役を引退する スポーツ界を引退する 役職を引退する 引退後の生活 引退を撤回する",
    # "悩む": "名詞 + 動詞（＝「～に悩む／～で悩む」など） 将来に悩む 問題に悩む 病気で悩む 人間関係で悩む 自分の性格に悩む 形容詞／修飾語 + 名詞（＝「～な悩み」） 深刻な悩み 日常的な悩み 小さな悩み 長年の悩み 心の悩み 名詞 + 名詞（＝「～の悩み／悩み～」など） 悩みの種 悩み相談 悩み解消 悩み事 悩み多き 動詞句／句動詞 + 名詞（動詞＋目的語的な組み合わせ） 悩みを抱える 悩みを聞く 悩みを話す 悩みを共有する 悩みを打ち明ける",
    # "悪魔": "名詞 + 動詞／動詞句（～が／～を～する） 悪魔がささやく 悪魔を召喚する 悪魔に取り憑かれる 悪魔を追い払う 悪魔と契約する 修飾語／形容詞 + 名詞（～な悪魔） 暗黒の悪魔 強大な悪魔 邪悪な悪魔 古の悪魔 醜い悪魔 名詞 + 名詞（～の悪魔／悪魔～など） 悪魔の使い 悪魔の契約 悪魔のささやき 悪魔の力 悪魔の影 動詞句 + 名詞（動詞＋目的語として使うパターン） 悪魔を呼ぶ 悪魔を封じる 悪魔を倒す 悪魔を封印する 悪魔を退ける",
    # "意思": "意思を伝える 意思を表示する 意思を示す 意思を確認する 意思を尊重する 形容詞／修飾語 + 名詞（＝「～な意思」） 明確な意思 強い意思 固い意思 自分の意思 自覚的な意思 名詞 + 名詞（＝「～の意思／意思～」など） 意思表示 意思決定 意思疎通 意思表示書（または 意思表明書） 意思確認 動詞句／句動詞 + 名詞（動詞＋目的語的な組み合わせ） 意思を固める 意思を曲げる 意思を持つ 意思表示を行う 意思決定を下す",
    # "憧れる": "動詞 + 憧れる （ある対象に憧れを抱くパターン） ○○に憧れる 若い時代に憧れる 遠い国に憧れる 都会生活に憧れる あの人に憧れる 名詞 + に憧れる （名詞を対象に憧れる表現） 芸能人に憧れる 有名人に憧れる あの人に憧れる モデルに憧れる 先輩に憧れる 憧れる + 対象を表す語句 （憧れる対象を補足する語句と共起） 憧れる存在 憧れる生活 憧れる夢 憧れるモデル 憧れる人物 〜に憧れを抱く／持つ （憧れを名詞形で述べる言い方） 憧れを抱く 憧れを持つ 憧れの念 憧れの的 憧れの人 憧れる + 状態／動詞補助表現 （憧れの気持ちの状態を表す語と共起） 憧れてやまない 憧れている 憧れて止まない 憧れていた 憧れて仕方がない",
    # "手術": "名詞 + を + 動詞（目的語＋動詞型） 手術を受ける 手術をする 手術を行う 手術を施す 手術を中止する 名詞 + の + 名詞（手術に関わる属性・関係語） 手術日 手術中 手術後 手術前 手術室 形容詞 + 手術 （手術を修飾する語） 緊急手術 大手術 小手術 外科手術 内視鏡手術 動詞 + の手術 / ～による手術 ～による手術 ～を目的とした手術 ～を伴う手術 ～後の手術 ～前の手術 手術 + に関する語句 手術成功 手術失敗 手術合併症 手術経過 手術計画",
    # "教科書": "名詞 + の + 教科書 国語の教科書 数学の教科書 英語の教科書 理科の教科書 歴史の教科書 教科書 + の + 名詞 教科書代 教科書代金 教科書編集 教科書本文 教科書採択 動詞 + 教科書 教科書を使う 教科書を読む 教科書を買う 教科書を配布する 教科書を改訂する 教科書 + を + 動詞 教科書を補足する 教科書を持参する 教科書を提出する 教科書を返却する 教科書を持ち運ぶ 形容詞 + 教科書 新しい教科書 分かりやすい教科書 厚い教科書 小型の教科書 高価な教科書",
    # "新鮮": "形容詞 + 名詞 （「新鮮な〜」という形で名詞を修飾する） 新鮮な空気 新鮮な野菜 新鮮な魚 新鮮な感覚 新鮮な印象 名詞 + が + 形容詞 （主語が「新鮮」である） 空気が新鮮だ 魚が新鮮だ 感覚が新鮮だ 香りが新鮮だ 出来事が新鮮だ 動詞 + 新鮮（に／だ） （動作と結びつく表現） 感じが新鮮だ 見ると新鮮だ 思いつきが新鮮だ 気分が新鮮になる 目に新鮮だ 新鮮 + 名詞（「〜の新鮮さ」など） （新鮮さ、性質を表す表現） 新鮮さ 新鮮味 新鮮感 新鮮度 新鮮性 動詞 + の + 新鮮（を／が） （「〜の新鮮を感じる」など） 味の新鮮を感じる 空気の新鮮を味わう 新鮮の感覚を得る 新鮮の息吹を感じる 新鮮の感動を覚える",
    # "日常": "名詞 + の + 日常 平凡の日常 忙しい日常 安定した日常 普通の日常 慣れた日常 日常 + 名詞 日常生活 日常会話 日常業務 日常風景 日常感覚 動詞 + 日常 日常を送る 日常を営む 日常を過ごす 日常に戻る 日常に溶け込む 日常 + に + 動詞／形容詞 日常に戻る 日常に過ごす 日常に欠かせない 日常に刺激を求める 日常に変化をつける 形容詞 + 日常 何気ない日常 忙しい日常 単調な日常 平凡な日常 幸せな日常",
    # "昇る": "名詞 + が / を + 昇る （主語が昇る、対象を通じて昇る） 太陽が昇る 月が昇る 煙が昇る 気球が昇る 霧が昇る 昇る + 目的語 / 到達点 （昇っていく対象・到達先） 地位に昇る 位が昇る 昇給する / 昇給を昇る（昇給を使うことが多い） 昇進する / 職に昇る 階段を昇る（機械などを用いる文脈で） ～に昇る （〜という段階・場所に昇る意味） エレベーターに昇る 展望台に昇る 天に昇る 神殿に昇る 舞台に昇る 名詞 + の + 昇る / ～の昇る （名詞＋「昇る」を組み込む表現） 昇竜 上昇（「昇」の語を使った派生語） 昇天 昇降（昇る＋降りるの複合語） 昇格 表現・慣用句 （「昇る」を含む定型表現） 天にも昇る気持ち 朝日が昇る 夕日が昇る（あまり一般的ではないが詩的表現） 昇給・昇進 昇段（武道などで階級が上がる）",
    # "暮らし": "名詞 + の + 暮らし 平穏な暮らし 豊かな暮らし 質素な暮らし 安定した暮らし 快適な暮らし 暮らし + の + 名詞 暮らし向き 暮らしぶり 暮らしぶりを見せる 暮らしぶりを伝える 暮らし向きを改善する 動詞 + 暮らし 暮らしを営む 暮らしを支える 暮らしを守る 暮らしを楽しむ 暮らしを立てる 暮らし + を + 動詞 暮らしを改善する 暮らしを豊かにする 暮らしを変える 暮らしを支える 暮らしを切り詰める 形容詞 + 暮らし 長い暮らし 暮らしやすい 暮らしにくい 暮らし向き良い 暮らし向き悪い",
    # "構う": "動詞 + 「構う」 気に 構う あまり 構う 全く 構わない／構うまい 干渉して 構う 世話を 構う 「構う」 + 名詞・句 構う 人 構う 態度 構う 必要 構う 価値 構う 時間 副詞 + 「構う」 全然 構わない 少しも 構わない まったく 構わない 決して 構うな 別に 構わない 「〜に構う」 （対象を示す格助詞「に」を伴うケース） 人 に構う 子ども に構う 他人 に構う 犬 に構う 存在 に構う",
    # "武器": "名詞 + 「武器」 主要 武器 重火器 武器 遠距離 武器 近接 武器 核 武器 「武器」 + 名詞・句 武器 として 武器 使用 武器 保有 武器 拡散 武器 製造 動詞 + 「武器」 武器を 持つ 武器を 携える 武器を 使用する 武器を 放棄する 武器を 改良する 副詞・修飾語 + 「武器」 強力な 武器 致命的な 武器 効果的な 武器 危険な 武器 高性能な 武器",
    # "歩道": "形容詞 + 歩道 狭い 歩道 広い 歩道 安全な 歩道 歩きやすい 歩道 車道と分離された 歩道 名詞 + 歩道 歩道 幅 歩道 設置 歩道 整備 歩道 舗装 歩道 区間 動詞 + 歩道 歩道を 歩く 歩道を 整備する 歩道を 拡張する 歩道を 確保する 歩道を 通行する 「〜の歩道」「歩道の〜」 都市 の歩道 学校 の歩道 商店街 の歩道 歩道 の段差 歩道 の縁石",
    # "消防": "名詞 + 「消防」 消防署 消防士 消防車 消防団 消防法 「消防」 + 名詞・句 消防 活動 消防 設備 消防 訓練 消防 計画 消防 組織 動詞 + 「消防」 消防を 強化する 消防を 整備する 消防を 支援する 消防を 拡充する 消防を 管理する 「〜の消防」／「消防の〜」 地域 の消防 市町村 の消防 都市 の消防 消防 の責任 消防 の役割",
    # "清潔": "形容詞 + 「清潔」 清潔 な 衣服 清潔 な 部屋 清潔 な 環境 清潔 な 台所 清潔 な 空気 名詞 + 「清潔」 清潔 感 清潔 度 清潔 さ 清潔 基準 清潔 意識 動詞 + 「清潔」 清潔を 保つ 清潔に する 清潔を 保てる 清潔を 心がける 清潔を 維持する 「〜の清潔」／「清潔の〜」 体 の清潔 手 の清潔 衣服 の清潔 食器 の清潔 室内 の清潔",
    # "減る": "名詞 + 減る 人口が 減る 売上が 減る 収入が 減る 支出が 減る 数が 減る 「〜が減る」 （主語を表す「が」を伴うパターン） 森林が 減る 森が 減る 森林面積が 減る 飲酒量が 減る 使用量が 減る 副詞 + 減る / 減っている だんだん 減る 徐々に 減る 大幅に 減る 急激に 減る 著しく 減っている 「〜を減らす」対応（関連コロケーション） （「減る」は自動詞だが、他動詞「減らす」と対になる表現も合わせて知っておくと便利） コストを 減らす 支出を 減らす 食費を 減らす 人員を 減らす ゴミを 減らす",
    # "物語": "名詞 + 物語 長編物語 短編物語 冒険物語 恋愛物語 伝説物語 「物語」 + 名詞・句 物語 風景 物語 世界 物語 舞台 物語 結末 物語 語り手 動詞 + 物語 物語を 紡ぐ 物語を 描く 物語を 語る 物語が 進む 物語が 展開する 形容詞・修飾語 + 物語 感動的な 物語 重厚な 物語 美しい 物語 古典的な 物語 幻想的な 物語 「〜の物語」「物語の〜」 人生 の物語 少年 の物語 家族 の物語 物語 の始まり 物語 の結末",
    # "現れる": "名詞 + 「現れる」 影が現れる 物音が現れる 現象が現れる 症状が現れる 本性が現れる 「〜に / 〜から」 + 「現れる」 （場所・起点を表す副詞的語を伴う） 闇に現れる 空から現れる 山中に現れる 森の中から現れる 舞台に現れる 副詞 / 副詞句 + 「現れる」 （動作を修飾する） 突然現れる 急に現れる 徐々に現れる ゆっくり現れる ぱっと現れる 形容詞 / 形容詞句 + 「現れる」 （どのように見えるかを強調） はっきり現れる ぼんやり現れる 微かに現れる 鮮明に現れる 明確に現れる 動詞句 + 「現れる」 （動詞句で条件や変化を表す） 姿を現れる → 姿を現す（※ただし「現す」が使われる方が一般的） 表情が現れる 表情に現れる 気持ちが現れる 著しく現れる",
    # "発行": "名詞 + 「発行」 雑誌を発行 刊行物を発行 債券を発行 免許を発行 証書を発行 「〜を／〜が」 + 「発行する／発行される」 （他動・自動の文型で使われる） 券を発行する 免許が発行される 通貨を発行する 証明書を発行する チケットを発行する 副詞 / 副詞句 + 「発行」 （発行の様態を表す副詞的語） 定期的に発行する 随時発行する 新たに発行する 大量に発行する 再度発行する 形容詞 / 形容詞句 + 「発行」 （発行に関する性質を表す語） 正式に発行する 無料で発行する 迅速に発行する 適切に発行する 正確に発行する 動詞句 + 「発行する」 （動詞句と結びついて使われる表現） 発行手続きをする 発行開始する 発行停止する 発行申請する 発行許可する",
    # "真っ赤": "形容詞 + 名詞（“真っ赤な X”） 真っ赤な花 真っ赤な夕焼け 真っ赤な顔 真っ赤なリンゴ 真っ赤な旗 名詞 + 「真っ赤になる／真っ赤に染まる／真っ赤にする」などの動詞句 顔が真っ赤になる 空が真っ赤に染まる 顔を真っ赤にする 葉が真っ赤に染まる 血が真っ赤になる 副詞 / 副詞句 + 「真っ赤に／真っ赤で」など とても真っ赤に じわじわ真っ赤に すっかり真っ赤に めっきり真っ赤に 完全に真っ赤で 動詞句 + 「真っ赤」＋ 補語（状態を表す語） 赤く真っ赤（になる） 染めて真っ赤 焼けて真っ赤 燃えて真っ赤 落ちて真っ赤 比喩的用法／転用表現 + 「真っ赤」 （比喩・強調表現） 真っ赤な嘘 真っ赤なウソ 真っ赤な嘘をつく 真っ赤な罪 真っ赤な間違い",
    # "睡眠": "名詞 + 「を／が」 + 動詞パターン（睡眠と動詞の結びつき） 睡眠を取る 睡眠をとる 睡眠をとるようにする 睡眠を確保する 睡眠を妨げる 睡眠を改善する 睡眠が浅い 睡眠が深い 睡眠が乱れる 睡眠が足りない 副詞 / 副詞句 + 「睡眠」 + 動詞／述語パターン 十分な睡眠を取る 適切な睡眠を確保する 良質な睡眠をとる 長時間の睡眠をとる 十分な睡眠が取れていない 質の高い睡眠を得る 「睡眠」 + 助詞 + 名詞／形容詞パターン 睡眠時間 睡眠不足 睡眠障害 睡眠状態 睡眠習慣 形容詞 + 「睡眠」 良質な睡眠 深い睡眠 快適な睡眠 十分な睡眠 睡眠不足の 動詞句 + 「睡眠」 + 補語／修飾語（睡眠の状態を表す語を伴う表現） 睡眠が妨げられる 睡眠が妨害される 睡眠の質が低下する 睡眠の質を高める 睡眠のサイクルを整える",
    # "確かめる": "名詞 + 「を」 + 「確かめる」 身分証明書を確かめる 住所を確かめる 電話番号を確かめる 予約内容を確かめる 情報を確かめる 「〜が」 + 「確かめられる／確かめる」 真偽が確かめられる 事実が確かめられる 安全が確かめられる 正確さが確かめられる 効果が確かめられる 副詞 / 副詞句 + 「確かめる」 きちんと確かめる 念入りに確かめる 再度確かめる もう一度確かめる ちゃんと確かめる 形容詞 / 形容詞句 + 「確かめる」 正しく確かめる 十分に確かめる 完全に確かめる 正確に確かめる 確実に確かめる 動詞句 + 「確かめる」 確かめたうえで〜する 確かめる手続きをする 確かめておく 確かめに行く 確かめ直す",
    # "礼儀": "名詞 + 「の」 + 「礼儀」 礼儀の作法 礼儀の欠如 礼儀の問題 礼儀の基本 礼儀の範囲 「礼儀」 + 助詞 + 名詞／語句 礼儀作法 礼儀正しさ 礼儀正しい態度 礼儀正しい言葉遣い 礼儀正しい対応 動詞 + 「礼儀を〜／礼儀を重んじる／礼儀を守る」など 礼儀を重んじる 礼儀を守る 礼儀を欠く 礼儀をわきまえる 礼儀を示す 副詞／副詞句 + 「礼儀正しく〜」など 非常に礼儀正しく きちんと礼儀正しく もっと礼儀正しく ある程度礼儀正しく 比較的礼儀正しく",
    # "穏やか": "形容詞 + 名詞（“穏やかな ××”） 穏やかな天気 穏やかな海 穏やかな表情 穏やかな笑顔 穏やかな暮らし 名詞 + 「で／な／に」 + 「穏やか」／変化表現 心が穏やかである 気持ちが穏やかになる 空気が穏やかだ 日差しが穏やかだ 波風が穏やかだ 副詞 / 副詞句 + 「穏やかに」 穏やかに話す 穏やかに暮らす 穏やかに見守る 穏やかに進む 穏やかに過ごす 動詞句 + 「穏やか」／「穏やかである」 穏やかさを感じる 穏やかさを保つ 穏やかな雰囲気を醸す ～を穏やかにする ～を穏やかに受け止める",
    # "立派": "形容詞 + 名詞（“立派な ××”） 立派な人物 立派な建物 立派な仕事 立派な成果 立派な態度 名詞 + 「だ／である／に」 + 「立派」／形容動詞化 人は立派だ 建物は立派である 彼は立派な人だ それは立派なことだ 立派なものである 動詞 + 「立派に〜する／立派になる」など 立派に育つ 立派に成長する 立派に振る舞う 立派に見える 立派に務める 副詞 / 副詞句 + 「立派に」 とても立派に 実に立派に なかなか立派に 立派に…している 見事に立派に 動詞句 + 「立派」 + 補語／修飾語 立派さを感じる 立派さを示す 立派さを備える 立派さを失う 立派な様子を呈する",
    # "筋肉": "名詞 + 「の」 + 「筋肉」 筋肉の収縮 筋肉の量 筋肉の発達 筋肉の疲労 筋肉の強さ 名詞 + 「筋肉」 + 助詞 / 動詞パターン 筋肉を鍛える 筋肉を増やす 筋肉をつける 筋肉を落とす 筋肉がつく 筋肉が減る 筋肉がつらい 筋肉が疲れる 形容詞 / 形容動詞 + 「筋肉」 太い筋肉 細い筋肉 強い筋肉 強靭な筋肉 弱い筋肉 副詞 / 副詞句 + 「筋肉」 + 動詞／形容語パターン しっかり筋肉を鍛える 効率よく筋肉をつける 無理せず筋肉を増やす 徐々に筋肉がつく 急激に筋肉が落ちる 動詞句 + 「筋肉」 + 補語／修飾語 筋肉の発達を促す 筋肉の維持を図る 筋肉の修復を助ける 筋肉痛を感じる 筋肉痛が残る",
    # "組む": "名詞 + 「を」 + 「組む」 チームを組む 組織を組む 計画を組む 時間割を組む 日程を組む 「〜と／〜で／〜に」 + 「組む」 彼と組む 二人で組む 他社と組む 同盟を組む 連携を組む 体の部分 + 「を組む」 腕を組む 足を組む 手を組む 膝を組む 指を組む 動詞句 + 「組む」 組んで協力する 組んで進める 組んで取り組む 組んで対応する 組んで解決する 形容詞 / 副詞 + 「組む」 密に組む 緊密に組む 綿密に組む 緻密に組む 強固に組む",
    # "芸能": "名詞 + 「の」 + 「芸能」 芸能界の 芸能活動の 芸能界の事情 芸能界の人脈 芸能界の常識 「芸能」 + 助詞 + 名詞／語句 芸能界 芸能人 芸能活動 芸能界入り 芸能界デビュー 動詞 + 「芸能を／芸能界を〜」 芸能界に入る 芸能活動をする 芸能界で活躍する 芸能界を引退する 芸能界を席巻する 副詞／副詞句 + 「芸能界で／芸能活動を〜」 精力的に芸能活動をする 本格的に芸能界入りする 積極的に芸能活動を行う 着々と芸能界でキャリアを積む 間もなく芸能界デビューする",
    # "見舞い": "名詞 + 「の」 + 「見舞い／お見舞い」 入院の見舞い 病気のお見舞い 火事の見舞い 出産のお見舞い 見舞い客 動詞 + 「見舞い／お見舞いを〜」 見舞いに行く／行う お見舞いをする お見舞いを渡す お見舞いを送る お見舞いを受ける 「〜に」 + 「お見舞い」 病院にお見舞い（に行く） 家にお見舞い（に行く） 入院先にお見舞い（に行く） 会社にお見舞い（をする） 形容詞 / 副詞 + 「お見舞い」／「見舞い」＋ 動詞 心からのお見舞い 心ばかりのお見舞い ちょっとしたお見舞い お見舞い申し上げます 早期のお見舞い 動詞句 + 「見舞い／お見舞い」 + 補語／修飾語 見舞い品を贈る お見舞い状を書く お見舞いの言葉を述べる お見舞い金を包む 見舞い金を送る",
    # "覚める": "名詞 + 「が」 + 「覚める／覚めて」 夢が覚める 酔いが覚める 意識が覚める 幻想が覚める 情熱が覚める 副詞 / 副詞句 + 「覚める／覚めて」 すっと覚める ぱっと覚める だんだん覚める 目が覚めるように覚める はっと覚める 動詞句 + 「覚める／覚めて」 + 補語／修飾語 覚めた目で見る 覚めた心で受け止める 酔いがすっかり覚める 夢から覚めて現実を見る 衝動が覚めて冷静になる",
    # "計算": "名詞 + 「を」 + 「計算する」 数字を計算する 面積を計算する 金額を計算する 税金を計算する 利益を計算する 名詞 + 「の」 + 「計算」 電卓の計算 パソコンの計算 数学の計算 機械の計算 複雑な計算 形容詞 / 副詞 + 「計算」 正確な計算 複雑な計算 簡単な計算 大まかな計算 精密な計算 「計算」+ 動詞句 計算が合う 計算が合わない 計算が早い 計算に強い 計算に弱い 慣用的表現 計算に入れる 計算から外す 計算を誤る 計算を間違える 計算通りに進む",
    # "許可": "名詞 + 「を」 + 「許可する」 入場を許可する 使用を許可する 建築を許可する 輸入を許可する 開発を許可する 名詞 + 「の」 + 「許可」 営業の許可 建築の許可 使用の許可 出国の許可 立ち入りの許可 形容詞 / 副詞 + 「許可」 特別に許可する 例外的に許可する 正式に許可する 一時的に許可する 厳格に許可する 「許可」+ 動詞句 許可が下りる 許可を得る 許可を申請する 許可を与える 許可を取り消す 慣用的表現 許可なく立ち入る 許可なしで使用する 許可のもとに行う 許可を待つ 許可を必要とする",
    # "診察": "名詞 + 「を」 + 「診察する」 患者を診察する 病人を診察する 子どもを診察する 外来を診察する 内科を診察する 名詞 + 「の」 + 「診察」 医師の診察 今日の診察 外来の診察 定期の診察 次回の診察 「診察」+ 動詞句 診察を受ける 診察を終える 診察を待つ 診察を始める 診察を希望する 形容詞 / 副詞 + 「診察」 丁寧に診察する 慎重に診察する しっかり診察する 正確に診察する 詳しく診察する 関連的表現 診察券を出す 診察室に入る 診察時間が始まる 診察順を待つ 診察日を決める",
    # "評価": "名詞 + 「を」 + 「評価する」 成績を評価する 能力を評価する 成果を評価する 業績を評価する 商品を評価する 名詞 + 「の」 + 「評価」 作品の評価 人の評価 市場の評価 社会の評価 科学的評価 形容詞 / 副詞 + 「評価」 高い評価 正当な評価 客観的な評価 厳しい評価 公平な評価 「評価」+ 動詞句 評価を受ける 評価を下す 評価を得る 評価を与える 評価が高まる 慣用的表現 評価に値する 評価の対象となる 評価を誤る 評価を左右する 評価を行う",
    # "詳しい": "名詞 + 「に詳しい」 歴史に詳しい 音楽に詳しい 経済に詳しい 法律に詳しい パソコンに詳しい 「詳しい」+ 名詞 詳しい説明 詳しい内容 詳しい情報 詳しい地図 詳しい手順 副詞 + 「詳しい」 とても詳しい かなり詳しい たいへん詳しい すごく詳しい わりと詳しい 動詞句 + 「詳しい」関連 詳しいことを知る 詳しい人に聞く 詳しい資料を読む 詳しい話をする 詳しい状況を調べる 慣用的な使い方 詳しいことは分からない 詳しい説明を求める 詳しい手続きを案内する 詳しい経緯を明らかにする 詳しい事情を説明する",
    # "誇り": "名詞 + 「を」 + 「誇りに思う」 日本を誇りに思う 仕事を誇りに思う 家族を誇りに思う 成績を誇りに思う 出身を誇りに思う 「誇り」+ 助詞 + 名詞／動詞句 誇りを持つ 誇りを抱く 誇りを失う 誇りを守る 誇りを傷つける 形容詞 / 副詞 + 「誇り」 大きな誇り 高い誇り 強い誇り 最大の誇り 無上の誇り 慣用的表現 誇り高い精神 誇りある伝統 誇り高き民族 誇りの象徴 誇りの源 動詞句 + 「誇り」 誇りに耐える 誇りをかけて戦う 誇りを胸に抱く 誇りを取り戻す 誇りを示す",
    # "豊富": "「豊富な」+ 名詞 豊富な知識 豊富な経験 豊富な資源 豊富な種類 豊富な食材 名詞 + 「が豊富だ」 栄養が豊富だ 水が豊富だ 資金が豊富だ 人材が豊富だ 情報が豊富だ 副詞 + 「豊富に」 非常に豊富に きわめて豊富に 比較的豊富に 最も豊富に かなり豊富に 動詞句 + 「豊富」 バリエーションが豊富 メニューが豊富 選択肢が豊富 品揃えが豊富 経験が豊富 慣用的な使い方 豊富な才能 豊富な感性 豊富な表現力 豊富な内容 豊富な色彩",
    # "資本": "名詞 + 「の」 + 「資本」 企業の資本 国家の資本 外国の資本 金融の資本 株主の資本 名詞 + 「を」 + 「資本にする／資本とする」 知識を資本にする 経験を資本にする 技術を資本にする 人脈を資本にする 労働力を資本とする 「資本」+ 名詞（複合語） 資本主義 資本政策 資本提携 資本金 資本市場 動詞 + 「資本」 資本を投入する 資本を集める 資本を導入する 資本を持つ 資本を運用する 形容詞 / 副詞 + 「資本」 外国資本 多額の資本 巨大な資本 豊富な資本 不足する資本",
    # "賢い": "賢い」+ 名詞 賢い人 賢い子ども 賢い選択 賢い方法 賢い判断 名詞 + 「が賢い」 犬が賢い 子どもが賢い 生徒が賢い 彼が賢い 彼女が賢い 副詞 + 「賢い」 とても賢い 実に賢い かなり賢い すごく賢い 比較的賢い 動詞句 + 「賢い」関連 賢く生きる 賢く使う 賢く判断する 賢く選ぶ 賢く対応する 慣用的表現 賢い考え 賢い生き方 賢い消費者 賢い買い物 賢い言い方",
    # "迎え": "名詞 + 動詞（迎え + 動詞） 迎えに行く 迎えに来る 迎えを頼む 迎えに向かう 迎えに出る 動詞 + 名詞（動詞 + 迎え） 出迎えをする 出迎えに出る 歓迎と迎え（例：歓迎と迎え） 迎え役を務める 迎えの車 名詞 + 名詞（迎え + 名詞） 迎えの車 迎えの時間 迎えの言葉 迎えの準備 迎えの連絡 形容詞/連体修飾 + 迎え（…な迎え、…の迎え） 温かい迎え 盛大な迎え 簡素な迎え 予期せぬ迎え 心強い迎え",
    # "過ごす": "名詞 + を + 過ごす（名詞を目的語にとるパターン） 時間を過ごす 日々を過ごす 一日を過ごす 休日を過ごす 平和な日々を過ごす ～を + どう過ごす（疑問文・補語を使ったパターン） どう過ごすか ～をどう過ごすか考える ～を有意義に過ごす ～を快適に過ごす ～を楽しく過ごす 副詞／形容詞 + に + 過ごす（様態を表す） ゆったりと過ごす のんびりと過ごす 静かに過ごす 穏やかに過ごす 無事に過ごす 〜と / 〜とともに + 過ごす（共同行動・時間を共にする） 家族と過ごす 友人と過ごす 恋人と過ごす 一緒に過ごす 大切な人と過ごす",
    # "金額": "名詞 + の + 金額 総額の金額 見積額の金額 支払う金額 請求金額 控除後の金額 金額 + を + 動詞（動詞を伴う） 金額を確認する 金額を提示する 金額を支払う 金額を記入する 金額を交渉する 副詞 / 形容詞 + に + 金額（様態修飾） 大きな金額 少額の金額 高額な金額 適正な金額 想定外の金額 （～と / ～によって） + 金額 + が / を + 動詞（関係・比較） 金額が異なる 金額による差 金額に応じて 金額が膨らむ 金額が決まる",
    # "鋭い": "形容詞 + 名詞（鋭い + 名詞） 鋭い 痛み 鋭い 視線 鋭い 批判 鋭い 感覚 鋭い 刃 名詞 + を + 形容詞 （名詞を修飾する形に展開） 切れ味が 鋭い ナイフ 批評が 鋭い 言葉 感覚が 鋭い 人 目つきが 鋭い 表情 洞察力が 鋭い 分析 動詞 + 形容詞化／副詞化（鋭く + 動詞など） 鋭く 切り込む 鋭く 指摘する 鋭く 反応する 鋭く 迫る 鋭く 感じる 副詞・連体修飾語 + に + 鋭い 非常に 鋭い 突然 鋭い 一段と 鋭い さらに 鋭い 異様に 鋭い",
    # "除く": "名詞 + を + 除く（対象を目的語にとるパターン） 不良品を除く 雑草を除く 障害を除く 例外を除く ごみを除く 名詞 + を除いて（範囲を限定する／除外を示す表現）【文法表現】 Bunpro +2 JSMORI +2 土日を除いて 税金を除いて アルバイトを除いて 彼を除いて 場所を除いて 動詞 + 名詞（…を除く + 名詞、修飾表現として使う） 除いた部分 除く対象 除くべきもの 除いてよい要素 除いたあとの状態 副詞・形容詞 + に + 除く（修飾形） 例外的に除く 明示的に除く 適切に除く 完全に除く 一部を除く",
    # "隠す": "名詞 + を + 隠す（目的語をとるパターン） 真実を隠す 秘密を隠す 欠点を隠す 本音を隠す 財産を隠す ～を + 隠して（手段・方法を表す形） 影で隠して 声を隠して 身を隠して 顔を隠して 本当の意図を隠して 動詞 + 隠す（動詞と組む表現） 意図を隠す 感情を隠す 事実を隠す 痛みを隠す 弱みを隠す 副詞／形容詞 + に + 隠す（様子・程度を表す修飾語） こっそり隠す 完全に隠す うまく隠す 巧妙に隠す 密かに隠す",
    # "集まり": "名詞 + 動詞 （名詞 “集まり” を目的語・主語にとる動詞との結びつき） 集まり がある 集まり が悪い 集まり に参加する 集まり を開く 集まり を設ける 形容詞 + 名詞 （“どのような集まりか” を表す） 少ない集まり 多い集まり 小さな集まり 親しい集まり 非公式な集まり 名詞 + の + 名詞 （“～の集まり” の形で使われる） 友達の集まり 家族の集まり 同窓会の集まり 趣味の集まり 地域の集まり 動詞 + の + 名詞 （“～の集まり” を動詞化した表現） 会う の集まり（口語的に使われることも） 集まる の集まり 開く の集まり 整える の集まり（例：準備を整えるの集まり、は少し不自然だが見かけることも） 決める の集まり 動詞句 + 名詞 （他の動詞句と組み合わせて使われる表現） 予定された集まり 突発的な集まり 定期的な集まり 昼間の集まり 夜の集まり",
    # "震える": "主語 + 震える （何が震えるのかを表す） 手が震える 足が震える 声が震える 全身が震える 肩が震える 震える + 助詞 + 名詞 / 震える + 名詞 （震える対象や原因などを続ける表現） 震える ほど 震える 音 震える 手で 震える 声で 震える 笑顔（比喩的表現） 震える + 助詞 + 動詞 （震えながら〜する、などの表現） 震えながら話す 震えながら立つ 震えながら笑う 震えながら祈る 震えながら歌う 原因・理由を表す語 + で / に + 震える （なぜ震えるかを表す構造） 寒さで震える 緊張で震える 恐怖で震える 感動で震える 驚きに震える 震える + 補助表現（連用形や副詞句） ぶるぶる震える ぞくぞく震える ふるふる震える 微かに震える 激しく震える",
    # "非常": "副詞 + 形容詞／動詞 （非常に～する／非常に～だ） 非常に 重要だ 非常に 難しい 非常に 助かる 非常に 驚く 非常に 幸運だ 形容動詞 + 名詞 （「非常な～」の形で使われることが多い） 非常な 事態 非常な 状態 非常な 危険 非常な 事態 非常な 努力 名詞 + の + 非常 （「～の非常（～の緊急性）」の意味合いで使われる） 命の非常 社会の非常 国の非常 状況の非常 地域の非常 （ただしこのパターンはやや文語的・限定的な用法） 非常 + 【動詞句】 （「非常～する／非常～させる」などの構造） 非常を 要する 非常を 呼びかける 非常を 宣言する 非常を 認める 非常を 維持する",
    # "頂く": "主語 + 頂く （目上の人から何かを「頂く」） 賞を頂く 承認を頂く ご支援を頂く ご連絡を頂く ご意見を頂く ～して + 頂く （「～してもらう」の謙譲表現：「～して頂く／頂きたい」など） ご教示して頂く ご協力して頂く ご連絡して頂く ご説明して頂く ご検討して頂く ～させて + 頂く （自分の行為について、相手の許可を得て行う表現） ご案内させて頂く ご説明させて頂く ご紹介させて頂く ご相談させて頂く ご決定させて頂く お／ご + 名詞 + 頂く （名詞に「お／ご」を付けて敬意を込め、「頂く」を使う） ご利用頂く ご理解頂く ご対応頂く ご連絡頂く ご評価頂く ～て + 頂く（補助動詞的用法） （「～てもらう」の丁寧な言い方） ご承知いただく ご覧いただく お聞きいただく お話しいただく ご確認いただく",
    # "頭痛": "主語 + 動詞（〜がする／〜が起こる など） 頭痛がする 頭痛が起こる 頭痛が続く 頭痛がひどい 頭痛が治る 名詞 + の + 名詞 偏頭痛の 発作 緊張型頭痛の 症状 片頭痛の 発作 頭痛の 原因 頭痛の 種類 形容詞／副詞 + 名詞 激しい頭痛 軽い頭痛 ひどい頭痛 頭痛のひどさ 頭痛の頻度 原因・理由 + で／に + 頭痛 ストレスで頭痛 疲れで頭痛 寝不足で頭痛 騒音に頭痛 光に頭痛 補助語句 + 頭痛 頭痛持ち（＝頭痛が起こりやすい人） 頭痛薬 頭痛外来 頭痛日記 頭痛予防",
    # "飽くまで": "飽くまで + 名詞 飽くまで 自説 飽くまで 方針 飽くまで 反対 飽くまで 可能性 飽くまで 意見 飽くまで + 動詞／句（～する／～しようとする etc.） 飽くまで 主張する 飽くまで 貫く 飽くまで 反対する 飽くまで 求める 飽くまで 守ろうとする 副詞＋飽くまで （「飽くまでも」の形で強調を加える） あくまでも 意見として あくまでも 私見だが あくまでも その範囲で あくまでも 形式的に あくまでも 名目上",
    # "髪の毛": "主語 + 動詞（～が～する／～を～する 等） 髪の毛 が抜ける 髪の毛 が伸びる 髪の毛 を切る 髪の毛 を染める 髪の毛 を洗う 形容詞／副詞 + 髪の毛 長い髪の毛 短い髪の毛 柔らかい髪の毛 ふさふさの髪の毛 傷んだ髪の毛 髪の毛 + の + 名詞 髪の毛 の量 髪の毛 の状態 髪の毛 の質 髪の毛 の色 髪の毛 の手入れ 原因・理由／条件 + で／に + 髪の毛 ストレス で 髪の毛が抜ける 加齢 に よる髪の毛の変化 栄養不足 で 髪の毛が細くなる 紫外線 に よって髪の毛が傷む シャンプー で 髪の毛を洗う 補助語＋髪の毛 髪の毛 先（例：毛先、先端） 髪の毛 束（例：一束／二束の髪の毛） 髪の毛 量（多い／少ない） 髪の毛 質（硬さ・柔らかさなど） 髪の毛 色（黒髪／白髪／茶髪など）",
    # "黒板": "主語 + 動詞（～する／～が～する など） 黒板 に書く 黒板 を消す 黒板 を拭く 黒板 を使う 黒板 が見える 形容詞／副詞 + 黒板 古い黒板 大きな黒板 壁の黒板 掲示用黒板 移動式黒板 黒板 + の + 名詞 黒板 文字 黒板 消しゴム（黒板消し） 黒板 面 黒板 板書 黒板 表面 用途・機能を表す語 + 黒板 授業用黒板 会議用黒板 電子黒板（デジタル黒板） 企業用デジタルホワイトボード 『ミライタッチBiz』 移動黒板 両面黒板 原因・手段 + で／に + 黒板 チョークで 黒板 に書く Wikipedia 指で 黒板 を触る 消しゴムで 黒板 を消す 雨で 黒板 が汚れる 風で 黒板 に粉が舞う",
    words_dict = {}

    for word, raw_collocations in words_dict.items():
        w = JPWord2(word, 3, raw_collocations=raw_collocations, ai_init=True)

    # w = JPWord2("かなり", 3, ai_init=True)
    # w = JPWord2.load_from_json("瞬間")

    # w.tts()
    # w.pptx_generation()
